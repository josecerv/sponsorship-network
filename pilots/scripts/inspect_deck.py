"""
inspect_deck.py — read-only walker that prints what's on every slide of
docs/UChicago-0410.pptx so we can see what manual edits exist before
patching anything.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

DECK = Path(__file__).resolve().parent.parent.parent / "docs" / "UChicago-0410.pptx"

def emu_to_in(emu):
    return round(emu / 914400.0, 2)

def shape_summary(shape):
    parts = [f"{shape.shape_type}"]
    try:
        parts.append(f"x={emu_to_in(shape.left)}\" y={emu_to_in(shape.top)}\" "
                     f"w={emu_to_in(shape.width)}\" h={emu_to_in(shape.height)}\"")
    except Exception:
        pass
    if shape.has_text_frame:
        text = shape.text_frame.text.strip().replace("\n", " | ")
        if text:
            text = text[:90] + ("…" if len(text) > 90 else "")
            parts.append(f"text={text!r}")
    if shape.shape_type == 13:  # PICTURE
        try:
            img_part = shape.image
            parts.append(f"img={img_part.filename or '[embedded]'}")
        except Exception:
            parts.append("img=?")
    return "  ".join(parts)

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    print(f"{len(prs.slides)} slides\n")
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== Slide {i} ({len(slide.shapes)} shapes, layout='{slide.slide_layout.name}') ===")
        for shape in slide.shapes:
            print(f"  - {shape_summary(shape)}")
        print()

if __name__ == "__main__":
    main()
