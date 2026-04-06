"""
build_uchicago_deck.py
======================
Programmatically build the UChicago talk deck (20 slides) from:
  - the Maurice deck as a base (for theme/masters)
  - rendered Stage 3 stimulus PNGs from render_stage3_screens.py
  - rendered RQ figures from make_talk_figures.R
  - the hook composite from render_hook_card.py
  - extracted Maurice slide screenshots (Stage 1 task imagery, candidate
    pair screen, etc.) for slides where Jose explicitly liked the original

Output: docs/UChicago-0410.pptx

Outline (20 slides):
   1. Cover (Maurice "Who Builds Social Capital..." restored, dated for UChicago)
   2. Hook — outcome-focused (two sponsor cards + outcome reveal banner)
   3. What is sponsorship?
   4. Sponsorship involves real risk for sponsors
   5. Research questions (3 RQ cards)
   6. Multi-stage lab experiment overview
   7. Stage 1: where the candidates come from (Maurice task screenshots)
   8. Stage 2: real sponsor endorsements (with embedded candidate-pair screen)
   9. Stage 3 design: process diagram (3 manipulations × evaluator flow)
  10. Stage 3 walkthrough — D1: initial wager
  11. Stage 3 walkthrough — D2: outcome reveal
  12. Stage 3 walkthrough — D3: second wager
  13. RQ1 result
  14. RQ2 — strong endorsers
  15. RQ2 — weak endorsers
  16. RQ3 result
  17. Synthesis
  18. Theoretical implications
  19. Limitations & next steps
  20. Thank you

Removed vs. previous draft (per Jose 2026-04-06):
  - 8-conditions-at-a-glance grid (was old slide 13)
  - Sample & exclusions slide (was old slide 14)
  - Coded condition names like "W_correct_strong" replaced with plain English
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

SCRIPT_DIR  = Path(__file__).resolve().parent
PROJECT     = SCRIPT_DIR.parent.parent
DOCS_DIR    = PROJECT / "docs"
OUT_DIR     = PROJECT / "pilots" / "output"

BASE_PPTX   = DOCS_DIR / "MauriceLab-1113.pptx"
OUTPUT_PPTX = DOCS_DIR / "UChicago-0410.pptx"

STAGE3_DIR  = OUT_DIR / "stage3_rendered_screens"
WALK_DIR    = STAGE3_DIR / "walkthrough"

FIG_DIR        = OUT_DIR / "talk_figures"
HOOK_PNG       = FIG_DIR / "hook_card_composite.png"
RQ1_PNG        = FIG_DIR / "rq1_gender_x_outcome.png"
RQ2_STRONG_PNG = FIG_DIR / "rq2_strong_only.png"
RQ2_WEAK_PNG   = FIG_DIR / "rq2_weak_only.png"
RQ3_PNG        = FIG_DIR / "rq3_initial_trust.png"
SYNTH_PNG      = FIG_DIR / "synthesis_outcome_sensitivity.png"
SUMMARY_JSON   = FIG_DIR / "talk_figures_summary.json"

CANDIDATE_SCREEN_PNG = FIG_DIR / "stage2_candidate_screen.png"
STAGE1_TASKS_PNG     = FIG_DIR / "stage1_tasks_cropped.png"
WHARTON_LOGO_PNG     = FIG_DIR / "wharton_logo_navy.png"

# Color palette (matches Maurice template + R figures)
NAVY     = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DK  = RGBColor(0x01, 0x14, 0x40)
RED      = RGBColor(0x99, 0x00, 0x00)
RED_LT   = RGBColor(0xC8, 0x1E, 0x1E)
GRAY     = RGBColor(0x6B, 0x72, 0x80)
GRAY_LT  = RGBColor(0x9C, 0xA3, 0xAF)
DARK     = RGBColor(0x11, 0x18, 0x27)
BG       = RGBColor(0xEE, 0xED, 0xEA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_DK = RGBColor(0x04, 0x7A, 0x4B)

# Slide dimensions (Maurice = 13.33" × 7.5", 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def delete_all_slides(prs):
    """Strip all slides from prs (preserving masters/layouts)."""
    sldIdLst = prs.slides._sldIdLst
    rId_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    for sl in list(sldIdLst):
        rId = sl.attrib[rId_attr]
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
        sldIdLst.remove(sl)

def add_blank_slide(prs):
    """Add a slide using layout 11 ('No Title') and strip its placeholders.

    The Maurice 'No Title' layout still inherits a Content Placeholder, a
    Footer Placeholder, and a Slide-Number Placeholder. In PowerPoint's edit
    view those show up as dashed boxes with 'Click to add text' and stock
    icons (the user reported this as 'clipping'). We strip them so each new
    slide is truly blank but still inherits the Maurice ribbon (Wharton logo,
    red diagonal mark) from the master.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide

def add_textbox(slide, left, top, width, height, text,
                size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font="Arial", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def add_bullets(slide, left, top, width, height, lines,
                size=20, bullet_color=NAVY, text_color=DARK, font="Arial",
                spacing=8):
    """lines: list of (str, optional bool 'bold')"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = "•  " + text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = text_color
    return box

def add_picture_centered(slide, path, top, max_width=None, max_height=None):
    """Insert image so it's horizontally centered, scaled to fit max bounds
    while preserving aspect ratio."""
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size
    aspect = iw / ih
    if max_width and max_height:
        if max_width / aspect <= max_height:
            w = max_width
            h = int(max_width / aspect)
        else:
            h = max_height
            w = int(max_height * aspect)
    elif max_width:
        w = max_width
        h = int(max_width / aspect)
    elif max_height:
        h = max_height
        w = int(max_height * aspect)
    else:
        w, h = iw, ih
    left = int((SLIDE_W - w) / 2)
    return slide.shapes.add_picture(str(path), left, top, width=w, height=h)

def set_slide_background(slide, color):
    """Set a solid background fill on the slide."""
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = color

def add_title_bar(slide, text, color=NAVY, size=32, top=0.35):
    """Standard title at the top of a content slide."""
    return add_textbox(
        slide,
        left=Inches(0.55), top=Inches(top),
        width=SLIDE_W - Inches(1.1), height=Inches(0.8),
        text=text, size=size, bold=True, color=color, font="Arial"
    )

def add_slide_number(slide, n, total=20):
    """Add the slide number as plain white text inside the Maurice ribbon.

    The Maurice slide master already has a navy ribbon at the bottom of
    every slide (with the Wharton logo on the left and a small red diagonal
    mark near the right). We just drop the number text *inside* that ribbon
    — no separate pill background, so nothing protrudes.
    """
    box = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.55), SLIDE_H - Inches(0.42),
        Inches(0.85), Inches(0.32)
    )
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n} / {total}"
    run.font.name = "Arial"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return box

def add_caption(slide, text, top, size=14, color=GRAY, italic=True):
    """Centered caption beneath an image."""
    return add_textbox(
        slide,
        left=Inches(0.6), top=top,
        width=SLIDE_W - Inches(1.2), height=Inches(0.4),
        text=text, size=size, italic=italic, color=color,
        align=PP_ALIGN.CENTER
    )

def add_rounded_box(slide, x, y, w, h, fill, line_color=None, line_width=1.0):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line_color is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line_color
        box.line.width = Pt(line_width)
    box.text_frame.margin_left = Inches(0.2)
    box.text_frame.margin_right = Inches(0.2)
    box.text_frame.margin_top = Inches(0.18)
    box.text_frame.margin_bottom = Inches(0.18)
    return box

def add_rect_box(slide, x, y, w, h, fill, line_color=None, line_width=1.0):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line_color is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line_color
        box.line.width = Pt(line_width)
    box.text_frame.margin_left = Inches(0.15)
    box.text_frame.margin_right = Inches(0.15)
    box.text_frame.margin_top = Inches(0.10)
    box.text_frame.margin_bottom = Inches(0.10)
    return box

def set_box_text(box, lines):
    """lines = list of dicts: {text, size, bold, color, italic, align, space_before, space_after}"""
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        if "space_before" in line: p.space_before = Pt(line["space_before"])
        if "space_after"  in line: p.space_after  = Pt(line["space_after"])
        run = p.add_run()
        run.text = line["text"]
        run.font.name = line.get("font", "Arial")
        run.font.size = Pt(line.get("size", 16))
        run.font.bold = line.get("bold", False)
        run.font.italic = line.get("italic", False)
        run.font.color.rgb = line.get("color", DARK)

# ---------------------------------------------------------------------------
# Load summary stats from R script
# ---------------------------------------------------------------------------

def load_summary():
    if SUMMARY_JSON.exists():
        with open(SUMMARY_JSON, encoding="utf-8") as f:
            return json.load(f)
    print("  WARN: talk_figures_summary.json not found; using fallback values")
    return {
        "N": 403,
        "rq1": {"delta_male": 14.98, "delta_female": 4.32, "d_interaction": 0.339,
                "b_interaction": -11.60, "p_interaction": 0.0622},
        "rq2_strong": {"delta_male": 16.77, "delta_female": 5.95, "d_interaction": 0.361,
                       "b_interaction": -13.01, "p_interaction": 0.1125, "n": 211},
        "rq2_weak":   {"delta_male": 12.59, "delta_female": 3.41, "d_interaction": 0.279,
                       "b_interaction": -9.50, "p_interaction": 0.3216, "n": 192},
        "rq3":        {"gender_main_b": -1.06, "gender_main_p": 0.81},
    }

# ===========================================================================
# SLIDE 1: COVER (Maurice title restored + UChicago date)
# ===========================================================================

def slide_01_cover(prs, summary, slide_num):
    """Maurice 'Who Builds Social Capital...' cover, with date/venue updated.

    We rebuild it natively in python-pptx (rather than embedding the Maurice
    PNG and masking the date) so the title text is selectable and the slide
    matches the rest of the deck's vector look. Layout faithfully echoes the
    Maurice cover: navy background, Wharton-styled banner, big white title,
    author line, italic gray date.
    """
    slide = add_blank_slide(prs)
    set_slide_background(slide, NAVY)

    # Decorative red diagonal stripe on the right (echoes Maurice template)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(11.6), Inches(-0.5), Inches(2.6), Inches(8.5)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(0x99, 0x12, 0x12)
    stripe.line.fill.background()
    stripe.adjustments[0] = 0.55  # angle

    # Wharton logo (extracted from Maurice cover slide as a navy-bg PNG)
    if WHARTON_LOGO_PNG.exists():
        slide.shapes.add_picture(
            str(WHARTON_LOGO_PNG),
            Inches(0.65), Inches(0.45),
            width=Inches(3.4), height=Inches(0.96)
        )
    else:
        # Fallback: text-only banner
        add_textbox(
            slide, Inches(0.85), Inches(0.55), Inches(8), Inches(0.7),
            "WHARTON", size=34, bold=True, color=WHITE, font="Arial"
        )
        add_textbox(
            slide, Inches(0.85), Inches(1.10), Inches(8), Inches(0.4),
            "UNIVERSITY of PENNSYLVANIA", size=14, bold=False, color=WHITE,
            italic=True, font="Arial"
        )

    # Big title — restored from Maurice
    add_textbox(
        slide, Inches(0.85), Inches(2.85), Inches(11.5), Inches(1.9),
        "Who Builds Social Capital",
        size=54, bold=True, color=WHITE, font="Arial"
    )
    add_textbox(
        slide, Inches(0.85), Inches(3.75), Inches(11.5), Inches(1.0),
        "Through Sponsorship?",
        size=54, bold=True, color=WHITE, font="Arial"
    )

    # Authors
    add_textbox(
        slide, Inches(0.85), Inches(4.95), Inches(11.5), Inches(0.6),
        "Jose Cervantez  ·  Erika Kirgios  ·  Rosalind Chow",
        size=24, color=WHITE, font="Arial"
    )

    # Venue and date — updated for UChicago
    add_textbox(
        slide, Inches(0.85), Inches(6.05), Inches(11.5), Inches(0.5),
        "University of Chicago",
        size=20, italic=True, color=RGBColor(0xCC, 0xD2, 0xE0), font="Arial"
    )
    add_textbox(
        slide, Inches(0.85), Inches(6.55), Inches(11.5), Inches(0.45),
        "April 10, 2026",
        size=18, italic=True, color=RGBColor(0xCC, 0xD2, 0xE0), font="Arial"
    )
    # Cover slide: no slide number
    return slide

# ===========================================================================
# SLIDE 2: HOOK (outcome-focused composite)
# ===========================================================================

def slide_02_hook(prs, summary, slide_num):
    """Editable native-python-pptx hook.

    Two manager cards (Mark Davies, Sarah Lin) with MATCHED 4/2/2 promotion
    track records — same scoreboard, different sponsor. Setup line + open
    question below + forecast about 'a variable neither manager chose.'

    Built entirely from python-pptx shapes (no embedded PNG) so the user
    can tweak names, titles, track-record numbers, and colors directly in
    PowerPoint.
    """
    slide = add_blank_slide(prs)
    set_slide_background(slide, BG)

    BLUE_AVATAR = RGBColor(0x3B, 0x82, 0xF6)
    PINK_AVATAR = RGBColor(0xEC, 0x48, 0x99)
    GREEN_OK    = RGBColor(0x04, 0x7A, 0x4B)
    RED_BAD     = RGBColor(0xC8, 0x1E, 0x1E)
    DIVIDER     = RGBColor(0xE5, 0xE7, 0xEB)

    # Top tagline
    add_textbox(
        slide, Inches(0.55), Inches(0.45),
        Inches(12.3), Inches(0.4),
        "AN EXAMPLE TO SET UP TODAY'S QUESTION",
        size=14, bold=True, color=GRAY,
        align=PP_ALIGN.CENTER, font="Arial"
    )

    # ---- Two manager cards ----
    card_w = Inches(5.55)
    card_h = Inches(4.05)
    card_gap = Inches(0.55)
    total_w = 2 * card_w + card_gap
    start_x = (SLIDE_W - total_w) / 2
    card_y = Inches(1.05)

    cards = [
        {"name": "Mark Davies",
         "title": "Senior Manager, Acme Corp",
         "avatar": BLUE_AVATAR,
         "initial": "M",
         "x": start_x},
        {"name": "Sarah Lin",
         "title": "Senior Manager, Acme Corp",
         "avatar": PINK_AVATAR,
         "initial": "S",
         "x": start_x + card_w + card_gap},
    ]

    for c in cards:
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            c["x"], card_y, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = DIVIDER
        card.line.width = Pt(1.25)
        # Don't use the card's text frame for layout; we lay shapes on top
        card.text_frame.text = ""

        # Avatar (colored circle with white initial)
        avatar_size = Inches(1.05)
        avatar_x = c["x"] + (card_w - avatar_size) / 2
        avatar_y = card_y + Inches(0.30)
        avatar = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            avatar_x, avatar_y, avatar_size, avatar_size
        )
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = c["avatar"]
        avatar.line.fill.background()
        atf = avatar.text_frame
        atf.margin_left = Emu(0); atf.margin_right = Emu(0)
        atf.margin_top  = Emu(0); atf.margin_bottom = Emu(0)
        atf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = c["initial"]
        ar.font.name = "Arial"
        ar.font.size = Pt(38)
        ar.font.bold = True
        ar.font.color.rgb = WHITE

        # Name
        add_textbox(
            slide, c["x"], avatar_y + avatar_size + Inches(0.10),
            card_w, Inches(0.5),
            c["name"], size=24, bold=True, color=DARK,
            align=PP_ALIGN.CENTER
        )
        # Title
        add_textbox(
            slide, c["x"], avatar_y + avatar_size + Inches(0.55),
            card_w, Inches(0.4),
            c["title"], size=14, italic=True, color=GRAY,
            align=PP_ALIGN.CENTER
        )

        # Thin divider line under the title
        div_y = avatar_y + avatar_size + Inches(1.00)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            c["x"] + Inches(0.55), div_y,
            card_w - Inches(1.10), Inches(0.012)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = DIVIDER
        divider.line.fill.background()

        # Track-record header
        stats_y = div_y + Inches(0.20)
        add_textbox(
            slide, c["x"] + Inches(0.5), stats_y,
            card_w - Inches(1.0), Inches(0.42),
            "Past promotions endorsed",
            size=14, bold=True, color=GRAY,
            align=PP_ALIGN.CENTER, font="Arial"
        )
        # Big "4" right under the header
        add_textbox(
            slide, c["x"] + Inches(0.5), stats_y + Inches(0.32),
            card_w - Inches(1.0), Inches(0.55),
            "4", size=34, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER
        )

        # Outcome rows (✓ Worked out: 2  ·  ✗ Didn't pan out: 2)
        row_y = stats_y + Inches(0.95)
        # Use two side-by-side textboxes so labels and counts line up
        col_w = (card_w - Inches(0.7)) / 2
        # Worked out
        ok_box = slide.shapes.add_textbox(
            c["x"] + Inches(0.35), row_y, col_w, Inches(0.45)
        )
        otf = ok_box.text_frame
        otf.margin_left = Emu(0); otf.margin_right = Emu(0)
        otf.margin_top  = Emu(0); otf.margin_bottom = Emu(0)
        op = otf.paragraphs[0]
        op.alignment = PP_ALIGN.CENTER
        r1 = op.add_run(); r1.text = "✓  "
        r1.font.name = "Arial"; r1.font.size = Pt(20); r1.font.bold = True
        r1.font.color.rgb = GREEN_OK
        r2 = op.add_run(); r2.text = "Worked out"
        r2.font.name = "Arial"; r2.font.size = Pt(15); r2.font.color.rgb = DARK
        op2 = otf.add_paragraph()
        op2.alignment = PP_ALIGN.CENTER
        op2.space_before = Pt(2)
        r3 = op2.add_run(); r3.text = "2"
        r3.font.name = "Arial"; r3.font.size = Pt(22); r3.font.bold = True
        r3.font.color.rgb = GREEN_OK

        # Didn't pan out
        bad_box = slide.shapes.add_textbox(
            c["x"] + Inches(0.35) + col_w, row_y, col_w, Inches(0.45)
        )
        btf = bad_box.text_frame
        btf.margin_left = Emu(0); btf.margin_right = Emu(0)
        btf.margin_top  = Emu(0); btf.margin_bottom = Emu(0)
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        r4 = bp.add_run(); r4.text = "✗  "
        r4.font.name = "Arial"; r4.font.size = Pt(20); r4.font.bold = True
        r4.font.color.rgb = RED_BAD
        r5 = bp.add_run(); r5.text = "Didn't pan out"
        r5.font.name = "Arial"; r5.font.size = Pt(15); r5.font.color.rgb = DARK
        bp2 = btf.add_paragraph()
        bp2.alignment = PP_ALIGN.CENTER
        bp2.space_before = Pt(2)
        r6 = bp2.add_run(); r6.text = "2"
        r6.font.name = "Arial"; r6.font.size = Pt(22); r6.font.bold = True
        r6.font.color.rgb = RED_BAD

    # ---- Setup line below cards ----
    add_textbox(
        slide, Inches(0.55), card_y + card_h + Inches(0.20),
        Inches(12.3), Inches(0.4),
        "Today, each one brings you a new recommendation for promotion.",
        size=18, italic=True, color=DARK, align=PP_ALIGN.CENTER
    )

    # ---- Big question ----
    add_textbox(
        slide, Inches(0.55), card_y + card_h + Inches(0.70),
        Inches(12.3), Inches(0.65),
        "How do you weigh each one?",
        size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER
    )

    # ---- Forecast (small italic) ----
    add_textbox(
        slide, Inches(0.55), card_y + card_h + Inches(1.40),
        Inches(12.3), Inches(0.4),
        "Today: even with identical track records, audiences treat one of these "
        "recommendations as a meaningfully stronger signal than the other.",
        size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER
    )

    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 3: WHAT IS SPONSORSHIP?
# ===========================================================================

def slide_03_what_is_sponsorship(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "What is sponsorship?")

    add_textbox(
        slide, Inches(0.55), Inches(1.45), Inches(12.3), Inches(1.6),
        "Sponsorship is support provided to influence other people's perceptions "
        "of a protégé's skills, abilities, performance, and potential — with the "
        "intention of helping the protégé in their career.",
        size=22, color=DARK, font="Arial"
    )

    # Use a left-aligned bullet list with extra spacing so nothing clips
    add_bullets(slide, Inches(0.95), Inches(3.6), Inches(12), Inches(3.2), [
        "Public advocacy and recommendations",
        "Endorsements to colleagues, clients, hiring committees",
        "Introductions to high-status others",
        "Vouching for skill, judgment, or fit",
    ], size=22, text_color=DARK, spacing=14)
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 4: TRADEOFFS (Maurice-style two-box)
# ===========================================================================

def slide_04_tradeoffs(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Sponsorship involves real risk for sponsors")

    add_textbox(
        slide, Inches(0.55), Inches(1.4), Inches(12.3), Inches(0.7),
        "When a sponsor stakes their reputation on a protégé:",
        size=20, color=GRAY, font="Arial", italic=True
    )

    col_w = Inches(5.7)
    col_h = Inches(3.65)

    # Success box (left)
    left_box = add_rounded_box(
        slide, Inches(0.65), Inches(2.35), col_w, col_h,
        fill=RGBColor(0xE0, 0xF2, 0xFE), line_color=NAVY, line_width=1.5
    )
    set_box_text(left_box, [
        {"text": "When the protégé succeeds", "size": 22, "bold": True,
         "color": NAVY, "align": PP_ALIGN.CENTER, "space_after": 12},
        {"text": " ", "size": 6},  # spacer
        {"text": "→  Sponsor gains credibility", "size": 18, "color": DARK,
         "align": PP_ALIGN.CENTER, "space_after": 8},
        {"text": "→  Reputation lifts with the protégé", "size": 18, "color": DARK,
         "align": PP_ALIGN.CENTER, "space_after": 8},
        {"text": "→  Network ties strengthen", "size": 18, "color": DARK,
         "align": PP_ALIGN.CENTER},
    ])

    # Failure box (right)
    right_box = add_rounded_box(
        slide, Inches(7.0), Inches(2.35), col_w, col_h,
        fill=RGBColor(0xFD, 0xE2, 0xE2), line_color=RED, line_width=1.5
    )
    set_box_text(right_box, [
        {"text": "When the protégé fails", "size": 22, "bold": True,
         "color": RED, "align": PP_ALIGN.CENTER, "space_after": 12},
        {"text": " ", "size": 6},
        {"text": "→  Sponsor loses credibility", "size": 18, "color": DARK,
         "align": PP_ALIGN.CENTER, "space_after": 8},
        {"text": "→  Future endorsements are discounted", "size": 18, "color": DARK,
         "align": PP_ALIGN.CENTER, "space_after": 8},
        {"text": "→  Network ties may weaken", "size": 18, "color": DARK,
         "align": PP_ALIGN.CENTER},
    ])

    add_textbox(
        slide, Inches(0.55), Inches(6.45), Inches(12.3), Inches(0.55),
        "But are these risks symmetric across the gender of the sponsor?",
        size=20, italic=True, bold=True, color=NAVY, align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 5: RESEARCH QUESTIONS
# ===========================================================================

def slide_05_research_questions(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Research questions")

    rqs = [
        ("RQ1",
         "Does endorser gender shape how outcomes update trust?",
         "After observing one success or failure, do audiences treat male and female sponsors symmetrically?"),
        ("RQ2",
         "Does endorsement strength moderate the bias?",
         "Are confident vs. hedged endorsements treated differently as a function of sponsor gender?"),
        ("RQ3",
         "Does the bias exist before any outcome is observed?",
         "Or does it emerge only after evaluators see the protégé succeed or fail?"),
    ]
    y = Inches(1.5)
    for label, q, sub in rqs:
        row = add_rounded_box(
            slide, Inches(0.55), y, Inches(12.25), Inches(1.55),
            fill=RGBColor(0xF5, 0xF7, 0xFA), line_color=RGBColor(0xCB, 0xD5, 0xE1),
            line_width=0.75
        )
        row.text_frame.margin_left = Inches(0.35)
        row.text_frame.margin_top  = Inches(0.22)

        p = row.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = label + "    "
        run.font.name = "Arial"; run.font.size = Pt(22); run.font.bold = True
        run.font.color.rgb = RED
        run2 = p.add_run()
        run2.text = q
        run2.font.name = "Arial"; run2.font.size = Pt(22); run2.font.bold = True
        run2.font.color.rgb = NAVY

        p2 = row.text_frame.add_paragraph()
        p2.space_before = Pt(8)
        run3 = p2.add_run()
        run3.text = sub
        run3.font.name = "Arial"; run3.font.size = Pt(15); run3.font.italic = True
        run3.font.color.rgb = GRAY

        y += Inches(1.75)
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 6: MULTI-STAGE OVERVIEW
# ===========================================================================

def slide_06_multistage(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Multi-stage lab experiment")

    stages = [
        ("Stage 1", "Protégé performance",
         ["Real participants take three",
          "performance tasks (Barron et al., 2025).",
          "We use the verifiable scores to",
          "build candidate profiles."]),
        ("Stage 2", "Sponsor endorsements",
         ["Real sponsors view candidate",
          "pairs and rate which one will",
          "perform better, and how confident",
          "they are."]),
        ("Stage 3", "Audience evaluation",
         ["Evaluators bet on a sponsor's",
          "endorsement, see the outcome,",
          "then bet on a second endorsement",
          "from the same sponsor."]),
    ]
    col_w = Inches(3.95)
    col_h = Inches(4.85)
    gap   = Inches(0.32)
    total_w = 3 * col_w + 2 * gap
    start_x = (SLIDE_W - total_w) / 2

    for i, (label, title, lines) in enumerate(stages):
        x = start_x + i * (col_w + gap)
        box = add_rounded_box(
            slide, x, Inches(1.55), col_w, col_h,
            fill=WHITE, line_color=NAVY, line_width=2.0
        )
        box.text_frame.margin_left  = Inches(0.28)
        box.text_frame.margin_right = Inches(0.28)
        box.text_frame.margin_top   = Inches(0.30)

        p1 = box.text_frame.paragraphs[0]
        r1 = p1.add_run(); r1.text = label
        r1.font.name = "Arial"; r1.font.size = Pt(20); r1.font.bold = True
        r1.font.color.rgb = RED
        p1.alignment = PP_ALIGN.CENTER

        p2 = box.text_frame.add_paragraph()
        p2.space_before = Pt(4); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = title
        r2.font.name = "Arial"; r2.font.size = Pt(22); r2.font.bold = True
        r2.font.color.rgb = NAVY

        for j, line in enumerate(lines):
            pp = box.text_frame.add_paragraph()
            pp.space_before = Pt(6 if j == 0 else 0)
            pp.alignment = PP_ALIGN.LEFT
            if j == 0:
                pp.space_before = Pt(18)
            r = pp.add_run(); r.text = line
            r.font.name = "Arial"; r.font.size = Pt(15)
            r.font.color.rgb = DARK

    add_textbox(
        slide, Inches(0.55), Inches(6.65), Inches(12.3), Inches(0.5),
        "Today's talk focuses on Stage 3.",
        size=16, italic=True, bold=True, color=NAVY, align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 7: STAGE 1 (Maurice screenshot of the actual tasks)
# ===========================================================================

def slide_07_stage1_tasks(prs, summary, slide_num):
    """Embed the Maurice slide-8 task screenshots (German word search,
    knowledge quiz, logical reasoning matrices) so the audience sees the
    actual stimuli protégés worked on."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Stage 1:  Where do the candidates come from?")

    add_textbox(
        slide, Inches(0.55), Inches(1.15), Inches(12.3), Inches(0.6),
        "Real participants completed three professional-ability tasks "
        "(Barron et al., 2025).",
        size=18, italic=True, color=GRAY
    )

    if STAGE1_TASKS_PNG.exists():
        add_picture_centered(
            slide, STAGE1_TASKS_PNG, top=Inches(1.85),
            max_width=Inches(12.5), max_height=Inches(4.95)
        )
    else:
        add_textbox(slide, Inches(3), Inches(3.5), Inches(7), Inches(1),
                    "[Stage 1 tasks image missing]", size=20, color=RED,
                    align=PP_ALIGN.CENTER)

    add_textbox(
        slide, Inches(0.55), Inches(6.75), Inches(12.3), Inches(0.4),
        "Word search  ·  General-knowledge quiz  ·  Logical-reasoning (Raven's matrices)",
        size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 8: STAGE 2 (with embedded candidate-pair screen)
# ===========================================================================

def slide_08_stage2(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Stage 2:  Real sponsors made real endorsements")

    add_textbox(
        slide, Inches(0.55), Inches(1.15), Inches(12.3), Inches(0.6),
        "Stage 2 was a stimulus-generation step. Sponsors saw real candidate "
        "pairs and picked who would perform better.",
        size=17, italic=True, color=GRAY
    )

    # Left column: small bullet description
    add_bullets(slide, Inches(0.55), Inches(2.05), Inches(4.85), Inches(4.8), [
        ("Real sponsors  →  real endorsements", True),
        "Each sponsor sees one candidate pair from Stage 1",
        "Sponsor moves a slider to express which candidate they favor and how confident they are",
        "Their sliders become the endorsements that Stage 3 evaluators see",
        "We sample sponsors across gender to populate every Stage 3 cell",
    ], size=15, spacing=10)

    # Right column: embed the actual candidate-pair screen
    if CANDIDATE_SCREEN_PNG.exists():
        from PIL import Image
        img = Image.open(CANDIDATE_SCREEN_PNG)
        iw, ih = img.size
        aspect = iw / ih
        max_w = Inches(7.5)
        max_h = Inches(4.7)
        if max_w / aspect <= max_h:
            w, h = max_w, int(max_w / aspect)
        else:
            h, w = max_h, int(max_h * aspect)
        slide.shapes.add_picture(
            str(CANDIDATE_SCREEN_PNG),
            Inches(5.55), Inches(2.05), width=w, height=h
        )
    else:
        add_textbox(slide, Inches(6), Inches(3.5), Inches(6), Inches(1),
                    "[candidate screen missing]", size=18, color=RED,
                    align=PP_ALIGN.CENTER)

    add_textbox(
        slide, Inches(5.55), Inches(6.85), Inches(7.5), Inches(0.4),
        "An example of what a Stage 2 sponsor sees.",
        size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 9: STAGE 3 PROCESS DIAGRAM
# ===========================================================================
# This replaces the old "Gender × Outcome × Strength = 8 cells" grid that used
# coded variable names like W_correct_strong. The new diagram explains the
# three manipulations and the evaluator flow in plain English.

def slide_09_stage3_design(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Stage 3:  What we manipulate, what evaluators do")

    add_textbox(
        slide, Inches(0.55), Inches(1.15), Inches(12.3), Inches(0.55),
        "Each evaluator sees ONE sponsor making ONE endorsement, then "
        "watches the outcome, then sees a second endorsement from the same sponsor.",
        size=15, italic=True, color=GRAY
    )

    # ---- Top half: 3 "we manipulate" cards ----
    manip_top = Inches(2.0)
    manip_h = Inches(2.0)
    card_w = Inches(4.05)
    card_gap = Inches(0.16)
    total_card_w = 3 * card_w + 2 * card_gap
    start_x = (SLIDE_W - total_card_w) / 2

    # Header strip above the 3 cards
    add_textbox(
        slide, Inches(0.55), Inches(1.75), Inches(12.3), Inches(0.35),
        "WE RANDOMLY VARY THREE THINGS",
        size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER
    )

    manip_cards = [
        {
            "tag": "1.  Sponsor gender",
            "lines": [
                ("Man", NAVY),
                ("Woman", RED),
            ],
            "border": NAVY,
        },
        {
            "tag": "2.  Endorsement strength",
            "lines": [
                ("Strong  (very confident)", DARK),
                ("Weak  (hedged / unsure)",  DARK),
            ],
            "border": NAVY,
        },
        {
            "tag": "3.  Outcome",
            "lines": [
                ("Endorser was correct", GREEN_DK),
                ("Endorser was wrong",   RED_LT),
            ],
            "border": NAVY,
        },
    ]

    for i, m in enumerate(manip_cards):
        x = start_x + i * (card_w + card_gap)
        box = add_rounded_box(
            slide, x, manip_top, card_w, manip_h,
            fill=WHITE, line_color=m["border"], line_width=2.0
        )
        box.text_frame.margin_top = Inches(0.18)
        box.text_frame.margin_left = Inches(0.25)
        # Tag header
        p1 = box.text_frame.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = m["tag"]
        r1.font.name = "Arial"; r1.font.size = Pt(18); r1.font.bold = True
        r1.font.color.rgb = NAVY
        # Two value lines
        for j, (line, color) in enumerate(m["lines"]):
            pp = box.text_frame.add_paragraph()
            pp.alignment = PP_ALIGN.LEFT
            pp.space_before = Pt(10 if j == 0 else 4)
            r = pp.add_run(); r.text = "•  " + line
            r.font.name = "Arial"; r.font.size = Pt(16); r.font.bold = True
            r.font.color.rgb = color

    # ---- Bottom half: 3-step evaluator flow ----
    add_textbox(
        slide, Inches(0.55), Inches(4.35), Inches(12.3), Inches(0.35),
        "WHAT EACH EVALUATOR DOES",
        size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER
    )

    flow_top = Inches(4.75)
    step_w = Inches(3.55)
    step_h = Inches(1.85)
    arrow_w = Inches(0.55)
    total_flow_w = 3 * step_w + 2 * arrow_w
    flow_x = (SLIDE_W - total_flow_w) / 2

    steps = [
        {
            "num": "1",
            "title": "Initial wager",
            "body": "Sees the sponsor's first endorsement; bets a fraction of a $0.50 bank that the sponsor is correct.",
            "fill": RGBColor(0xE7, 0xEF, 0xFA),
            "border": NAVY,
        },
        {
            "num": "2",
            "title": "Outcome reveal",
            "body": "Learns whether the sponsor was correct. The bank is updated.",
            "fill": RGBColor(0xFD, 0xF1, 0xCC),
            "border": RGBColor(0xC9, 0xA2, 0x27),
        },
        {
            "num": "3",
            "title": "Second wager",
            "body": "Sees the SAME sponsor's second endorsement (different candidate pair); bets again.",
            "fill": RGBColor(0xE7, 0xEF, 0xFA),
            "border": NAVY,
        },
    ]
    for i, s in enumerate(steps):
        x = flow_x + i * (step_w + arrow_w)
        box = add_rounded_box(
            slide, x, flow_top, step_w, step_h,
            fill=s["fill"], line_color=s["border"], line_width=2.0
        )
        box.text_frame.margin_top  = Inches(0.18)
        box.text_frame.margin_left = Inches(0.25)
        # Step badge + title
        p = box.text_frame.paragraphs[0]
        r = p.add_run(); r.text = s["num"] + "    "
        r.font.name = "Arial"; r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = RED
        r2 = p.add_run(); r2.text = s["title"]
        r2.font.name = "Arial"; r2.font.size = Pt(20); r2.font.bold = True
        r2.font.color.rgb = NAVY
        # Body
        p2 = box.text_frame.add_paragraph()
        p2.space_before = Pt(8)
        rb = p2.add_run(); rb.text = s["body"]
        rb.font.name = "Arial"; rb.font.size = Pt(13)
        rb.font.color.rgb = DARK

        # Arrow connector to next step
        if i < len(steps) - 1:
            arrow_x = flow_x + (i + 1) * step_w + i * arrow_w
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x + Inches(0.05),
                flow_top + step_h / 2 - Inches(0.22),
                arrow_w - Inches(0.10),
                Inches(0.44)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = NAVY
            arrow.line.fill.background()

    # Outcome variable callout
    add_textbox(
        slide, Inches(0.55), Inches(6.85), Inches(12.3), Inches(0.45),
        "Outcome variable:  Trust change  =  (second wager)  −  (initial wager)",
        size=15, italic=True, bold=True, color=NAVY, align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDES 10–12: WALKTHROUGH (D1, D2, D3) — bigger images, no caption
# ===========================================================================

def slide_walkthrough(prs, png_path, label, sublabel, slide_num):
    """Walkthrough slide: tight title, big image, no footer caption."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, label, size=28, top=0.30)
    add_textbox(
        slide, Inches(0.55), Inches(0.95), Inches(12.3), Inches(0.45),
        sublabel, size=15, italic=True, color=GRAY
    )
    if png_path.exists():
        # With aspect ~2.6, max_width=12.5 → height 4.81 → bottlenecked by width
        add_picture_centered(
            slide, png_path, top=Inches(1.55),
            max_width=Inches(12.6), max_height=Inches(5.6)
        )
    else:
        add_textbox(slide, Inches(3), Inches(3.5), Inches(7), Inches(1),
                    f"[{png_path.name} missing]", size=20, color=RED,
                    align=PP_ALIGN.CENTER)
    add_slide_number(slide, slide_num)
    return slide

def slide_10_walk_d1(prs, summary, slide_num):
    return slide_walkthrough(
        prs,
        WALK_DIR / "W_correct_strong_d1.png",
        "Step 1:  Initial wager",
        "An evaluator sees one sponsor's confident endorsement and bets a fraction of their $0.50 bank.",
        slide_num,
    )

def slide_11_walk_d2(prs, summary, slide_num):
    return slide_walkthrough(
        prs,
        WALK_DIR / "W_correct_strong_d2_correct.png",
        "Step 2:  Outcome reveal",
        "The endorser was CORRECT — the evaluator's bank goes up.  (We also run a 'wrong' variant.)",
        slide_num,
    )

def slide_12_walk_d3(prs, summary, slide_num):
    return slide_walkthrough(
        prs,
        WALK_DIR / "W_correct_strong_d3.png",
        "Step 3:  Second wager",
        "The same sponsor, a different candidate pair.  How much does the evaluator wager now?",
        slide_num,
    )

# ===========================================================================
# SLIDES 13-16: RESULTS — full-bleed figures
# ===========================================================================

def slide_results_fullbleed(prs, png_path, title, slide_num, subtitle=None):
    """Results slide: small title, near-full-bleed figure underneath."""
    slide = add_blank_slide(prs)
    add_title_bar(slide, title, size=26, top=0.30)
    if subtitle:
        add_textbox(
            slide, Inches(0.55), Inches(0.92), Inches(12.3), Inches(0.42),
            subtitle, size=14, italic=True, color=GRAY
        )
        top = Inches(1.45)
    else:
        top = Inches(1.20)
    if png_path.exists():
        add_picture_centered(
            slide, png_path, top=top,
            max_width=Inches(12.7), max_height=Inches(5.85)
        )
    else:
        add_textbox(slide, Inches(3), Inches(3.5), Inches(7), Inches(1),
                    f"[{png_path.name} missing]", size=20, color=RED,
                    align=PP_ALIGN.CENTER)
    add_slide_number(slide, slide_num)
    return slide

def slide_13_rq1(prs, summary, slide_num):
    return slide_results_fullbleed(
        prs, RQ1_PNG,
        "RQ1:  Gender × Outcome on trust update",
        slide_num
    )

def slide_14_rq2_strong(prs, summary, slide_num):
    return slide_results_fullbleed(
        prs, RQ2_STRONG_PNG,
        "RQ2:  Strong endorsers",
        slide_num,
        subtitle="When sponsors are very confident, the gender × outcome gap is biggest."
    )

def slide_15_rq2_weak(prs, summary, slide_num):
    return slide_results_fullbleed(
        prs, RQ2_WEAK_PNG,
        "RQ2:  Weak endorsers",
        slide_num,
        subtitle="When sponsors hedge, the same pattern persists at smaller magnitudes."
    )

def slide_16_rq3(prs, summary, slide_num):
    return slide_results_fullbleed(
        prs, RQ3_PNG,
        "RQ3:  Does the bias exist before any outcome?",
        slide_num,
        subtitle="Before the outcome arrives, audiences trust male and female sponsors equally."
    )

# ===========================================================================
# SLIDE 17: SYNTHESIS
# ===========================================================================

def slide_17_synthesis(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Synthesis")
    rq1 = summary["rq1"]

    # Headline (left)
    add_textbox(
        slide, Inches(0.55), Inches(1.4), Inches(8.0), Inches(1.4),
        "The bias is real — but it shows up post-outcome, not pre-outcome.",
        size=24, bold=True, color=NAVY
    )

    add_bullets(slide, Inches(0.55), Inches(2.85), Inches(8.0), Inches(3.8), [
        f"Men's endorsements move trust by Δ = {rq1['delta_male']:+.1f}",
        f"Women's endorsements move trust by Δ = {rq1['delta_female']:+.1f}",
        ("Pre-outcome trust does NOT differ by gender", True),
        "Stronger endorsements amplify the asymmetry",
        f"Cohen's d for Gender × Outcome  ≈  {rq1['d_interaction']:.2f}",
    ], size=18, spacing=12)

    # Right side: synthesis mini-chart
    if SYNTH_PNG.exists():
        from PIL import Image
        img = Image.open(SYNTH_PNG)
        iw, ih = img.size
        aspect = iw / ih
        max_w = Inches(4.4)
        max_h = Inches(4.6)
        if max_w / aspect <= max_h:
            w, h = max_w, int(max_w / aspect)
        else:
            h, w = max_h, int(max_h * aspect)
        slide.shapes.add_picture(
            str(SYNTH_PNG),
            Inches(8.85), Inches(1.55), width=w, height=h
        )

    # Reframe footer
    add_textbox(
        slide, Inches(0.55), Inches(6.55), Inches(12.3), Inches(0.7),
        "Reframe: this is not 'women are punished more harshly.' "
        "It is 'men's endorsements get updated against more — in BOTH directions.'",
        size=15, italic=True, bold=True, color=RED, align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 18: THEORETICAL IMPLICATIONS
# ===========================================================================

def slide_18_implications(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Theoretical implications")
    add_bullets(slide, Inches(0.55), Inches(1.55), Inches(12.3), Inches(5.5), [
        ("Endorser informativeness is gendered.", True),
        "Audiences treat a single outcome as more diagnostic of a man's reliability than of a woman's.",
        ("This is not a 'taste-based' bias against women.", True),
        "Initial trust is symmetric.  The asymmetry arrives only after outcome feedback.",
        ("Implications for sponsorship as a career-mobility lever.", True),
        "If outcome learning is dampened for women sponsors, their reputational ROI per endorsement is structurally lower — even when they pick equally well.",
    ], size=18, spacing=14)
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 19: LIMITATIONS & NEXT STEPS
# ===========================================================================

def slide_19_limitations(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Limitations and next steps")
    add_bullets(slide, Inches(0.55), Inches(1.55), Inches(12.3), Inches(5.5), [
        "Single-outcome design — multi-trial accumulation may differ",
        "Prolific sample limits generalizability to organizational contexts",
        "Endorsement strength was rated by sponsors, not held to a single calibration",
        ("Next:", True),
        "Pre-registered confirmatory study with N ≈ 1,200 (powered for the 3-way interaction)",
        "Field replication: real LinkedIn-style endorsements from professional reviewers",
        "Mechanism: ask evaluators to attribute success/failure to skill vs. luck and test mediation",
    ], size=17, spacing=12)
    add_slide_number(slide, slide_num)
    return slide

# ===========================================================================
# SLIDE 20: THANKS
# ===========================================================================

def slide_20_thanks(prs, summary, slide_num):
    slide = add_blank_slide(prs)
    set_slide_background(slide, BG)
    add_textbox(
        slide, Inches(0.6), Inches(2.4), Inches(12), Inches(1.0),
        "Thank you", size=64, bold=True, color=NAVY, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.6),
        "Jose Cervantez · Erika Kirgios · Rosalind Chow",
        size=24, color=DARK, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
        "University of Chicago  ·  April 10, 2026",
        size=18, italic=True, color=GRAY, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
        "Questions?", size=28, italic=True, bold=True, color=RED,
        align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, slide_num)
    return slide

# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------

BUILDERS = [
    slide_01_cover,                # 1
    slide_02_hook,                 # 2
    slide_03_what_is_sponsorship,  # 3
    slide_04_tradeoffs,            # 4
    slide_05_research_questions,   # 5
    slide_06_multistage,           # 6
    slide_07_stage1_tasks,         # 7
    slide_08_stage2,               # 8
    slide_09_stage3_design,        # 9
    slide_10_walk_d1,              # 10
    slide_11_walk_d2,              # 11
    slide_12_walk_d3,              # 12
    slide_13_rq1,                  # 13
    slide_14_rq2_strong,           # 14
    slide_15_rq2_weak,             # 15
    slide_16_rq3,                  # 16
    slide_17_synthesis,            # 17
    slide_18_implications,         # 18
    slide_19_limitations,          # 19
    slide_20_thanks,               # 20
]

TOTAL_SLIDES = len(BUILDERS)

def main():
    print("=" * 72)
    print(f"  Building UChicago talk deck  ({TOTAL_SLIDES} slides)")
    print("=" * 72)
    print(f"  Base:   {BASE_PPTX}")
    print(f"  Output: {OUTPUT_PPTX}")

    summary = load_summary()
    print(f"  Summary: N={summary['N']}, "
          f"Δ_M={summary['rq1']['delta_male']:.2f}, "
          f"Δ_F={summary['rq1']['delta_female']:.2f}, "
          f"p_int={summary['rq1']['p_interaction']:.4f}")

    prs = Presentation(str(BASE_PPTX))
    print(f"  Loaded base ({len(prs.slides)} slides)")

    delete_all_slides(prs)
    print(f"  Cleared all base slides")

    for i, builder in enumerate(BUILDERS, 1):
        builder(prs, summary, i)
        print(f"  [{i:2d}/{TOTAL_SLIDES}] {builder.__name__}")

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"\n  Saved: {OUTPUT_PPTX}")

    # ---- Verification ----
    prs2 = Presentation(str(OUTPUT_PPTX))
    nslides = len(prs2.slides)
    print(f"\n  Verification: {nslides} slides")
    assert nslides == TOTAL_SLIDES, f"expected {TOTAL_SLIDES} slides, got {nslides}"
    pic_count = 0
    for i, slide in enumerate(prs2.slides, 1):
        n_shapes = len(slide.shapes)
        n_pics = sum(1 for s in slide.shapes if s.shape_type == 13)  # PICTURE
        pic_count += n_pics
        print(f"    slide {i:2d}: {n_shapes} shapes, {n_pics} pictures")
    print(f"\n  Total pictures embedded: {pic_count}")
    print("\n" + "=" * 72)
    print("  DONE")
    print("=" * 72)

if __name__ == "__main__":
    main()
