"""
patch_2026_04_08_polish_v2.py — second polish pass.

Items addressed:
  1. Slide 22 — strength badge caption "hedged / unsure" -> "unsure".
  2. Slide 12 (Stage 2 incentive) — full rebuild parallel to Stage 3:
     bank-style lead, green/red right/wrong cards, slider, twin reveal,
     3-col table. Drops the stale "Name of Initiative" footer.
  3. Slides 19 & 20 — descriptive RQ titles in plain English (kept the
     RQ1/RQ2 prefix for organizational continuity with slide 8).
  4. Slide 15 (Stage 3 randomization) — rebuild the 3 variant cards
     with native mini-stimuli below each: pink/blue avatars (gender),
     navy/red strength bars (strength), green/red outcome pills.
  5. Slide 14 (Stage 3 setup) — full rebuild adding sample size, an
     evaluator-eye-view bullet list, and the SAME-sponsor / two-decision
     framing. Stale "Name of Initiative" footer dropped.

Hooks: this calls clear_non_placeholder_shapes on slides 12, 14, 15.
After it runs, call fix_fld_guids.py + fix_orphaned_timing.py.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"


# ---------------------------------------------------------------------------
# Palette (matches patch_2026_04_08_polish.py)
# ---------------------------------------------------------------------------

NAVY         = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP    = RGBColor(0x00, 0x14, 0x3D)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT   = RGBColor(0x1F, 0x1F, 0x1F)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY     = RGBColor(0x5B, 0x65, 0x7A)
LIGHT_GRAY   = RGBColor(0x7C, 0x86, 0x99)
TITLE_GRAY   = RGBColor(0x4B, 0x55, 0x63)
LIGHT_RULE   = RGBColor(0x9A, 0xA7, 0xBD)
CARD_OUTLINE = RGBColor(0xC7, 0xCF, 0xDD)
BAR_FILL     = RGBColor(0xEA, 0xF0, 0xFB)
BAR_TRACK    = RGBColor(0xE5, 0xEA, 0xF3)
PILL_GREEN   = RGBColor(0xDF, 0xF4, 0xE8)
PILL_RED     = RGBColor(0xFD, 0xE2, 0xE2)
PILL_TEXT_G  = RGBColor(0x11, 0x6B, 0x45)
PILL_TEXT_R  = RGBColor(0xB1, 0x21, 0x21)
TABLE_HILITE = RGBColor(0xDB, 0xE3, 0xF1)
RIGHT_GREEN  = RGBColor(0x10, 0x6B, 0x45)
WRONG_RED    = RGBColor(0xB1, 0x21, 0x21)
WOMAN_PINK   = RGBColor(0xEC, 0x48, 0x99)
WOMAN_BG     = RGBColor(0xFC, 0xE7, 0xF3)
MAN_BLUE     = RGBColor(0x3B, 0x82, 0xF6)
MAN_BG       = RGBColor(0xDB, 0xEA, 0xFE)

FONT = "Calibri"


# ---------------------------------------------------------------------------
# Shape helpers (copied from patch_2026_04_08_polish.py for self-containment)
# ---------------------------------------------------------------------------

def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size,
                bold=False, italic=False, color=BLACK_TEXT,
                align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color)
    return box


def add_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
             rounded=False, corner=0.5):
    """Add a rectangle. If rounded=True, use ROUNDED_RECTANGLE with the
    given corner adjustment (0-1). corner=0.5 produces a lozenge/pill;
    corner~0.05-0.10 produces a softly-rounded panel."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, in_(x), in_(y), in_(w), in_(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    if rounded:
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


def add_oval(slide, x, y, w, h, fill, outline=None, outline_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, in_(x), in_(y), in_(w), in_(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    shape.text_frame.text = ""
    return shape


def clear_non_placeholder_shapes(slide):
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        spTree.remove(shape._element)
        removed += 1
    return removed


def remove_name_of_initiative_placeholder(slide):
    """Drop the stale layout 'Name of Initiative' footer placeholder, keep
    the slide-number placeholder (idx=11)."""
    spTree = slide.shapes._spTree
    removed = 0
    for ph in list(slide.placeholders):
        try:
            idx = ph.placeholder_format.idx
        except Exception:
            continue
        if idx == 11:
            continue
        txt = ph.text_frame.text.strip() if ph.has_text_frame else ""
        if txt == "Name of Initiative" or idx == 12:
            spTree.remove(ph._element)
            removed += 1
    return removed


# ---------------------------------------------------------------------------
# Phase 1: Slide 22 — "hedged / unsure" -> "unsure"
# ---------------------------------------------------------------------------

def patch_slide22_badge_caption(slide):
    fixed = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == "hedged / unsure":
            for para in shape.text_frame.paragraphs:
                if para.runs:
                    para.runs[0].text = "unsure"
                    for r in para.runs[1:]:
                        r.text = ""
            fixed += 1
    print(f"  slide 22: {fixed} caption(s) updated")


# ---------------------------------------------------------------------------
# Phase 2: Slide 12 — Stage 2 incentive parallel rebuild
# ---------------------------------------------------------------------------

def rebuild_slide12_stage2_incentive(slide):
    """Rewrite Stage 2 incentive in the same shape language as Stage 3.
    Lead with the per-endorsement maximum, parallel right/wrong cards,
    slider + twin reveal + 3-col table. Drops 'Name of Initiative'.
    """
    n_shapes = clear_non_placeholder_shapes(slide)
    n_phs = remove_name_of_initiative_placeholder(slide)
    print(f"  slide 12: cleared {n_shapes} shapes + {n_phs} stale placeholder(s)")

    # Title
    add_textbox(slide, 0.55, 0.30, 12.20, 0.70,
                "Stage 2 incentive",
                size=32, bold=True, color=NAVY_DEEP)

    # Lead — bank-style headline
    lead_box = slide.shapes.add_textbox(in_(0.55), in_(1.18), in_(12.20), in_(0.55))
    tf = lead_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "Each endorsement could earn the sponsor up to "
    style_run(r1, 26, color=DARK_GRAY)
    r2 = p.add_run()
    r2.text = "$3.00."
    style_run(r2, 26, bold=True, color=NAVY_DEEP)

    # Sub — explain the dependency
    add_textbox(slide, 0.55, 1.78, 12.20, 0.42,
                "The amount depends on how confidently a downstream evaluator wagers on it — and only if the endorsed candidate truly scored higher.",
                size=15, italic=True, color=MID_GRAY)

    # Two outcome cards (right / wrong) — parallel to Stage 3
    card_y = 2.30
    card_h = 0.92
    card_w = 5.85
    gap = 0.50
    card_left_x = 0.55 + (12.20 - 2 * card_w - gap) / 2

    # RIGHT outcome
    right_card = add_rect(slide, card_left_x, card_y, card_w, card_h,
                          fill=WHITE, outline=RIGHT_GREEN, outline_width=2.0)
    rc_tf = right_card.text_frame
    rc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rc_tf.margin_left = Inches(0.20)
    rc_tf.margin_right = Inches(0.20)
    rc_tf.text = ""
    rcp = rc_tf.paragraphs[0]
    rcp.alignment = PP_ALIGN.LEFT
    r1 = rcp.add_run()
    r1.text = "If "
    style_run(r1, 21, color=DARK_GRAY)
    r2 = rcp.add_run()
    r2.text = "right "
    style_run(r2, 21, bold=True, color=RIGHT_GREEN)
    r3 = rcp.add_run()
    r3.text = "→  $3 × evaluator's wager"
    style_run(r3, 21, color=DARK_GRAY)

    # WRONG outcome
    wrong_card = add_rect(slide, card_left_x + card_w + gap, card_y, card_w, card_h,
                          fill=WHITE, outline=WRONG_RED, outline_width=2.0)
    wc_tf = wrong_card.text_frame
    wc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    wc_tf.margin_left = Inches(0.20)
    wc_tf.margin_right = Inches(0.20)
    wc_tf.text = ""
    wcp = wc_tf.paragraphs[0]
    wcp.alignment = PP_ALIGN.LEFT
    r1 = wcp.add_run()
    r1.text = "If "
    style_run(r1, 21, color=DARK_GRAY)
    r2 = wcp.add_run()
    r2.text = "wrong "
    style_run(r2, 21, bold=True, color=WRONG_RED)
    r3 = wcp.add_run()
    r3.text = "→  nothing"
    style_run(r3, 21, color=DARK_GRAY)

    # Worked example: slider showing 50% wager
    add_textbox(slide, 0.55, 3.45, 12.20, 0.32,
                "evaluator wagers 50%",
                size=13, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    track_x = 3.50
    track_y = 3.85
    track_w = 6.33
    track = add_rect(slide, track_x, track_y, track_w, 0.22,
                     fill=BAR_FILL, rounded=False)
    track.line.color.rgb = LIGHT_RULE
    track.line.width = Pt(0.75)
    fill_w = track_w / 2
    add_rect(slide, track_x, track_y, fill_w, 0.22, fill=NAVY, rounded=False)
    thumb_size = 0.42
    add_oval(slide, track_x + fill_w - thumb_size / 2,
             track_y + 0.11 - thumb_size / 2,
             thumb_size, thumb_size,
             fill=WHITE, outline=NAVY, outline_width=2.5)
    add_textbox(slide, track_x - 0.75, track_y + 0.32, 0.70, 0.30,
                "0%", size=12, color=MID_GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, track_x + track_w + 0.05, track_y + 0.32, 0.80, 0.30,
                "100%", size=12, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Twin reveal
    reveal_y = 4.55
    add_textbox(slide, 0.55, reveal_y, 5.85, 0.85,
                "right:  $1.50",
                size=32, bold=True, color=RIGHT_GREEN,
                align=PP_ALIGN.RIGHT)
    add_textbox(slide, 6.95, reveal_y, 5.85, 0.85,
                "wrong:  $0",
                size=32, bold=True, color=WRONG_RED,
                align=PP_ALIGN.LEFT)

    # Payout table — 3 cols (parallel to Stage 3)
    table_top = 5.55
    row_h = 0.34
    table_left_x = 2.40
    col_widths = [1.95, 2.55, 2.55, 2.55]
    col_lefts = []
    cx = table_left_x
    for w in col_widths:
        col_lefts.append(cx)
        cx += w

    add_textbox(slide, col_lefts[0], table_top, col_widths[0], row_h,
                "wager", size=15, bold=True, color=MID_GRAY,
                align=PP_ALIGN.RIGHT)
    for i, label in enumerate(["0%", "50%", "100%"], start=1):
        add_textbox(slide, col_lefts[i], table_top, col_widths[i], row_h,
                    label, size=15, bold=True, color=MID_GRAY,
                    align=PP_ALIGN.CENTER)
    add_rect(slide, col_lefts[2], table_top + row_h - 0.02,
             col_widths[2], row_h * 2 + 0.04,
             fill=TABLE_HILITE, rounded=False)

    add_textbox(slide, col_lefts[0], table_top + row_h, col_widths[0], row_h,
                "if right", size=15, bold=True, color=RIGHT_GREEN,
                align=PP_ALIGN.RIGHT)
    for i, val in enumerate(["$0", "$1.50", "$3.00"], start=1):
        add_textbox(slide, col_lefts[i], table_top + row_h, col_widths[i], row_h,
                    val, size=15, color=DARK_GRAY,
                    align=PP_ALIGN.CENTER)

    add_textbox(slide, col_lefts[0], table_top + 2 * row_h, col_widths[0], row_h,
                "if wrong", size=15, bold=True, color=WRONG_RED,
                align=PP_ALIGN.RIGHT)
    for i, val in enumerate(["$0", "$0", "$0"], start=1):
        add_textbox(slide, col_lefts[i], table_top + 2 * row_h, col_widths[i], row_h,
                    val, size=15, color=DARK_GRAY,
                    align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, 0.55, 6.67, 12.20, 0.30,
                "Paid on one randomly selected endorsement.",
                size=12, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    print("  slide 12: rebuilt parallel to Stage 3")


# ---------------------------------------------------------------------------
# Phase 3: Slides 19, 20 — descriptive RQ titles
# ---------------------------------------------------------------------------

RQ_TITLE_REWRITES = {
    "RQ1 — Pre-outcome trust by gender":
        "RQ1  —  How much did evaluators wager on the first decision?",
    "RQ2 — Gender × outcome on trust update":
        "RQ2  —  How much did one outcome change evaluator wagers?",
}


def rewrite_rq_title(slide, slide_label):
    """Rewrite the slide's RQ title shape, preserving styling. Also expands
    the title bbox so the longer title can wrap to 2 lines without colliding
    with the stats corner."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        current = shape.text_frame.text or ""
        for old, new in RQ_TITLE_REWRITES.items():
            if current.strip() == old:
                para = shape.text_frame.paragraphs[0]
                if not para.runs:
                    continue
                para.runs[0].text = new
                # Allow wrap
                shape.height = Inches(1.40)
                # Stats corner is at x=9.38; keep title to x<8.60
                if shape.width > Inches(8.60):
                    shape.width = Inches(8.60)
                else:
                    shape.width = max(shape.width, Inches(7.60))
                print(f"  {slide_label}: {old!r}\n    -> {new!r}")
                return
    print(f"  {slide_label}: WARN no matching RQ title")


# ---------------------------------------------------------------------------
# Phase 4: Slide 15 — variant cards with mini-stimuli
# ---------------------------------------------------------------------------

def build_mini_strength_bar(slide, x, y, w, h, fill_pct, fill_color):
    """Track + filled bar for the strength stim. fill_pct in 0-100."""
    track = add_rect(slide, x, y, w, h, fill=BAR_TRACK, rounded=True)
    track.line.color.rgb = LIGHT_RULE
    track.line.width = Pt(0.75)
    if fill_pct > 0:
        fill = add_rect(slide, x, y, w * (fill_pct / 100.0), h,
                        fill=fill_color, rounded=True)
        fill.line.fill.background()


def build_pill(slide, x, y, w, h, text, fill, text_color):
    pill = add_rect(slide, x, y, w, h, fill=fill, rounded=True)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 14, bold=True, color=text_color)
    return pill


def rebuild_slide15_randomization(slide):
    """Rebuild the 3 variant cards on slide 15, replacing the bullet list
    inside each with a mini-stimulus pair (live native shapes)."""
    # Capture & remove only the 3 variant cards (keep title + WE VARY header)
    cards_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt.startswith("1.  Sponsor gender") or \
               txt.startswith("2.  Endorsement strength") or \
               txt.startswith("3.  Outcome"):
                cards_to_remove.append(shape)
    spTree = slide.shapes._spTree
    for shape in cards_to_remove:
        spTree.remove(shape._element)
    print(f"  slide 15: removed {len(cards_to_remove)} old variant cards")

    # New card geometry — taller to fit the mini-stim
    card_y = 1.95
    card_w = 4.05
    card_h = 4.85
    card_xs = [0.43, 4.64, 8.85]

    # Background cards — softly rounded (corner=0.06) so they read as panels
    # rather than the lozenge shape that corner=0.5 produces.
    for x in card_xs:
        add_rect(slide, x, card_y, card_w, card_h,
                 fill=WHITE, outline=CARD_OUTLINE, outline_width=1.25,
                 rounded=True, corner=0.06)

    # ---- Card 1: Sponsor gender ------------------------------------------
    cx = card_xs[0]
    add_textbox(slide, cx + 0.20, card_y + 0.15, card_w - 0.40, 0.45,
                "1.  Sponsor gender",
                size=18, bold=True, color=NAVY_DEEP)
    add_textbox(slide, cx + 0.20, card_y + 0.65, card_w - 0.40, 0.40,
                "(name + colored avatar)",
                size=12, italic=True, color=MID_GRAY)

    # Stim 1 — Woman
    av_size = 0.70
    av_x = cx + 0.30
    stim1_y = card_y + 1.30
    add_oval(slide, av_x, stim1_y, av_size, av_size,
             fill=WOMAN_BG, outline=WOMAN_PINK, outline_width=4.0)
    add_textbox(slide, av_x + av_size + 0.18, stim1_y + 0.10, 1.80, 0.45,
                "Woman", size=18, bold=True, color=WOMAN_PINK)

    # Stim 2 — Man
    stim2_y = card_y + 2.55
    add_oval(slide, av_x, stim2_y, av_size, av_size,
             fill=MAN_BG, outline=MAN_BLUE, outline_width=4.0)
    add_textbox(slide, av_x + av_size + 0.18, stim2_y + 0.10, 1.80, 0.45,
                "Man", size=18, bold=True, color=MAN_BLUE)

    # ---- Card 2: Endorsement strength ------------------------------------
    cx = card_xs[1]
    add_textbox(slide, cx + 0.20, card_y + 0.15, card_w - 0.40, 0.45,
                "2.  Endorsement strength",
                size=18, bold=True, color=NAVY_DEEP)
    add_textbox(slide, cx + 0.20, card_y + 0.65, card_w - 0.40, 0.40,
                "(how confident the sponsor was)",
                size=12, italic=True, color=MID_GRAY)

    # Stim 1 — Strong (90% navy bar)
    bar_w = 2.50
    bar_h = 0.20
    bar_x = cx + (card_w - bar_w) / 2
    bar1_y = card_y + 1.35
    build_mini_strength_bar(slide, bar_x, bar1_y, bar_w, bar_h, 90, NAVY)
    add_textbox(slide, cx + 0.20, bar1_y + 0.30, card_w - 0.40, 0.40,
                "Strong", size=18, bold=True, color=NAVY_DEEP,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, cx + 0.20, bar1_y + 0.65, card_w - 0.40, 0.30,
                "very confident", size=12, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # Stim 2 — Weak (22% red bar)
    bar2_y = card_y + 2.85
    build_mini_strength_bar(slide, bar_x, bar2_y, bar_w, bar_h, 22, WRONG_RED)
    add_textbox(slide, cx + 0.20, bar2_y + 0.30, card_w - 0.40, 0.40,
                "Weak", size=18, bold=True, color=WRONG_RED,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, cx + 0.20, bar2_y + 0.65, card_w - 0.40, 0.30,
                "unsure", size=12, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # ---- Card 3: Outcome --------------------------------------------------
    cx = card_xs[2]
    add_textbox(slide, cx + 0.20, card_y + 0.15, card_w - 0.40, 0.45,
                "3.  Outcome",
                size=18, bold=True, color=NAVY_DEEP)
    add_textbox(slide, cx + 0.20, card_y + 0.65, card_w - 0.40, 0.40,
                "(was the endorsement right?)",
                size=12, italic=True, color=MID_GRAY)

    pill_w = 2.00
    pill_h = 0.50
    pill_x = cx + (card_w - pill_w) / 2

    # Stim 1 — CORRECT
    pill1_y = card_y + 1.40
    build_pill(slide, pill_x, pill1_y, pill_w, pill_h,
               "CORRECT", PILL_GREEN, PILL_TEXT_G)
    add_textbox(slide, cx + 0.20, pill1_y + 0.60, card_w - 0.40, 0.32,
                "endorsement was right", size=12, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # Stim 2 — WRONG
    pill2_y = card_y + 2.65
    build_pill(slide, pill_x, pill2_y, pill_w, pill_h,
               "WRONG", PILL_RED, PILL_TEXT_R)
    add_textbox(slide, cx + 0.20, pill2_y + 0.60, card_w - 0.40, 0.32,
                "endorsement was wrong", size=12, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    print("  slide 15: rebuilt with mini-stimuli for all 3 variants")


# ---------------------------------------------------------------------------
# Phase 5: Slide 14 — Stage 3 setup full rebuild with more info
# ---------------------------------------------------------------------------

def rebuild_slide14_setup(slide):
    """Rebuild the Stage 3 setup slide with sample size, what evaluators
    actually see, and the two-decision flow with the same sponsor."""
    n_shapes = clear_non_placeholder_shapes(slide)
    n_phs = remove_name_of_initiative_placeholder(slide)
    print(f"  slide 14: cleared {n_shapes} shapes + {n_phs} stale placeholder(s)")

    # Title
    add_textbox(slide, 0.55, 0.30, 12.20, 0.70,
                "Stage 3 setup",
                size=32, bold=True, color=NAVY_DEEP)

    # Sample size sub-line
    add_textbox(slide, 0.55, 1.05, 12.20, 0.40,
                "n = 403 Prolific evaluators (passed 5/5 comprehension checks)",
                size=16, italic=True, color=MID_GRAY)

    # "What each evaluator sees" header
    add_textbox(slide, 0.55, 1.55, 12.20, 0.40,
                "Each evaluator is randomly assigned to ONE sponsor and sees:",
                size=18, bold=True, color=NAVY_DEEP)

    # Bullets — what they see
    bullets = [
        "the sponsor's name + colored avatar (signals gender)",
        "that sponsor's endorsement of one candidate from a pair",
        "how confident the sponsor was on that endorsement (Strong / Weak)",
        "a slider to wager part of a $0.50 bank on the endorser being right",
    ]
    bullet_x = 0.95
    bullet_y = 2.05
    bullet_h = 0.42
    for i, bullet in enumerate(bullets):
        box = slide.shapes.add_textbox(in_(bullet_x), in_(bullet_y + i * bullet_h),
                                       in_(11.80), in_(bullet_h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.03)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = "•  "
        style_run(r1, 18, bold=True, color=NAVY)
        r2 = p.add_run()
        r2.text = bullet
        style_run(r2, 18, color=DARK_GRAY)

    # Italic line — back-to-back two decisions
    after_bullets_y = bullet_y + len(bullets) * bullet_h + 0.10
    add_textbox(slide, 0.55, after_bullets_y, 12.20, 0.42,
                "Each evaluator sees TWO decisions back-to-back, with the SAME sponsor on different candidate pairs:",
                size=15, italic=True, color=MID_GRAY)

    # Flow boxes: First wager → outcome revealed → Second wager
    flow_y = after_bullets_y + 0.50
    flow_h = 0.78
    box_w = 2.10
    arrow_w = 0.65
    total_w = 3 * box_w + 2 * arrow_w
    flow_left = (13.33 - total_w) / 2

    box_x = [
        flow_left,
        flow_left + box_w + arrow_w,
        flow_left + 2 * (box_w + arrow_w),
    ]
    labels = ["First wager", "outcome revealed", "Second wager"]
    fills = [BAR_FILL, PILL_GREEN, BAR_FILL]
    text_colors = [NAVY_DEEP, PILL_TEXT_G, NAVY_DEEP]
    outlines = [NAVY, RIGHT_GREEN, NAVY]

    for i in range(3):
        rect = add_rect(slide, box_x[i], flow_y, box_w, flow_h,
                        fill=fills[i], outline=outlines[i], outline_width=2.0,
                        rounded=True)
        tf = rect.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.text = labels[i]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style_run(p.runs[0], 17, bold=True, color=text_colors[i])

    # Arrow connectors
    for i in range(2):
        ax = box_x[i] + box_w + 0.05
        ay = flow_y + flow_h / 2 - 0.10
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            in_(ax), in_(ay), in_(arrow_w - 0.10), in_(0.20)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NAVY
        arrow.line.fill.background()
        arrow.text_frame.text = ""

    # DV formula
    formula_y = flow_y + flow_h + 0.30
    formula_box = slide.shapes.add_textbox(in_(0.55), in_(formula_y), in_(12.20), in_(0.55))
    tf = formula_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "trust change  =  (Second wager)  −  (First wager)"
    style_run(r1, 24, bold=True, color=NAVY_DEEP)

    print("  slide 14: rebuilt with sample size + bullet list + flow + DV")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides")

    s_12 = slides[11]   # Stage 2 incentive
    s_14 = slides[13]   # Stage 3 setup
    s_15 = slides[14]   # Stage 3 randomization
    s_19 = slides[18]   # RQ1
    s_20 = slides[19]   # RQ2
    s_22 = slides[21]   # RQ3 weak (badge text fix)

    print("\n[Phase 1] Slide 22: 'hedged / unsure' -> 'unsure'")
    patch_slide22_badge_caption(s_22)

    print("\n[Phase 2] Slide 12: rebuild Stage 2 incentive parallel to Stage 3")
    rebuild_slide12_stage2_incentive(s_12)

    print("\n[Phase 3] Slides 19 & 20: descriptive RQ titles")
    rewrite_rq_title(s_19, "slide 19")
    rewrite_rq_title(s_20, "slide 20")

    print("\n[Phase 4] Slide 15: variant cards with mini-stimuli")
    rebuild_slide15_randomization(s_15)

    print("\n[Phase 5] Slide 14: Stage 3 setup full rebuild")
    rebuild_slide14_setup(s_14)

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
