"""
patch_2026_04_07_content.py — content patch for slides 2, 9, 10, 11, 12.

This runs AFTER patch_2026_04_07.py (which swapped stimuli + figures).
It touches a different set of slides and different kinds of content:

  SLIDE  2 — Hook (letter-of-recommendation framing)
             REMOVE the 19 Mark Davies / Sarah Lin manager-card shapes
             BUILD a new hook: two parallel letter cards (Nick Calder,
             Jose Cervantez) both signed by "Susan Keller" (or "Bob Keller"
             if Jose swaps it manually), a right-pointing chevron in the
             middle with a SUCCEEDED pill, a bottom question bar with
             "How much does one outcome change your trust in the next letter?"
             Preserves the slide-number placeholder at (12.62, 7.09).

  SLIDE  9 — Stage 2 incentive (sponsors)
             REMOVE the blank "Incentive slides (Stage 2)" content placeholder.
             BUILD title + big equation + condition subcue + slider widget
             with handle at 50% + big $1.50 reveal + 5-column payout table
             with the 50% column highlighted + footnote.
             Keeps the footer and slide-number placeholders.

  SLIDE 10 — Stage 3 incentive (evaluators)
             Same treatment as slide 9, but with two stacked formulas
             (correct/wrong), twin reveal, and a 3-row payout table showing
             wager 0% / 50% / 100% × {correct, wrong}.

  SLIDE 11 — Stage 3 setup (what evaluators actually do)
             REMOVE the empty content placeholder.
             BUILD title "Stage 3" + three facts summarizing the instructions
             the comprehension checks tested (tasks, endorser role,
             definition of "correct") + lead-in + D1 → outcome → D2 flow
             boxes with arrows + DV equation.

  SLIDE 12 — Stage 3 randomization (trim)
             REMOVE the "WHAT EACH EVALUATOR DOES" flow section (header,
             three flow cards, two arrow shapes, and the "Outcome variable"
             caption) — everything with top >= 4.3 inches.
             KEEP the title, "WE VARY THREE THINGS" header, the three
             randomization cards, and the slide-number placeholder.

Never touches any other slide. Back up the deck first.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PROJECT = Path(__file__).resolve().parent.parent.parent
DECK    = PROJECT / "docs" / "UChicago-0410.pptx"

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------

NAVY         = RGBColor(0x01, 0x1F, 0x5B)   # Wharton navy
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT   = RGBColor(0x1F, 0x1F, 0x1F)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY     = RGBColor(0x5B, 0x65, 0x7A)
TITLE_GRAY   = RGBColor(0x4B, 0x55, 0x63)
LIGHT_RULE   = RGBColor(0x9A, 0xA7, 0xBD)
CARD_OUTLINE = RGBColor(0xC7, 0xCF, 0xDD)
LIGHT_FILL   = RGBColor(0xE5, 0xE7, 0xEB)
BAR_FILL     = RGBColor(0xEA, 0xF0, 0xFB)
PILL_BG      = RGBColor(0xDF, 0xF4, 0xE8)
PILL_TEXT    = RGBColor(0x11, 0x6B, 0x45)
RED_WRONG    = RGBColor(0x99, 0x00, 0x00)

FONT = "Calibri"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def in_(v):
    return Inches(v)

def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def add_textbox(slide, x, y, w, h, text, size,
                bold=False, italic=False, color=BLACK_TEXT,
                align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color)
    return box

def add_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
             rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, in_(x), in_(y), in_(w), in_(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    if rounded:
        try:
            shape.adjustments[0] = 0.5
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape

def add_rule(slide, x, y, w, color=LIGHT_RULE):
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, in_(x), in_(y), in_(w), Inches(0.018)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()
    rule.text_frame.text = ""
    return rule

def add_right_arrow(slide, x, y, w, h, fill=BAR_FILL, outline=NAVY,
                    outline_width=1.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, in_(x), in_(y), in_(w), in_(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = outline
    shape.line.width = Pt(outline_width)
    shape.text_frame.text = ""
    return shape


# ---------------------------------------------------------------------------
# Slide 2 — Hook (letter-of-rec)
# ---------------------------------------------------------------------------

def clear_slide2(slide):
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        spTree.remove(shape._element)
        removed += 1
    return removed


LETTER_BODY_TEMPLATE = (
    "I am pleased to recommend {name} for admission to your doctoral "
    "program in organizational behavior. In my seminar and as a research "
    "assistant, {first} distinguished himself through careful analytical "
    "reasoning, dependable follow-through, and a mature willingness to "
    "revise his views when the evidence changed. He writes clearly, "
    "collaborates well, and consistently raises the quality of the work "
    "around him. I recommend him with enthusiasm."
)


def build_letter_card(slide, x_in, student, signer,
                      signer_title="Professor, Department of Behavioral Science"):
    """Build one letter-of-rec card at the given x position."""
    CARD_TOP   = 0.55
    CARD_W     = 4.70
    CARD_H     = 5.55
    BAND_H     = 0.48

    # 1. Card background (white fill, gray outline)
    add_rect(slide, x_in, CARD_TOP, CARD_W, CARD_H,
             fill=WHITE, outline=CARD_OUTLINE, outline_width=1.25)

    # 2. Letterhead band (navy) — added AFTER card so it's on top
    band = add_rect(slide, x_in, CARD_TOP, CARD_W, BAND_H, fill=NAVY)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.text = "GRAYBRIDGE UNIVERSITY  \u00b7  Department of Behavioral Science"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 12, bold=True, color=WHITE)

    # 3. Body text area (salutation + body paragraph)
    body_left = x_in + 0.32
    body_w    = CARD_W - 0.64  # 4.06
    body_top  = CARD_TOP + BAND_H + 0.18  # 1.21
    body_h    = 3.30

    first_name = student.split()[0]

    body = slide.shapes.add_textbox(
        in_(body_left), in_(body_top), in_(body_w), in_(body_h)
    )
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    # Salutation
    tf.text = "To the Admissions Committee,"
    p1 = tf.paragraphs[0]
    style_run(p1.runs[0], 13, color=BLACK_TEXT)
    p1.space_after = Pt(8)

    # Body paragraph
    p2 = tf.add_paragraph()
    p2.text = LETTER_BODY_TEMPLATE.format(name=student, first=first_name)
    style_run(p2.runs[0], 13, color=BLACK_TEXT)
    p2.line_spacing = 1.10

    # 4. Horizontal rule before signature
    rule_y = 4.68
    add_rule(slide, body_left, rule_y, body_w - 0.30)

    # 5. Signature block
    sig = slide.shapes.add_textbox(
        in_(body_left), in_(4.80), in_(body_w), in_(1.30)
    )
    tf = sig.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)

    tf.text = "Sincerely,"
    p_sinc = tf.paragraphs[0]
    style_run(p_sinc.runs[0], 12, color=BLACK_TEXT)
    p_sinc.space_after = Pt(6)

    p_name = tf.add_paragraph()
    p_name.text = signer
    style_run(p_name.runs[0], 13, bold=True, italic=True, color=NAVY)

    p_title = tf.add_paragraph()
    p_title.text = signer_title
    style_run(p_title.runs[0], 10.5, color=TITLE_GRAY)


def build_slide2_hook(slide):
    # === Letter 1 — Nick Calder ===
    build_letter_card(slide, x_in=0.55,
                      student="Nick Calder",
                      signer="Susan Keller")

    # === Letter 2 — Jose Cervantez ===
    build_letter_card(slide, x_in=8.08,
                      student="Jose Cervantez",
                      signer="Susan Keller")

    # === Transition chevron ===
    add_right_arrow(slide, 5.47, 2.15, 2.36, 1.70,
                    fill=BAR_FILL, outline=NAVY, outline_width=1.75)

    # Top label "ONE OUTCOME" inside the chevron body portion
    add_textbox(slide, 5.47, 2.30, 1.80, 0.25,
                "ONE OUTCOME",
                size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Outcome pill "SUCCEEDED"
    pill = add_rect(slide, 5.62, 2.65, 1.50, 0.48,
                    fill=PILL_BG, rounded=True)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = "SUCCEEDED"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 14, bold=True, color=PILL_TEXT)

    # Bottom label below the pill (still inside chevron)
    add_textbox(slide, 5.47, 3.22, 1.80, 0.26,
                "new letter about Jose",
                size=9, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # === Bottom question bar ===
    bar = add_rect(slide, 0.55, 6.30, 11.70, 0.60,
                   fill=BAR_FILL, rounded=True)
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.text = "How much does one outcome change your trust in the next letter?"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 22, bold=True, color=NAVY)


# ---------------------------------------------------------------------------
# Slides 9/10/11 — clear content placeholder
# ---------------------------------------------------------------------------

def clear_content_placeholder(slide):
    """Remove the main content placeholder (idx=1), keep footer (idx=10)
    and slide-number (idx=11)."""
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        try:
            idx = shape.placeholder_format.idx
        except Exception:
            continue
        if idx in (10, 11):
            continue
        spTree.remove(shape._element)
        removed += 1
    return removed


# ---------------------------------------------------------------------------
# Slider widget used on slides 9 and 10
# ---------------------------------------------------------------------------

def build_slider_widget(slide, track_y, handle_pct=0.5, above_label="evaluator wagers 50%"):
    """Build a horizontal slider widget centered on the slide canvas.

    track at y=track_y, height 0.22", width 6.33" centered horizontally.
    A filled navy portion covers the left portion up to the handle.
    A circular handle sits at the boundary.
    Labels '0%' and '100%' sit below the track endpoints.
    A label above the track names the current value.
    """
    SLIDER_LEFT  = 3.50
    SLIDER_WIDTH = 6.33
    TRACK_H      = 0.22
    HANDLE_SIZE  = 0.42

    # Above-track label
    add_textbox(slide, 0.55, track_y - 0.50, 12.20, 0.32,
                above_label, size=13, color=DARK_GRAY,
                align=PP_ALIGN.CENTER)

    # Track (light gray pill)
    track = add_rect(slide, SLIDER_LEFT, track_y, SLIDER_WIDTH, TRACK_H,
                     fill=LIGHT_FILL, outline=CARD_OUTLINE, outline_width=0.75,
                     rounded=True)

    # Filled portion (navy, up to the handle)
    filled_w = SLIDER_WIDTH * handle_pct
    if filled_w > 0:
        add_rect(slide, SLIDER_LEFT, track_y, filled_w, TRACK_H,
                 fill=NAVY, rounded=True)

    # Handle (white circle with navy outline)
    handle_cx = SLIDER_LEFT + filled_w
    handle_x  = handle_cx - HANDLE_SIZE / 2
    handle_y  = track_y + TRACK_H / 2 - HANDLE_SIZE / 2
    handle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, in_(handle_x), in_(handle_y),
        in_(HANDLE_SIZE), in_(HANDLE_SIZE)
    )
    handle.fill.solid()
    handle.fill.fore_color.rgb = WHITE
    handle.line.color.rgb = NAVY
    handle.line.width = Pt(1.75)
    handle.text_frame.text = ""

    # 0% label at left
    add_textbox(slide, SLIDER_LEFT - 0.75, track_y + 0.32, 0.70, 0.30,
                "0%", size=12, color=MID_GRAY, align=PP_ALIGN.RIGHT)
    # 100% label at right
    add_textbox(slide, SLIDER_LEFT + SLIDER_WIDTH + 0.05, track_y + 0.32, 0.80, 0.30,
                "100%", size=12, color=MID_GRAY, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 9 — Stage 2 incentive
# ---------------------------------------------------------------------------

def build_slide9_stage2_incentive(slide):
    # Title
    add_textbox(slide, 0.55, 0.40, 12.20, 0.60,
                "Stage 2 incentive",
                size=28, bold=True, color=NAVY)

    # Equation
    add_textbox(slide, 0.55, 1.25, 12.20, 0.80,
                "sponsor earns  =  $3  \u00d7  evaluator wager",
                size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Condition subcue
    add_textbox(slide, 0.55, 2.10, 12.20, 0.35,
                "only if the endorsed candidate scored higher on logical reasoning",
                size=14, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # --- Highlight column (add BEFORE table so it's behind) ---
    # Table geometry: 6 columns, first is label, 5 data columns
    TABLE_LEFT  = 2.15
    LABEL_W     = 1.75
    CELL_W      = 1.70
    HEADER_Y    = 5.30
    VALUE_Y     = 5.75
    ROW_H       = 0.40

    # 50% column is index 2 in data (0%, 25%, 50%, 75%, 100%)
    hl_col_x = TABLE_LEFT + LABEL_W + 2 * CELL_W
    add_rect(slide, hl_col_x, HEADER_Y - 0.06, CELL_W,
             2 * ROW_H + 0.12,
             fill=BAR_FILL)

    # --- Slider widget ---
    build_slider_widget(slide, track_y=3.10, handle_pct=0.5,
                        above_label="evaluator wagers 50%")

    # --- Big $1.50 reveal ---
    add_textbox(slide, 4.67, 3.85, 4.00, 1.00,
                "$1.50",
                size=56, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # --- Payout table ---
    headers = ["wager", "0%", "25%", "50%", "75%", "100%"]
    values  = ["bonus", "$0", "$0.75", "$1.50", "$2.25", "$3.00"]

    def _cell(x, y, w, h, text, bold=False, color=DARK_GRAY):
        add_textbox(slide, x, y, w, h, text,
                    size=15, bold=bold, color=color,
                    align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

    # Header row
    _cell(TABLE_LEFT, HEADER_Y, LABEL_W, ROW_H, headers[0],
          bold=True, color=MID_GRAY)
    for i, txt in enumerate(headers[1:]):
        x = TABLE_LEFT + LABEL_W + i * CELL_W
        _cell(x, HEADER_Y, CELL_W, ROW_H, txt, bold=True, color=NAVY)

    # Value row
    _cell(TABLE_LEFT, VALUE_Y, LABEL_W, ROW_H, values[0],
          bold=True, color=MID_GRAY)
    for i, txt in enumerate(values[1:]):
        x = TABLE_LEFT + LABEL_W + i * CELL_W
        _cell(x, VALUE_Y, CELL_W, ROW_H, txt, bold=False, color=DARK_GRAY)

    # Footnote
    add_textbox(slide, 0.55, 6.55, 12.20, 0.35,
                "Paid on one randomly selected endorsement.",
                size=12, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 10 — Stage 3 incentive
# ---------------------------------------------------------------------------

def build_slide10_stage3_incentive(slide):
    # Title
    add_textbox(slide, 0.55, 0.40, 12.20, 0.60,
                "Stage 3 incentive",
                size=28, bold=True, color=NAVY)

    # Equation 1: correct
    add_textbox(slide, 0.55, 1.15, 12.20, 0.55,
                "correct:     $0.50  +  $0.50  \u00d7  wager",
                size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Equation 2: wrong
    add_textbox(slide, 0.55, 1.75, 12.20, 0.55,
                "wrong:     $0.50  \u2212  $0.50  \u00d7  wager",
                size=26, bold=True, color=RED_WRONG, align=PP_ALIGN.CENTER)

    # --- Highlight column behind 50% in table ---
    TABLE_LEFT = 3.00
    LABEL_W    = 1.50
    CELL_W     = 2.20
    HEADER_Y   = 4.95
    CORRECT_Y  = 5.40
    WRONG_Y    = 5.85
    ROW_H      = 0.40

    # 50% column is index 1 in data (0%, 50%, 100%)
    hl_col_x = TABLE_LEFT + LABEL_W + 1 * CELL_W
    add_rect(slide, hl_col_x, HEADER_Y - 0.06, CELL_W,
             3 * ROW_H + 0.12,
             fill=BAR_FILL)

    # --- Slider widget ---
    build_slider_widget(slide, track_y=2.90, handle_pct=0.5,
                        above_label="evaluator wagers 50%")

    # --- Twin reveal ---
    add_textbox(slide, 1.50, 3.75, 5.00, 0.80,
                "correct:  $0.75",
                size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 6.85, 3.75, 5.00, 0.80,
                "wrong:  $0.25",
                size=32, bold=True, color=RED_WRONG, align=PP_ALIGN.CENTER)

    # --- Payout table (3 rows, 4 columns: label + 0%/50%/100%) ---
    headers      = ["wager",   "0%",    "50%",   "100%"]
    correct_row  = ["correct", "$0.50", "$0.75", "$1.00"]
    wrong_row    = ["wrong",   "$0.50", "$0.25", "$0.00"]

    def _cell(x, y, w, h, text, bold=False, color=DARK_GRAY, size=15):
        add_textbox(slide, x, y, w, h, text,
                    size=size, bold=bold, color=color,
                    align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

    def _row(y, cells, label_color, value_color, value_bold=False):
        _cell(TABLE_LEFT, y, LABEL_W, ROW_H, cells[0],
              bold=True, color=label_color)
        for i, txt in enumerate(cells[1:]):
            x = TABLE_LEFT + LABEL_W + i * CELL_W
            _cell(x, y, CELL_W, ROW_H, txt,
                  bold=value_bold, color=value_color)

    _row(HEADER_Y,  headers,     MID_GRAY, NAVY,      value_bold=True)
    _row(CORRECT_Y, correct_row, NAVY,     DARK_GRAY)
    _row(WRONG_Y,   wrong_row,   RED_WRONG, DARK_GRAY)

    # Footnote
    add_textbox(slide, 0.55, 6.50, 12.20, 0.35,
                "$0.50 bank per decision  \u00b7  one randomly selected decision determines the bonus",
                size=12, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 11 — Stage 3 setup
# ---------------------------------------------------------------------------

def build_slide11_stage3_setup(slide):
    # Title
    add_textbox(slide, 0.55, 0.40, 12.20, 0.60,
                "Stage 3",
                size=28, bold=True, color=NAVY)

    # Fact 1
    add_textbox(slide, 0.80, 1.20, 11.70, 0.55,
                "Candidates completed three tasks: word-search, general-knowledge, and logical-reasoning.",
                size=18, color=DARK_GRAY)

    # Fact 2
    add_textbox(slide, 0.80, 1.95, 11.70, 0.80,
                "Endorsers viewed two candidates and used a slider to predict which one would score higher on logical-reasoning.",
                size=18, color=DARK_GRAY)

    # Fact 3
    add_textbox(slide, 0.80, 2.85, 11.70, 0.90,
                "An endorsement is \u201ccorrect\u201d when the endorsed candidate outscores a randomly selected other candidate on logical-reasoning.",
                size=18, color=DARK_GRAY)

    # Lead-in
    add_textbox(slide, 0.80, 3.90, 11.70, 0.45,
                "Each evaluator sees two endorsements from the SAME sponsor:",
                size=18, italic=True, color=MID_GRAY)

    # Flow boxes
    box_y = 4.55
    box_w = 2.00
    box_h = 0.95
    box1_x = 2.82
    box2_x = 5.67
    box3_x = 8.52

    def _flow_box(x, text):
        box = add_rect(slide, x, box_y, box_w, box_h,
                       fill=BAR_FILL, outline=NAVY, outline_width=1.25,
                       rounded=True)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.text = text
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style_run(p.runs[0], 16, bold=True, color=NAVY)

    _flow_box(box1_x, "D1 wager")
    _flow_box(box2_x, "outcome reveal")
    _flow_box(box3_x, "D2 wager")

    # Arrows between boxes
    arrow1 = add_right_arrow(slide, 4.87, 4.92, 0.75, 0.20,
                             fill=NAVY, outline=NAVY, outline_width=0.5)
    arrow2 = add_right_arrow(slide, 7.72, 4.92, 0.75, 0.20,
                             fill=NAVY, outline=NAVY, outline_width=0.5)

    # DV equation
    add_textbox(slide, 0.55, 6.05, 12.20, 0.60,
                "trust change  =  (D2 wager)  \u2212  (D1 wager)",
                size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 12 — Trim flow section
# ---------------------------------------------------------------------------

def trim_slide12(slide):
    """Remove anything on slide 12 with top >= 4.3 inches, except placeholders."""
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        top_emu = shape.top or 0
        top_in = top_emu / 914400.0
        if top_in >= 4.3:
            spTree.remove(shape._element)
            removed += 1
    return removed


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("=" * 72)
    print(f"  Content patch: {DECK}")
    print("=" * 72)

    prs = Presentation(str(DECK))
    print(f"  Loaded ({len(prs.slides)} slides)")

    # Slide 2 — Hook
    print("\n[1/5] Slide 2 — rebuild letter-of-rec hook...")
    s2 = prs.slides[1]
    removed = clear_slide2(s2)
    print(f"  cleared {removed} non-placeholder shape(s) from old hook")
    build_slide2_hook(s2)
    print(f"  built new hook (shape count now {len(s2.shapes)})")

    # Slide 9 — Stage 2 incentive
    print("\n[2/5] Slide 9 — Stage 2 incentive...")
    s9 = prs.slides[8]
    removed = clear_content_placeholder(s9)
    print(f"  cleared {removed} placeholder(s)")
    build_slide9_stage2_incentive(s9)
    print(f"  built Stage 2 incentive (shape count now {len(s9.shapes)})")

    # Slide 10 — Stage 3 incentive
    print("\n[3/5] Slide 10 — Stage 3 incentive...")
    s10 = prs.slides[9]
    removed = clear_content_placeholder(s10)
    print(f"  cleared {removed} placeholder(s)")
    build_slide10_stage3_incentive(s10)
    print(f"  built Stage 3 incentive (shape count now {len(s10.shapes)})")

    # Slide 11 — Stage 3 setup
    print("\n[4/5] Slide 11 — Stage 3 setup...")
    s11 = prs.slides[10]
    removed = clear_content_placeholder(s11)
    print(f"  cleared {removed} placeholder(s)")
    build_slide11_stage3_setup(s11)
    print(f"  built Stage 3 setup (shape count now {len(s11.shapes)})")

    # Slide 12 — trim flow
    print("\n[5/5] Slide 12 — trim 'WHAT EACH EVALUATOR DOES' flow section...")
    s12 = prs.slides[11]
    removed = trim_slide12(s12)
    print(f"  removed {removed} flow-section shape(s) (shape count now {len(s12.shapes)})")

    prs.save(str(DECK))
    print(f"\n  Saved {DECK}")
    print("=" * 72)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
