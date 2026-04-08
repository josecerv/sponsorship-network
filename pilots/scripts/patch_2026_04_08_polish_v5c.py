"""
patch_2026_04_08_polish_v5c.py — follow-up fix for v5.

The v5 patch deleted the RQ3 titles on slides 24 and 25 along with the
old badge because the delete rect x_max (3.00) accidentally covered the
right-pushed title's x=2.95 corner. This patch re-adds the titles as
left-aligned textboxes at the new position.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"

NAVY = RGBColor(0x01, 0x1F, 0x5B)


TITLES = {
    23: "RQ3  \u2014  Strong endorsements:  how did one outcome shift wagers?",
    24: "RQ3  \u2014  Weak endorsements:  how did one outcome shift wagers?",
}


def has_title(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip().startswith("RQ3"):
            return True
    return False


def add_left_title(slide, text):
    box = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.30), Inches(6.70), Inches(1.40),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.name = "Arial"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = NAVY


def main():
    prs = Presentation(str(DECK))
    slides = list(prs.slides)

    for idx, title_text in TITLES.items():
        slide = slides[idx]
        if has_title(slide):
            print(f"  slide {idx+1}: title already present, skipping")
            continue
        add_left_title(slide, title_text)
        print(f"  slide {idx+1}: added left-aligned title")

    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
