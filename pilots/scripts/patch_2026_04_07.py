"""
patch_2026_04_07.py — focused, one-shot patcher for the UChicago deck.

What this does (and ONLY this):
  1. Swaps walkthrough PNGs on slides 13, 14, 15 (Step 1 / Step 2 / Step 3)
     to pick up the new Stage 3 renders with real Qualtrics avatars + the
     full GENDER_STYLE CSS (pink/blue borders, tinted backgrounds, glow
     rings, badge pills, card accents).
  2. Swaps result-figure PNGs on slides 16, 17, 18 (RQ1 / RQ2 strong /
     RQ2 weak) to pick up the transposed plots (gender on x-axis, outcome
     as the fill legend: Success=NAVY, Failure=RED).

What this does NOT do:
  - Touch slides 1, 2 (cover / hook), 3-12 (intro + incentive slides),
    19 (RQ3), 20 (Summary), 21-23 (Running story / Thanks / new appendix).
  - Re-run the full patch_deck.py (whose SLIDE_* indices are stale and
    would clobber the wrong slides if run blind — see memory).
  - Add or remove any slide-number placeholders (those are already in
    place on every slide except the cover).

Preconditions:
  - Deck is at docs/UChicago-0410.pptx. Back it up first.
  - Fresh PNGs exist:
      pilots/output/stage3_rendered_screens/walkthrough/W_correct_strong_d{1,2_correct,3}.png
      pilots/output/talk_figures/rq{1_gender_x_outcome, 2_strong_only, 2_weak_only}.png
  - Current deck has 23 slides; the target slides still hold a single
    large picture that is the one we want to replace.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

PROJECT = Path(__file__).resolve().parent.parent.parent
DECK    = PROJECT / "docs" / "UChicago-0410.pptx"

WALKTHROUGH_DIR = PROJECT / "pilots" / "output" / "stage3_rendered_screens" / "walkthrough"
FIG_DIR         = PROJECT / "pilots" / "output" / "talk_figures"

# (slide_index_1based, png_path, label)
SWAPS = [
    # Walkthrough stimuli — preserve the original picture geometry on each slide
    (13, WALKTHROUGH_DIR / "W_correct_strong_d1.png",         "Step 1 walkthrough"),
    (14, WALKTHROUGH_DIR / "W_correct_strong_d2_correct.png", "Step 2 walkthrough"),
    (15, WALKTHROUGH_DIR / "W_correct_strong_d3.png",         "Step 3 walkthrough"),
    # Result figures — uniform 13.13" x 5.25" geometry at (0.10", 1.80")
    (16, FIG_DIR / "rq1_gender_x_outcome.png", "RQ1 transposed"),
    (17, FIG_DIR / "rq2_strong_only.png",      "RQ2 strong transposed"),
    (18, FIG_DIR / "rq2_weak_only.png",        "RQ2 weak transposed"),
]

# Slides where we force a uniform slide-fill geometry (matches patch_deck.py's
# original RQ1-4 layout choice). For walkthrough slides we preserve whatever
# geometry Jose last set.
FORCE_GEOMETRY = {
    16: (Inches(0.10), Inches(1.80), Inches(13.13), Inches(5.25)),
    17: (Inches(0.10), Inches(1.80), Inches(13.13), Inches(5.25)),
    18: (Inches(0.10), Inches(1.80), Inches(13.13), Inches(5.25)),
}


def find_largest_picture(slide):
    """Return the picture shape with the largest area on the slide, or None."""
    best, best_area = None, 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            area = (shape.width or 0) * (shape.height or 0)
            if area > best_area:
                best = shape
                best_area = area
    return best


def replace_picture(slide, old_shape, new_png_path, forced_geom=None):
    """Remove old_shape and add a new picture at the same (or forced) geometry."""
    if forced_geom is not None:
        left, top, width, height = forced_geom
    else:
        left   = old_shape.left
        top    = old_shape.top
        width  = old_shape.width
        height = old_shape.height

    spTree = slide.shapes._spTree
    spTree.remove(old_shape._element)
    return slide.shapes.add_picture(
        str(new_png_path), left, top, width=width, height=height
    )


def emu_in(x):
    return round(x / 914400.0, 2)


def main():
    print("=" * 72)
    print(f"  Patching {DECK}")
    print("=" * 72)

    # Preflight — every PNG must exist
    missing = [str(p) for _, p, _ in SWAPS if not p.exists()]
    if missing:
        print("ERROR: missing PNG inputs:")
        for m in missing:
            print(f"  - {m}")
        return 1

    prs = Presentation(str(DECK))
    n_slides = len(prs.slides)
    print(f"  Loaded deck with {n_slides} slides\n")

    if n_slides != 23:
        print(f"  WARN: deck has {n_slides} slides, expected 23. "
              f"Double-check slide indices before/after this patch.")

    for slide_idx, png_path, label in SWAPS:
        slide = prs.slides[slide_idx - 1]
        pic = find_largest_picture(slide)
        if pic is None:
            print(f"  slide {slide_idx:2d} ({label:28s}): NO picture found — SKIPPED")
            continue

        old_geom = (emu_in(pic.left), emu_in(pic.top),
                    emu_in(pic.width), emu_in(pic.height))

        forced = FORCE_GEOMETRY.get(slide_idx)
        new_pic = replace_picture(slide, pic, png_path, forced_geom=forced)
        new_geom = (emu_in(new_pic.left), emu_in(new_pic.top),
                    emu_in(new_pic.width), emu_in(new_pic.height))

        print(f"  slide {slide_idx:2d} ({label:28s}): "
              f"swapped picture (old={old_geom}, new={new_geom})")

    prs.save(str(DECK))
    print(f"\n  Saved {DECK}")
    print("=" * 72)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
