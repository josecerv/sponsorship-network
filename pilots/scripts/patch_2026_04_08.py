"""
patch_2026_04_08.py — multi-slide hook + RQ reorder + RQ3 transpose.

Run AFTER `make_talk_figures.R` so the transposed rq3_initial_trust.png is on disk.
Back up `docs/UChicago-0410.pptx` BEFORE running (the script also skips work on
shapes it doesn't recognize, but you still want a safety net).

Phases
======

1. RQ text rewrite on slide 5 — reorder the three RQ cards' run[1] text:
      RQ1 = pre-outcome (was RQ3)
      RQ2 = gender × outcome (was RQ1)
      RQ3 = strength (was RQ2)
   Formatting (colors, bold, size) is preserved because we only touch run.text.

2. Rebuild slide 2 as hook 2a — centered single letter:
      Robert Keller -> Nick Calder
   Clears every non-placeholder shape on slide 2 and builds one big centered
   letter card. No chevron, no second letter, no bottom bar. Just the letter.

3. Add NEW slides for hook 2b and 2c using the "No Title" layout (which has
   the native slide-number placeholder baked in — inherited correctly).
   2b = two letters (Robert -> Nick SUCCEEDED, Robert -> Jose new letter)
   2c = identical layout but signatures flipped to SUSAN Keller + question bar

4. Rewrite result-slide titles in place:
      old slide 16: "RQ1 — Gender × outcome..."  ->  "RQ2 — Gender × outcome..."
      old slide 17: "RQ2 — Strong endorsers"     ->  "RQ3 — Strong endorsers"
      old slide 18: "RQ2 — Weak endorsers"       ->  "RQ3 — Weak endorsers"
      old slide 19: "RQ3 — Pre-outcome..."       ->  "RQ1 — Pre-outcome..."

5. Swap the RQ3/new-RQ1 figure on old slide 19 with the freshly re-rendered
   (transposed) rq3_initial_trust.png.

6. Reorder slides via sldIdLst:
      - Move the two new hook slides to positions 3 and 4 (0-indexed 2, 3)
        so the deck goes:  cover  |  2a  |  2b  |  2c  |  what is sponsorship ...
      - Move old slide 19 (pre-outcome) so it sits right before the RQ2
        gender×outcome slide, making it the first result slide.

Uses python-pptx NATIVE APIs everywhere except one place: sldIdLst child
reordering (unavoidable for slide moves). Doesn't clone raw XML, doesn't
touch relationships, doesn't duplicate shape IDs — this is the safe path
for avoiding PowerPoint "repair" prompts.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
RQ3_FIG = PROJECT / "pilots" / "output" / "talk_figures" / "rq3_initial_trust.png"


# ---------------------------------------------------------------------------
# Palette (matches patch_2026_04_07_content.py so the hook slides look
# consistent with the rest of the deck)
# ---------------------------------------------------------------------------

NAVY         = RGBColor(0x01, 0x1F, 0x5B)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT   = RGBColor(0x1F, 0x1F, 0x1F)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY     = RGBColor(0x5B, 0x65, 0x7A)
TITLE_GRAY   = RGBColor(0x4B, 0x55, 0x63)
LIGHT_RULE   = RGBColor(0x9A, 0xA7, 0xBD)
CARD_OUTLINE = RGBColor(0xC7, 0xCF, 0xDD)
BAR_FILL     = RGBColor(0xEA, 0xF0, 0xFB)
PILL_BG      = RGBColor(0xDF, 0xF4, 0xE8)
PILL_TEXT    = RGBColor(0x11, 0x6B, 0x45)
HIGHLIGHT_PK = RGBColor(0xFC, 0xE7, 0xF3)   # subtle pink for Susan signature highlight

FONT = "Calibri"


# ---------------------------------------------------------------------------
# Small shape helpers
# ---------------------------------------------------------------------------

def in_(v): return Inches(v)

def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size,
                bold=False, italic=False, color=BLACK_TEXT,
                align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color)
    return box


def add_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
             rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, in_(x), in_(y), in_(w), in_(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    if rounded:
        try:
            shape.adjustments[0] = 0.5
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


def add_rule(slide, x, y, w, color=LIGHT_RULE):
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, in_(x), in_(y), in_(w), Inches(0.018)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()
    rule.text_frame.text = ""
    return rule


def add_right_arrow(slide, x, y, w, h, fill=BAR_FILL, outline=NAVY,
                    outline_width=1.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, in_(x), in_(y), in_(w), in_(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = outline
    shape.line.width = Pt(outline_width)
    shape.text_frame.text = ""
    return shape


# ---------------------------------------------------------------------------
# Letter-of-rec body template (same voice across all three hook slides)
# ---------------------------------------------------------------------------

LETTER_BODY_TEMPLATE = (
    "I am pleased to recommend {name} for admission to your doctoral "
    "program in organizational behavior. In my seminar and as a research "
    "assistant, {first} distinguished himself through careful analytical "
    "reasoning, dependable follow-through, and a mature willingness to "
    "revise his views when the evidence changed. He writes clearly, "
    "collaborates well, and consistently raises the quality of the work "
    "around him. I recommend him with enthusiasm."
)


def build_letter_card(slide, x_in, y_in, w_in, h_in,
                      student, signer,
                      signer_title="Professor, Department of Behavioral Science",
                      body_font_size=13,
                      signer_highlight=False):
    """Build one letter-of-rec card at the given position and size.

    Parameters
    ----------
    x_in, y_in, w_in, h_in : float
        Card bounding box in inches.
    student : str
        Name of the student the letter is about.
    signer : str
        Full signer name shown above the signer title.
    signer_highlight : bool
        If True, draw a subtle pink rectangle behind the signer name
        (used on the Susan-Keller flip to draw the eye to the change).
    """
    BAND_H = 0.48
    body_left = x_in + 0.32
    body_w = w_in - 0.64
    body_top = y_in + BAND_H + 0.18
    body_h = h_in - BAND_H - 0.18 - 2.0  # leave room for divider + signature

    # 1. Card background
    add_rect(slide, x_in, y_in, w_in, h_in,
             fill=WHITE, outline=CARD_OUTLINE, outline_width=1.25)

    # 2. Letterhead band (navy)
    band = add_rect(slide, x_in, y_in, w_in, BAND_H, fill=NAVY)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.text = "GRAYBRIDGE UNIVERSITY  \u00b7  Department of Behavioral Science"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 12, bold=True, color=WHITE)

    # 3. Body
    first_name = student.split()[0]
    body = slide.shapes.add_textbox(
        in_(body_left), in_(body_top), in_(body_w), in_(body_h)
    )
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    tf.text = "To the Admissions Committee,"
    p1 = tf.paragraphs[0]
    style_run(p1.runs[0], body_font_size, color=BLACK_TEXT)
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = LETTER_BODY_TEMPLATE.format(name=student, first=first_name)
    style_run(p2.runs[0], body_font_size, color=BLACK_TEXT)
    p2.line_spacing = 1.10

    # 4. Divider above signature (anchor to CARD bottom, not body top)
    rule_y = y_in + h_in - 1.42
    add_rule(slide, body_left, rule_y, body_w - 0.30)

    # 5. Optional pink highlight behind signer name
    sig_top = y_in + h_in - 1.30
    if signer_highlight:
        hl_w = min(1.60, body_w - 0.20)
        add_rect(slide, body_left + 0.02, sig_top + 0.24, hl_w, 0.30,
                 fill=HIGHLIGHT_PK)

    # 6. Signature block
    sig = slide.shapes.add_textbox(
        in_(body_left), in_(sig_top), in_(body_w), in_(1.25)
    )
    tf = sig.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.02)

    tf.text = "Sincerely,"
    p_sinc = tf.paragraphs[0]
    style_run(p_sinc.runs[0], 12, color=BLACK_TEXT)
    p_sinc.space_after = Pt(4)

    p_name = tf.add_paragraph()
    p_name.text = signer
    style_run(p_name.runs[0], 13, bold=True, italic=True, color=NAVY)

    p_title = tf.add_paragraph()
    p_title.text = signer_title
    style_run(p_title.runs[0], 10.5, color=TITLE_GRAY)


# ---------------------------------------------------------------------------
# Hook slide builders
# ---------------------------------------------------------------------------

def clear_non_placeholder_shapes(slide):
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        spTree.remove(shape._element)
        removed += 1
    return removed


def build_hook_2a(slide):
    """Slide 2a — single centered Robert Keller -> Nick Calder letter.

    No chevron, no extras. Just the letter, big and dignified.
    """
    # Card: wider + taller than the 2-letter versions, centered on a 13.33"
    # slide. Positioning: card centered at x=3.92, y=0.90, 5.50 wide, 6.20 tall.
    build_letter_card(
        slide,
        x_in=3.92, y_in=0.90, w_in=5.50, h_in=6.20,
        student="Nick Calder",
        signer="Robert Keller",
        body_font_size=14,
    )


def build_hook_2b(slide):
    """Slide 2b — two letters side-by-side, SUCCEEDED chevron in middle.

    Robert Keller -> Nick Calder  (with SUCCEEDED one year later)
    Robert Keller -> Jose Cervantez  (the new letter)
    Bottom question bar.
    """
    # Same card geometry as the existing (old) hook so it feels continuous.
    CARD_TOP = 0.55
    CARD_W = 4.70
    CARD_H = 5.55

    build_letter_card(
        slide,
        x_in=0.55, y_in=CARD_TOP, w_in=CARD_W, h_in=CARD_H,
        student="Nick Calder",
        signer="Robert Keller",
    )
    build_letter_card(
        slide,
        x_in=8.08, y_in=CARD_TOP, w_in=CARD_W, h_in=CARD_H,
        student="Jose Cervantez",
        signer="Robert Keller",
    )

    # Transition chevron
    add_right_arrow(slide, 5.47, 2.15, 2.36, 1.70,
                    fill=BAR_FILL, outline=NAVY, outline_width=1.75)

    add_textbox(slide, 5.47, 2.30, 1.80, 0.25,
                "ONE OUTCOME",
                size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    pill = add_rect(slide, 5.62, 2.65, 1.50, 0.48,
                    fill=PILL_BG, rounded=True)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = "SUCCEEDED"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 14, bold=True, color=PILL_TEXT)

    add_textbox(slide, 5.47, 3.22, 1.80, 0.26,
                "new letter about Jose",
                size=9, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Bottom question bar
    bar = add_rect(slide, 0.55, 6.30, 11.70, 0.60,
                   fill=BAR_FILL, rounded=True)
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.text = "How do you read the new letter after one successful outcome?"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 22, bold=True, color=NAVY)


def build_hook_2c(slide):
    """Slide 2c — identical to 2b but signatures flipped to Susan Keller.

    Pink highlight behind each Susan signature draws the audience's eye to
    the one thing that changed. New bottom bar asks the gender question.
    """
    CARD_TOP = 0.55
    CARD_W = 4.70
    CARD_H = 5.55

    build_letter_card(
        slide,
        x_in=0.55, y_in=CARD_TOP, w_in=CARD_W, h_in=CARD_H,
        student="Nick Calder",
        signer="Susan Keller",
        signer_highlight=True,
    )
    build_letter_card(
        slide,
        x_in=8.08, y_in=CARD_TOP, w_in=CARD_W, h_in=CARD_H,
        student="Jose Cervantez",
        signer="Susan Keller",
        signer_highlight=True,
    )

    # Same chevron + SUCCEEDED pill (the outcome is the same — only the signer changed)
    add_right_arrow(slide, 5.47, 2.15, 2.36, 1.70,
                    fill=BAR_FILL, outline=NAVY, outline_width=1.75)
    add_textbox(slide, 5.47, 2.30, 1.80, 0.25,
                "ONE OUTCOME",
                size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    pill = add_rect(slide, 5.62, 2.65, 1.50, 0.48,
                    fill=PILL_BG, rounded=True)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = "SUCCEEDED"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 14, bold=True, color=PILL_TEXT)
    add_textbox(slide, 5.47, 3.22, 1.80, 0.26,
                "new letter about Jose",
                size=9, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Gender-flip question bar
    bar = add_rect(slide, 0.55, 6.30, 11.70, 0.60,
                   fill=BAR_FILL, rounded=True)
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.text = "Does one success buy Susan the same credibility it bought Robert?"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 22, bold=True, color=NAVY)


# ---------------------------------------------------------------------------
# New slide provisioning — uses python-pptx add_slide with the "No Title"
# layout so the slide-number placeholder is inherited from the layout.
# ---------------------------------------------------------------------------

def add_blank_hook_slide(prs):
    """Add a new slide using the 'No Title' layout and strip the unwanted
    content/footer placeholders, leaving only the slide-number placeholder
    (which we want to keep for auto-updating numbering)."""
    layout = None
    for lyt in prs.slide_layouts:
        if lyt.name == "No Title":
            layout = lyt
            break
    if layout is None:
        # Fallback — shouldn't happen with this template
        layout = prs.slide_layouts[11]

    slide = prs.slides.add_slide(layout)

    # Remove content and footer placeholders, keep slide-number (idx=11)
    spTree = slide.shapes._spTree
    for shape in list(slide.placeholders):
        try:
            idx = shape.placeholder_format.idx
        except Exception:
            continue
        if idx == 11:
            continue  # keep slide number
        spTree.remove(shape._element)

    return slide


def move_slide_to(prs, slide, target_index):
    """Move slide to target_index (0-based) in the deck's slide order.

    Works by relocating the <p:sldId> element inside <p:sldIdLst>. This is
    the standard python-pptx workaround for moving slides. No XML is
    rewritten — the slide element's position in the parent list is just
    shifted, which PowerPoint interprets as a reorder.
    """
    sldIdLst = prs.slides._sldIdLst
    slide_id_value = slide.slide_id
    match = None
    for sldId in sldIdLst:
        # sldId element has an 'id' attribute
        if int(sldId.get("id")) == slide_id_value:
            match = sldId
            break
    if match is None:
        raise ValueError(f"Could not locate sldId for slide_id={slide_id_value}")

    sldIdLst.remove(match)

    children = list(sldIdLst)
    if target_index >= len(children):
        sldIdLst.append(match)
    else:
        sldIdLst.insert(target_index, match)


# ---------------------------------------------------------------------------
# Slide-5 RQ text rewrite
# ---------------------------------------------------------------------------

NEW_RQ_BODIES = [
    # run[1] text for RQ1, RQ2, RQ3 in the NEW order
    "Are male and female sponsors trusted equally before any outcome is observed?",
    "When a sponsor\u2019s prot\u00e9g\u00e9 succeeds or fails, does the sponsor\u2019s gender shape how audiences react?",
    "Does it matter how confident the sponsor was when they made the endorsement?",
]


def rewrite_slide5_rqs(slide):
    """Rewrite the 3 RQ cards' bodies in place.

    Structure assumed (verified 2026-04-08):
    - Each RQ card is an AUTO_SHAPE whose text_frame has one paragraph
      with two runs:  run[0] = "RQx    " (red bold), run[1] = question text (navy bold).
    - We only touch run[1].text, preserving all styling.
    """
    rq_cards = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith("RQ"):
            rq_cards.append(shape)
    if len(rq_cards) != 3:
        raise RuntimeError(
            f"Expected 3 RQ cards on slide 5, found {len(rq_cards)}. "
            f"Aborting to avoid writing to the wrong shapes."
        )
    for card, new_body in zip(rq_cards, NEW_RQ_BODIES):
        para = card.text_frame.paragraphs[0]
        if len(para.runs) < 2:
            raise RuntimeError(
                f"RQ card {card.text_frame.text!r} has fewer than 2 runs; "
                f"can't safely rewrite without losing formatting."
            )
        # Leave run[0] ("RQx    ") alone, replace run[1] body text
        para.runs[1].text = new_body


# ---------------------------------------------------------------------------
# Result-slide title rewrite (one run each, straightforward swap)
# ---------------------------------------------------------------------------

TITLE_REWRITES = {
    # Key: old title prefix we search for, Value: new title text
    "RQ1 — Gender \u00d7 outcome on trust update":
        "RQ2 — Gender \u00d7 outcome on trust update",
    "RQ2 — Strong endorsers":
        "RQ3 — Strong endorsers",
    "RQ2 — Weak endorsers":
        "RQ3 — Weak endorsers",
    "RQ3 — Pre-outcome trust by gender":
        "RQ1 — Pre-outcome trust by gender",
}


def rewrite_result_slide_title(slide):
    """Find the slide's title textbox (first shape whose text starts with 'RQ')
    and rewrite it according to TITLE_REWRITES. Preserves font styling by
    only touching run[0].text."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        current = shape.text_frame.text or ""
        for old_prefix, new_title in TITLE_REWRITES.items():
            if current.strip() == old_prefix:
                para = shape.text_frame.paragraphs[0]
                if not para.runs:
                    continue
                para.runs[0].text = new_title
                return old_prefix, new_title
    return None, None


# ---------------------------------------------------------------------------
# Picture swap for the new-RQ1 / old-slide-19 figure
# ---------------------------------------------------------------------------

def swap_main_picture(slide, new_png_path):
    """Find the largest picture on the slide, record its geometry, remove it,
    and add the new PNG at the same bounding box. This matches the pattern
    used in patch_2026_04_07.py."""
    from pptx.shapes.picture import Picture

    pics = [s for s in slide.shapes if isinstance(s, Picture)]
    if not pics:
        raise RuntimeError("No picture found on slide to swap")

    # Largest by area
    pics.sort(key=lambda s: s.width * s.height, reverse=True)
    target = pics[0]

    left, top, width, height = target.left, target.top, target.width, target.height
    spTree = slide.shapes._spTree
    spTree.remove(target._element)

    slide.shapes.add_picture(
        str(new_png_path),
        left=left, top=top, width=width, height=height,
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides before patch")

    # Grab references to slides we'll touch, using ORIGINAL 0-based indices
    s_hook    = slides[1]   # slide 2 (hook — will be rebuilt as 2a)
    s_rqs     = slides[4]   # slide 5 (research questions)
    s_old_rq1 = slides[15]  # old slide 16 — gender×outcome  -> new RQ2
    s_old_rq2_strong = slides[16]  # old slide 17 — strong  -> new RQ3 strong
    s_old_rq2_weak   = slides[17]  # old slide 18 — weak    -> new RQ3 weak
    s_old_rq3 = slides[18]  # old slide 19 — pre-outcome   -> new RQ1

    # --- Phase 1: rewrite slide 5 RQ bodies -----------------------------------
    print("\n[Phase 1] Rewriting slide 5 research questions (new order)")
    rewrite_slide5_rqs(s_rqs)
    print("  OK  — RQ bodies swapped in place")

    # --- Phase 2: rebuild slide 2 as 2a (single centered letter) --------------
    print("\n[Phase 2] Rebuilding slide 2 as hook 2a (Robert -> Nick, single letter)")
    removed = clear_non_placeholder_shapes(s_hook)
    print(f"  cleared {removed} shapes")
    build_hook_2a(s_hook)
    print("  built single centered letter card")

    # --- Phase 3: rewrite result-slide titles in place ------------------------
    print("\n[Phase 3] Rewriting result-slide titles")
    for s in (s_old_rq1, s_old_rq2_strong, s_old_rq2_weak, s_old_rq3):
        old, new = rewrite_result_slide_title(s)
        if old is None:
            print(f"  WARN: no matching title on a result slide (skipped)")
        else:
            print(f"  {old!r}\n    -> {new!r}")

    # --- Phase 4: swap RQ3 figure with the transposed version -----------------
    print(f"\n[Phase 4] Swapping picture on old slide 19 with {RQ3_FIG.name}")
    if not RQ3_FIG.exists():
        raise FileNotFoundError(f"Missing {RQ3_FIG} — run make_talk_figures.R first")
    swap_main_picture(s_old_rq3, RQ3_FIG)
    print("  picture swapped")

    # --- Phase 5: add 2 new hook slides at the end of the deck ----------------
    print("\n[Phase 5] Adding hook slides 2b and 2c (appended to deck for now)")
    s_2b = add_blank_hook_slide(prs)
    build_hook_2b(s_2b)
    print("  built 2b (Robert -> Nick + Jose, SUCCEEDED)")
    s_2c = add_blank_hook_slide(prs)
    build_hook_2c(s_2c)
    print("  built 2c (signatures flipped to Susan Keller)")

    # --- Phase 6: reorder via sldIdLst ---------------------------------------
    print("\n[Phase 6] Reordering slides")
    # Move 2b to position 2 (after the cover+hook-2a)
    move_slide_to(prs, s_2b, 2)
    print("  moved 2b -> index 2")
    # Move 2c to position 3 (right after 2b)
    move_slide_to(prs, s_2c, 3)
    print("  moved 2c -> index 3")

    # Now move the pre-outcome slide to the first result-slide position.
    # After the hook inserts, old slide 19 (pre-outcome) is at index 20.
    # We want it at index 17 (before old slide 16, which is now at 17).
    move_slide_to(prs, s_old_rq3, 17)
    print("  moved pre-outcome (new RQ1) -> index 17")

    # --- Save -----------------------------------------------------------------
    out_path = DECK
    print(f"\nSaving {out_path}")
    prs.save(str(out_path))
    print("Done.")


if __name__ == "__main__":
    main()
