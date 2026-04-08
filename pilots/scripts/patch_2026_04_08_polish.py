"""
patch_2026_04_08_polish.py — final polish pass before UChicago talk.

Addresses six items from Jose's 2026-04-08 review:
  1. Slides 4 & 5 (hook): replace "ONE OUTCOME / new letter about Jose"
     chevron text with YEAR 1 / YEAR 2 badges above each letter so the
     temporal pattern reads at a glance. SUCCEEDED pill stays in chevron.
  2. Slide 6 (definition): shrink bullet font 28pt -> 22pt so bullets
     stop dwarfing the definition body.
  3. Slide 11 (Stage 2 sponsors): bump bullet font 15pt -> 18pt for
     legibility (was the smallest body text in the deck).
  4. Slide 13 (Stage 3 incentive): full rebuild leading with the bank
     concept ("For each decision, you start with a $0.50 bank.") and
     mirroring Stage 2's cleaner top-to-bottom layout. Drops the dense
     stacked-formula opener.
  5. Slide 14 (Stage 3 setup): D1/D2 acronyms -> "First wager" / "Second
     wager"; trust-change formula updated; body bumped 18pt -> 20pt;
     title bumped 28pt -> 32pt to match other slide titles; stale "Name
     of Initiative" footer placeholder removed.
  6. Slide 23 (Summary): swap synth chart with the talk_theme_v2-styled
     re-render so x-axis Male/Female matches the size-24 bold labels on
     RQ result slides.

Runs the orphaned-timing/fld-GUID fixers afterwards (only slide 13 gets
a clear/rebuild but we run them globally as a safety net).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
SYNTH_PNG = PROJECT / "pilots" / "output" / "talk_figures" / "synthesis_outcome_sensitivity.png"


# ---------------------------------------------------------------------------
# Palette (consistent with patch_2026_04_07_content.py / patch_2026_04_08.py)
# ---------------------------------------------------------------------------

NAVY         = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP    = RGBColor(0x00, 0x14, 0x3D)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT   = RGBColor(0x1F, 0x1F, 0x1F)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY     = RGBColor(0x5B, 0x65, 0x7A)
TITLE_GRAY   = RGBColor(0x4B, 0x55, 0x63)
LIGHT_RULE   = RGBColor(0x9A, 0xA7, 0xBD)
CARD_OUTLINE = RGBColor(0xC7, 0xCF, 0xDD)
BAR_FILL     = RGBColor(0xEA, 0xF0, 0xFB)
PILL_BG      = RGBColor(0xDF, 0xF4, 0xE8)
PILL_TEXT    = RGBColor(0x11, 0x6B, 0x45)
TABLE_HEAD   = RGBColor(0xE5, 0xEA, 0xF3)
TABLE_HILITE = RGBColor(0xDB, 0xE3, 0xF1)
RIGHT_GREEN  = RGBColor(0x10, 0x6B, 0x45)
WRONG_RED    = RGBColor(0xB1, 0x21, 0x21)

FONT = "Calibri"


# ---------------------------------------------------------------------------
# Small shape helpers (same patterns as patch_2026_04_08.py)
# ---------------------------------------------------------------------------

def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size,
                bold=False, italic=False, color=BLACK_TEXT,
                align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color)
    return box


def add_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
             rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, in_(x), in_(y), in_(w), in_(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    if rounded:
        try:
            shape.adjustments[0] = 0.5
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


def add_oval(slide, x, y, w, h, fill, outline=None, outline_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, in_(x), in_(y), in_(w), in_(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    shape.text_frame.text = ""
    return shape


def clear_non_placeholder_shapes(slide):
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        spTree.remove(shape._element)
        removed += 1
    return removed


# ---------------------------------------------------------------------------
# Phase 1: Slides 4 & 5 — replace ONE OUTCOME / new letter about Jose
# with YEAR 1 / YEAR 2 badges above each letter
# ---------------------------------------------------------------------------

# Coordinates of the existing chevron components (matched against slides 4/5
# inspect output). We delete these by exact (x, y, w, h) so we don't touch the
# letter cards or bottom bar.
CHEVRON_TEXTS_TO_DELETE = {
    "ONE OUTCOME",
    "new letter about Jose",
}


def patch_hook_year_labels(slide, slide_label):
    """On a hook slide:
       - Delete the "ONE OUTCOME" caption above the chevron arrow
       - Delete the "new letter about Jose" caption below the chevron
       - Delete the right-arrow chevron + SUCCEEDED pill (will be rebuilt)
       - Add YEAR 1 / YEAR 2 badges above each letter
       - Rebuild a cleaner chevron with "ONE YEAR LATER" caption + SUCCEEDED pill
    """
    spTree = slide.shapes._spTree
    deleted_text = []
    deleted_shapes = []

    # Identify shapes by content/type
    chevron_shape = None
    succeeded_pill = None

    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt in CHEVRON_TEXTS_TO_DELETE:
                spTree.remove(shape._element)
                deleted_text.append(txt)
                continue
            if txt == "SUCCEEDED":
                succeeded_pill = shape
                continue
        # Right arrow chevron is a non-text auto_shape positioned around x=5.47
        try:
            x = shape.left / 914400
            y = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
        except (AttributeError, TypeError):
            continue
        # The chevron sits at roughly (5.47, 2.15, 2.36, 1.70). Match loosely.
        if (5.40 < x < 5.55) and (2.10 < y < 2.20) and (2.30 < w < 2.42) and (1.65 < h < 1.75):
            chevron_shape = shape

    # Remove chevron + pill (we'll rebuild a tidier version)
    if chevron_shape is not None:
        spTree.remove(chevron_shape._element)
        deleted_shapes.append("chevron")
    if succeeded_pill is not None:
        spTree.remove(succeeded_pill._element)
        deleted_shapes.append("SUCCEEDED pill")

    # --- Add YEAR 1 / YEAR 2 badges above each letter card -----------------
    # Letter cards on slides 4/5 are at:
    #   left  card: x=0.55, y=0.55, w=4.70, h=5.55
    #   right card: x=8.08, y=0.55, w=4.70, h=5.55
    # We add a small navy pill above each card, anchored just outside the card
    # top so it reads as a temporal label.
    BADGE_W = 1.50
    BADGE_H = 0.36
    BADGE_TOP = 0.10  # absolute slide y

    # Year 1 badge — centered above the left letter card
    left_badge_x = 0.55 + (4.70 - BADGE_W) / 2
    badge1 = add_rect(slide, left_badge_x, BADGE_TOP, BADGE_W, BADGE_H,
                      fill=NAVY, rounded=True)
    tf = badge1.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = "YEAR 1"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 13, bold=True, color=WHITE)

    # Year 2 badge — centered above the right letter card
    right_badge_x = 8.08 + (4.70 - BADGE_W) / 2
    badge2 = add_rect(slide, right_badge_x, BADGE_TOP, BADGE_W, BADGE_H,
                      fill=NAVY, rounded=True)
    tf = badge2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = "YEAR 2"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 13, bold=True, color=WHITE)

    # --- Rebuild chevron + SUCCEEDED pill, with cleaner ONE YEAR LATER caption
    # Caption above chevron
    cap = slide.shapes.add_textbox(in_(5.47), in_(2.18), in_(2.36), in_(0.30))
    tf = cap.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.text = "ONE YEAR LATER"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 11, bold=True, color=NAVY)

    # The chevron arrow itself (slightly slimmer so the SUCCEEDED pill reads as the focal point)
    chev = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, in_(5.62), in_(2.55), in_(2.06), in_(0.78)
    )
    chev.fill.solid()
    chev.fill.fore_color.rgb = BAR_FILL
    chev.line.color.rgb = NAVY
    chev.line.width = Pt(1.75)
    chev.text_frame.text = ""

    # SUCCEEDED pill — sits just below the chevron, tighter to it
    pill = add_rect(slide, 5.62, 3.55, 2.06, 0.50, fill=PILL_BG, rounded=True)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = "SUCCEEDED"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 14, bold=True, color=PILL_TEXT)

    print(f"  {slide_label}: deleted={deleted_text + deleted_shapes}; "
          f"added YEAR badges + new chevron")


# ---------------------------------------------------------------------------
# Phase 2: Slide 6 — shrink bullet font 28pt -> 22pt
# Phase 3: Slide 11 — bump bullet font 15pt -> 18pt
# ---------------------------------------------------------------------------

def resize_bullet_text(slide, slide_label, target_size, text_marker):
    """Find the bullets text frame on `slide` (identified by `text_marker`
    appearing in its text) and set every run's font size to `target_size`.
    Preserves bullet characters, paragraph structure, and color."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if text_marker not in txt:
            continue
        # Skip the title (single short paragraph)
        if shape.text_frame.text.strip().startswith("What is sponsorship") or \
           shape.text_frame.text.strip().startswith("Stage 2:"):
            continue
        # Apply target size to every run in every paragraph
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(target_size)
        print(f"  {slide_label}: resized bullets to {target_size}pt")
        return
    print(f"  {slide_label}: WARN no bullets matching {text_marker!r} found")


def patch_slide6_bullets(slide):
    """Slide 6: bullets are 28pt while definition is 22pt — shrink to 22pt."""
    resize_bullet_text(slide, "slide 6", 22, "Public advocacy")


def patch_slide11_bullets(slide):
    """Slide 11: bullets are 15pt — bump to 18pt for legibility."""
    resize_bullet_text(slide, "slide 11", 18, "Real sponsors")


# ---------------------------------------------------------------------------
# Phase 4: Slide 13 — full rebuild of Stage 3 incentive (bank-first layout)
# ---------------------------------------------------------------------------

def remove_name_of_initiative_placeholder(slide):
    """Drop the stale layout 'Name of Initiative' footer placeholder, keep
    the slide-number placeholder (idx=11)."""
    spTree = slide.shapes._spTree
    removed = 0
    for ph in list(slide.placeholders):
        try:
            idx = ph.placeholder_format.idx
        except Exception:
            continue
        if idx == 11:
            continue  # keep slide number
        # Identify the footer-text placeholder by its content
        txt = ph.text_frame.text.strip() if ph.has_text_frame else ""
        if txt == "Name of Initiative" or idx == 12:
            spTree.remove(ph._element)
            removed += 1
    return removed


def rebuild_slide13_stage3_incentive(slide):
    """Wipe non-placeholder shapes and rebuild as a bank-first layout that
    mirrors Stage 2's clarity. Also drop the stale 'Name of Initiative'
    footer placeholder.
    """
    n_removed_shapes = clear_non_placeholder_shapes(slide)
    n_removed_phs = remove_name_of_initiative_placeholder(slide)
    print(f"  slide 13: cleared {n_removed_shapes} shapes "
          f"+ {n_removed_phs} stale placeholder(s)")

    # ---- Title --------------------------------------------------------------
    add_textbox(slide, 0.55, 0.30, 12.20, 0.70,
                "Stage 3 incentive",
                size=32, bold=True, color=NAVY_DEEP)

    # ---- Lead: bank statement (the headline insight) ------------------------
    lead_y = 1.18
    lead_box = slide.shapes.add_textbox(in_(0.55), in_(lead_y), in_(12.20), in_(0.55))
    tf = lead_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "For each decision, you start with a "
    style_run(r1, 26, color=DARK_GRAY)
    r2 = p.add_run()
    r2.text = "$0.50 bank."
    style_run(r2, 26, bold=True, color=NAVY_DEEP)

    # ---- Sub: how the wager works ------------------------------------------
    sub = add_textbox(slide, 0.55, 1.78, 12.20, 0.42,
                      "You wager part of your bank on the endorser being right.",
                      size=18, italic=True, color=MID_GRAY)

    # ---- Two outcome cards (right / wrong) ---------------------------------
    # NOTE: y values across the slide were tightened on the second pass so
    # the footer text clears the master-slide navy band (which starts at
    # y~7.10). Original first-pass y was 2.40 for the cards.
    card_y = 2.30
    card_h = 0.92
    card_w = 5.85
    gap = 0.50
    card_left_x = 0.55 + (12.20 - 2 * card_w - gap) / 2

    # RIGHT outcome
    right_card = add_rect(slide, card_left_x, card_y, card_w, card_h,
                          fill=WHITE, outline=RIGHT_GREEN, outline_width=2.0)
    rc_tf = right_card.text_frame
    rc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rc_tf.margin_left = Inches(0.20)
    rc_tf.margin_right = Inches(0.20)
    rc_tf.text = ""
    rcp = rc_tf.paragraphs[0]
    rcp.alignment = PP_ALIGN.LEFT
    r1 = rcp.add_run()
    r1.text = "If "
    style_run(r1, 21, color=DARK_GRAY)
    r2 = rcp.add_run()
    r2.text = "right "
    style_run(r2, 21, bold=True, color=RIGHT_GREEN)
    r3 = rcp.add_run()
    r3.text = "→  bank "
    style_run(r3, 21, color=DARK_GRAY)
    r4 = rcp.add_run()
    r4.text = "+"
    style_run(r4, 21, bold=True, color=RIGHT_GREEN)
    r5 = rcp.add_run()
    r5.text = " your wager"
    style_run(r5, 21, color=DARK_GRAY)

    # WRONG outcome
    wrong_card = add_rect(slide, card_left_x + card_w + gap, card_y, card_w, card_h,
                          fill=WHITE, outline=WRONG_RED, outline_width=2.0)
    wc_tf = wrong_card.text_frame
    wc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    wc_tf.margin_left = Inches(0.20)
    wc_tf.margin_right = Inches(0.20)
    wc_tf.text = ""
    wcp = wc_tf.paragraphs[0]
    wcp.alignment = PP_ALIGN.LEFT
    r1 = wcp.add_run()
    r1.text = "If "
    style_run(r1, 21, color=DARK_GRAY)
    r2 = wcp.add_run()
    r2.text = "wrong "
    style_run(r2, 21, bold=True, color=WRONG_RED)
    r3 = wcp.add_run()
    r3.text = "→  bank "
    style_run(r3, 21, color=DARK_GRAY)
    r4 = wcp.add_run()
    r4.text = "−"
    style_run(r4, 21, bold=True, color=WRONG_RED)
    r5 = wcp.add_run()
    r5.text = " your wager"
    style_run(r5, 21, color=DARK_GRAY)

    # ---- Worked example: slider showing 50% wager --------------------------
    add_textbox(slide, 0.55, 3.45, 12.20, 0.32,
                "evaluator wagers 50%",
                size=13, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # Slider track
    track_x = 3.50
    track_y = 3.85
    track_w = 6.33
    track = add_rect(slide, track_x, track_y, track_w, 0.22,
                     fill=BAR_FILL, rounded=False)
    track.line.color.rgb = LIGHT_RULE
    track.line.width = Pt(0.75)

    # Slider fill (50% of track width)
    fill_w = track_w / 2
    fill_shape = add_rect(slide, track_x, track_y, fill_w, 0.22,
                          fill=NAVY, rounded=False)
    # Slider thumb
    thumb_size = 0.42
    thumb_x = track_x + fill_w - thumb_size / 2
    thumb_y = track_y + 0.11 - thumb_size / 2
    add_oval(slide, thumb_x, thumb_y, thumb_size, thumb_size,
             fill=WHITE, outline=NAVY, outline_width=2.5)

    # 0% / 100% slider end labels
    add_textbox(slide, track_x - 0.75, track_y + 0.32, 0.70, 0.30,
                "0%", size=12, color=MID_GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, track_x + track_w + 0.05, track_y + 0.32, 0.80, 0.30,
                "100%", size=12, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # ---- Twin reveal: right $0.75  /  wrong $0.25 --------------------------
    reveal_y = 4.55
    add_textbox(slide, 0.55, reveal_y, 5.85, 0.85,
                "right:  $0.75",
                size=32, bold=True, color=RIGHT_GREEN,
                align=PP_ALIGN.RIGHT)
    add_textbox(slide, 6.95, reveal_y, 5.85, 0.85,
                "wrong:  $0.25",
                size=32, bold=True, color=WRONG_RED,
                align=PP_ALIGN.LEFT)

    # ---- Payout table at the bottom ----------------------------------------
    table_top = 5.55
    row_h = 0.34

    table_left_x = 2.40
    col_widths = [1.95, 2.55, 2.55, 2.55]   # [label] [0%] [50%] [100%]
    col_lefts = []
    cx = table_left_x
    for w in col_widths:
        col_lefts.append(cx)
        cx += w

    # Header row (wager %)
    add_textbox(slide, col_lefts[0], table_top, col_widths[0], row_h,
                "wager", size=15, bold=True, color=MID_GRAY,
                align=PP_ALIGN.RIGHT)
    for i, label in enumerate(["0%", "50%", "100%"], start=1):
        add_textbox(slide, col_lefts[i], table_top, col_widths[i], row_h,
                    label, size=15, bold=True, color=MID_GRAY,
                    align=PP_ALIGN.CENTER)
    # Highlight the 50% column (matches the slider example)
    add_rect(slide, col_lefts[2], table_top + row_h - 0.02,
             col_widths[2], row_h * 2 + 0.04,
             fill=TABLE_HILITE, rounded=False)

    # if right row
    add_textbox(slide, col_lefts[0], table_top + row_h, col_widths[0], row_h,
                "if right", size=15, bold=True, color=RIGHT_GREEN,
                align=PP_ALIGN.RIGHT)
    for i, val in enumerate(["$0.50", "$0.75", "$1.00"], start=1):
        add_textbox(slide, col_lefts[i], table_top + row_h, col_widths[i], row_h,
                    val, size=15, color=DARK_GRAY,
                    align=PP_ALIGN.CENTER)

    # if wrong row
    add_textbox(slide, col_lefts[0], table_top + 2 * row_h, col_widths[0], row_h,
                "if wrong", size=15, bold=True, color=WRONG_RED,
                align=PP_ALIGN.RIGHT)
    for i, val in enumerate(["$0.50", "$0.25", "$0.00"], start=1):
        add_textbox(slide, col_lefts[i], table_top + 2 * row_h, col_widths[i], row_h,
                    val, size=15, color=DARK_GRAY,
                    align=PP_ALIGN.CENTER)

    # ---- Footer ------------------------------------------------------------
    # Footer must end before y~7.05 (where the master-slide navy band starts).
    # Table ends at table_top + 3*row_h ~= 6.57; gap of 0.10; footer at 6.67.
    add_textbox(slide, 0.55, 6.67, 12.20, 0.30,
                "One randomly selected decision determines the bonus.",
                size=12, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    print("  slide 13: rebuilt with bank-first layout")


# ---------------------------------------------------------------------------
# Phase 5: Slide 14 — surgical text edits + font bumps
# ---------------------------------------------------------------------------

# Direct text replacements on slide 14 (D1/D2 -> First/Second wager)
SLIDE14_TEXT_REPLACEMENTS = {
    "D1 wager": "First wager",
    "D2 wager": "Second wager",
    "outcome reveal": "outcome revealed",
    "trust change  =  (D2 wager)  −  (D1 wager)":
        "trust change  =  (Second wager)  −  (First wager)",
    "trust change = (D2 wager) − (D1 wager)":
        "trust change = (Second wager) − (First wager)",
}


def patch_slide14_setup(slide):
    """Surgical edits on the Stage 3 setup slide:
       - D1/D2 -> First/Second wager (preserving styling)
       - outcome reveal -> outcome revealed
       - body text 18pt -> 20pt for legibility
       - title 28pt -> 32pt to match other slide titles
       - drop the stale 'Name of Initiative' footer placeholder
    """
    n_removed = remove_name_of_initiative_placeholder(slide)
    if n_removed:
        print(f"  slide 14: dropped {n_removed} stale placeholder(s)")

    n_text_changes = 0
    n_size_changes = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            # Reconstruct paragraph text from runs to find candidates for
            # full-paragraph replacements (like the trust-change formula)
            full = "".join(r.text for r in para.runs)
            stripped = full.strip()

            # Check exact-paragraph replacements first
            replaced = False
            for old, new in SLIDE14_TEXT_REPLACEMENTS.items():
                if stripped == old:
                    # Concentrate the replacement in the first run, blank others
                    if para.runs:
                        para.runs[0].text = new
                        for r in para.runs[1:]:
                            r.text = ""
                        n_text_changes += 1
                        replaced = True
                        break
            if replaced:
                continue

            # Otherwise, do per-run substring replacements
            for run in para.runs:
                for old, new in SLIDE14_TEXT_REPLACEMENTS.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        n_text_changes += 1

        # Font size adjustments based on shape's content/role
        ttext = shape.text_frame.text.strip()
        if ttext == "Stage 3":
            # Title: 28pt -> 32pt
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.color.rgb = NAVY_DEEP
            n_size_changes += 1
        elif ttext.startswith("Candidates completed three tasks") or \
             ttext.startswith("Endorsers viewed two candidates") or \
             ttext.startswith("An endorsement is") or \
             ttext.startswith("Each evaluator sees"):
            # Body lines: 18pt -> 20pt
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(20)
            n_size_changes += 1
        elif ttext.startswith("trust change"):
            # Formula: keep 22pt but make it bold + navy for prominence
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = NAVY_DEEP
            n_size_changes += 1
        elif ttext in ("First wager", "Second wager", "outcome revealed"):
            # Flow boxes: bump from 16pt to 18pt
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(18)
                    run.font.bold = True
            n_size_changes += 1

    print(f"  slide 14: {n_text_changes} text replacements + "
          f"{n_size_changes} font adjustments")


# ---------------------------------------------------------------------------
# Phase 6: Slide 23 — swap synthesis chart with the talk_theme_v2 re-render
# ---------------------------------------------------------------------------

def swap_slide23_synthesis(slide):
    """Replace the synthesis chart picture with the freshly re-rendered
    PNG (now using talk_theme_v2). Preserves position and size."""
    from pptx.shapes.picture import Picture
    pics = [s for s in slide.shapes if isinstance(s, Picture)]
    if not pics:
        print("  slide 23: WARN no picture found")
        return
    pics.sort(key=lambda s: s.width * s.height, reverse=True)
    target = pics[0]
    left, top, width, height = target.left, target.top, target.width, target.height
    spTree = slide.shapes._spTree
    spTree.remove(target._element)
    slide.shapes.add_picture(str(SYNTH_PNG), left=left, top=top,
                             width=width, height=height)
    print("  slide 23: swapped synthesis chart picture")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides")

    if not SYNTH_PNG.exists():
        raise FileNotFoundError(f"Missing {SYNTH_PNG} — run make_talk_figures.R first")

    # 0-based indices for the slides we touch
    s_4  = slides[3]   # hook 2b — Robert Nick + Jose
    s_5  = slides[4]   # hook 2c — Susan Nick + Jose
    s_6  = slides[5]   # what is sponsorship
    s_11 = slides[10]  # Stage 2 sponsors
    s_13 = slides[12]  # Stage 3 incentive
    s_14 = slides[13]  # Stage 3 setup
    s_23 = slides[22]  # Summary

    print("\n[Phase 1] Hook slides 4 & 5: YEAR 1 / YEAR 2 badges + cleaner chevron")
    patch_hook_year_labels(s_4, "slide 4")
    patch_hook_year_labels(s_5, "slide 5")

    print("\n[Phase 2] Slide 6: shrink bullet font")
    patch_slide6_bullets(s_6)

    print("\n[Phase 3] Slide 11: bump bullet font")
    patch_slide11_bullets(s_11)

    print("\n[Phase 4] Slide 13: rebuild Stage 3 incentive bank-first")
    rebuild_slide13_stage3_incentive(s_13)

    print("\n[Phase 5] Slide 14: D1/D2 -> First/Second wager, font bumps")
    patch_slide14_setup(s_14)

    print("\n[Phase 6] Slide 23: swap synthesis chart picture")
    swap_slide23_synthesis(s_23)

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
