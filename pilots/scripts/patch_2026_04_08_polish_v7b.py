"""
patch_2026_04_08_polish_v7b.py — slide 27 layout fix.

v7 set bullet_h=1.20 and header=20pt with LEFT_W=6.10. At 20pt bold the
longer headers wrapped to 2 lines, blowing past the per-bullet height
budget and overlapping the next bullet. v7b widens the left column to
7.70" (shrinking the right-column figure proportionally), shrinks the
header to 19pt so the longest header still fits 1 line, and preserves
all of Jose's bullet text without further trimming.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR = PROJECT / "pilots" / "output" / "talk_figures"
PRIOR_FIG = FIG_DIR / "prior_distributions.png"


NAVY       = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP  = RGBColor(0x00, 0x14, 0x3D)
DARK_GRAY  = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY   = RGBColor(0x5B, 0x65, 0x7A)
BLACK_TEXT = RGBColor(0x11, 0x18, 0x27)

FONT = "Arial"


# Hardcoded bullets — preserves Jose's edits to bullets 2 and 3 from
# the current deck state. NO em dashes anywhere. Bullet 1 + bullet 4
# headers shortened slightly so they fit on a single line at 19pt bold
# without colliding with the right-column figure.
BULLETS = [
    (
        "People hold different priors for men vs women.",
        "Men are assumed critical; women are assumed nicer on average.",
    ),
    (
        "Women's endorsements cluster high + narrow.",
        "A generous prior is hard to read signal from, and evaluators "
        "discount their endorsements (consistent with the effects "
        "among sponsors in the stronger confidence cells).",
    ),
    (
        "Men's endorsements spread wide.",
        "When men endorse with confidence, it carries information which "
        "can be both praised or penalized.",
    ),
    (
        "Next direction: telegraph their standards.",
        "Signaling \"I only recommend the top 10%\" could close the gap "
        "without years of seniority.",
    ),
]

# Layout — left column ends at 7.75 so the figure at x=7.85 has a 0.10
# breathing gap. v7b had LEFT_W=7.70 which made the textbox end at 8.25
# and the bullet body text bled UNDER the figure.
LEFT_X     = 0.55
LEFT_Y     = 1.55
LEFT_W     = 7.20
BULLET_H   = 1.20
BULLET_GAP = 0.10
HEADER_PT  = 19
BODY_PT    = 16

FIG_X = 7.85
FIG_Y = 1.55
FIG_W = 5.40
FIG_H = 3.32  # = 5.40 / (7.8/4.8 aspect ratio of source PNG)


def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
                color=BLACK_TEXT, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color, font=font)
    return box


def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)


def main():
    prs = Presentation(str(DECK))
    slide = prs.slides[26]  # slide 27

    # Drop everything except the slidenum placeholder; rebuild from scratch
    dropped = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        remove_shape(slide, shape)
        dropped += 1
    print(f"  slide 27: dropped {dropped} shape(s)")

    # ---- Title ----
    add_textbox(
        slide, 0.55, 0.30, 12.23, 0.70,
        "Running story",
        size=32, bold=True, color=NAVY,
    )

    # ---- Subtitle (no em dash) ----
    add_textbox(
        slide, 0.55, 1.05, 12.23, 0.42,
        "Different priors, not a gender penalty",
        size=20, italic=True, color=MID_GRAY,
    )

    # ---- Left column: 4 bullets at 19pt header / 16pt body ----
    for i, (head, body) in enumerate(BULLETS):
        y = LEFT_Y + i * (BULLET_H + BULLET_GAP)

        # Navy accent square
        mark = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            in_(LEFT_X), in_(y + 0.14),
            in_(0.16), in_(0.16),
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = NAVY
        mark.line.fill.background()
        mark.text_frame.text = ""

        box = slide.shapes.add_textbox(
            in_(LEFT_X + 0.28), in_(y),
            in_(LEFT_W - 0.28), in_(BULLET_H),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_top = Inches(0.0)

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = head
        style_run(r1, HEADER_PT, bold=True, color=NAVY_DEEP)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = body
        style_run(r2, BODY_PT, bold=False, color=DARK_GRAY)

    print(f"  slide 27: added {len(BULLETS)} bullet(s) at {HEADER_PT}/{BODY_PT}pt")

    # ---- Right column figure ----
    if PRIOR_FIG.exists():
        slide.shapes.add_picture(
            str(PRIOR_FIG),
            left=in_(FIG_X), top=in_(FIG_Y),
            width=in_(FIG_W), height=in_(FIG_H),
        )
        add_textbox(
            slide, FIG_X, FIG_Y + FIG_H + 0.02, FIG_W, 0.30,
            "Hypothesized prior distributions",
            size=12, italic=True, color=MID_GRAY,
            align=PP_ALIGN.CENTER,
        )
        print("  slide 27: prior figure + caption added")

    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
