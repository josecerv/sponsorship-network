"""
patch_2026_04_08_polish_v4b.py — strength range badge + RQ1 two-panel.

Items:
  1. Build a reusable "endorser confidence scale" badge that shows the
     full 0-100 slider range with the actual Strong (77-88 display) and
     Weak (12-23 display) bands highlighted. This replaces the old
     single-level badge that Jose had on slides 21/22 before the v3
     cleanup.

  2. Slide 17 (Stage 3 randomization): inside the "Endorsement strength"
     card, replace the two mini fill bars + "Strong / very confident"
     and "Weak / unsure" labels with the new combined range badge.

  3. Slide 23 (RQ3 strong) & slide 24 (RQ3 weak): re-add the badge in
     the top-left corner so the audience has the range context when
     looking at the result figures. Shrink the descriptive title from
     w=8.60 back to w=6.30 to make room.

  4. Slide 21 (RQ1): rebuild the body to be a 2-panel composite ("first
     show the overall, then split by gender"). Left panel =
     rq1_overall.png (2 bars: Strong 56 / Weak 27). Right panel =
     rq3_initial_trust.png (the transposed 4-bar chart). Adds a simple
     "OVERALL" / "BY GENDER" caption above each.

Surgical-only. fix_fld_guids/orphaned_timing not needed (no
clear_non_placeholder_shapes calls). Range values come from actual
study_data_clean.csv analysis: Q1 display strengths.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR = PROJECT / "pilots" / "output" / "talk_figures"
RQ1_OVERALL = FIG_DIR / "rq1_overall.png"
RQ1_SPLIT = FIG_DIR / "rq3_initial_trust.png"


# Actual Q1 display strength ranges (from study_data_clean.csv analysis)
STRONG_LO, STRONG_HI = 77, 88
WEAK_LO,   WEAK_HI   = 12, 23


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------

NAVY         = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP    = RGBColor(0x00, 0x14, 0x3D)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT   = RGBColor(0x1F, 0x1F, 0x1F)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY     = RGBColor(0x5B, 0x65, 0x7A)
LIGHT_GRAY   = RGBColor(0x9A, 0xA7, 0xBD)
CARD_OUTLINE = RGBColor(0xC7, 0xCF, 0xDD)
BAR_TRACK    = RGBColor(0xE5, 0xEA, 0xF3)
WRONG_RED    = RGBColor(0xB1, 0x21, 0x21)

FONT = "Calibri"


def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size,
                bold=False, italic=False, color=BLACK_TEXT,
                align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color)
    return box


def add_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
             rounded=False, corner=0.5):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, in_(x), in_(y), in_(w), in_(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    if rounded:
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


# ---------------------------------------------------------------------------
# Strength range badge helper
# ---------------------------------------------------------------------------

def draw_strength_range_badge(slide, x, y, w, h,
                              show_frame=True, header_size=10, label_size=13,
                              range_size=10):
    """Draw an 'endorser confidence scale' badge that visualizes the full
    0-100 slider range with highlighted Weak and Strong bands.

    Layout (proportional to h):
      top 20%: "ENDORSER CONFIDENCE" header (small gray)
      middle 30%: scale bar (track + 2 highlighted bands)
      bottom 50%: labels ("Weak 12-23" and "Strong 77-88") centered
                  under their respective bands

    Coordinates are in inches. Returns the list of shape objects added.
    """
    # Optional background frame
    if show_frame:
        add_rect(slide, x, y, w, h,
                 fill=WHITE, outline=CARD_OUTLINE, outline_width=1.0,
                 rounded=True, corner=0.10)

    margin = 0.12
    inner_x = x + margin
    inner_w = w - 2 * margin

    # Header
    header_y = y + 0.08
    header_h = 0.22
    add_textbox(slide, inner_x, header_y, inner_w, header_h,
                "ENDORSER CONFIDENCE",
                size=header_size, bold=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # Scale bar track
    bar_y = header_y + header_h + 0.12
    bar_h = 0.22
    track = add_rect(slide, inner_x, bar_y, inner_w, bar_h,
                     fill=BAR_TRACK, rounded=True, corner=0.5)
    track.line.color.rgb = LIGHT_GRAY
    track.line.width = Pt(0.75)

    # Weak band: positioned at WEAK_LO to WEAK_HI % of bar width
    weak_x = inner_x + inner_w * (WEAK_LO / 100.0)
    weak_w = inner_w * ((WEAK_HI - WEAK_LO) / 100.0)
    add_rect(slide, weak_x, bar_y, weak_w, bar_h,
             fill=WRONG_RED, rounded=True, corner=0.5)

    # Strong band
    strong_x = inner_x + inner_w * (STRONG_LO / 100.0)
    strong_w = inner_w * ((STRONG_HI - STRONG_LO) / 100.0)
    add_rect(slide, strong_x, bar_y, strong_w, bar_h,
             fill=NAVY, rounded=True, corner=0.5)

    # Scale endpoints: "0" and "100" below the bar
    endpoint_y = bar_y + bar_h + 0.02
    add_textbox(slide, inner_x - 0.10, endpoint_y, 0.35, 0.22,
                "0", size=9, color=MID_GRAY, align=PP_ALIGN.LEFT)
    add_textbox(slide, inner_x + inner_w - 0.25, endpoint_y, 0.35, 0.22,
                "100", size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    # Labels under the bands: "Weak 12-23" and "Strong 77-88"
    label_y = endpoint_y + 0.22
    # Center of each band, in inches
    weak_center = inner_x + inner_w * ((WEAK_LO + WEAK_HI) / 2 / 100.0)
    strong_center = inner_x + inner_w * ((STRONG_LO + STRONG_HI) / 2 / 100.0)
    # Label box width (approximate, allows the band-range text to fit)
    lw = 1.15
    add_textbox(slide, weak_center - lw / 2, label_y, lw, 0.32,
                "Weak", size=label_size, bold=True, color=WRONG_RED,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, weak_center - lw / 2, label_y + 0.28, lw, 0.26,
                f"{WEAK_LO}\u2013{WEAK_HI}",
                size=range_size, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, strong_center - lw / 2, label_y, lw, 0.32,
                "Strong", size=label_size, bold=True, color=NAVY_DEEP,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, strong_center - lw / 2, label_y + 0.28, lw, 0.26,
                f"{STRONG_LO}\u2013{STRONG_HI}",
                size=range_size, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Phase 1: Slide 17 — replace strength card content with combined badge
# ---------------------------------------------------------------------------

# Identify shapes inside the strength card (card x ~ 4.64-8.69) that
# represent the old mini-bars and stim labels we want to remove.
SLIDE17_STRENGTH_DELETE_TEXTS = {
    "Strong", "Weak", "very confident", "unsure",
}


def patch_slide17_strength_card(slide):
    """Remove the old 2-bar mini stim + labels inside the strength card
    and add a combined range badge below the italic sub-caption."""
    # The strength card sits at x in [4.64, 8.69]
    CARD_X_MIN = 4.64
    CARD_X_MAX = 8.69
    CARD_Y_MIN = 1.95
    CARD_Y_MAX = 6.80

    spTree = slide.shapes._spTree
    dropped = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        try:
            x = shape.left / 914400
            y = shape.top / 914400
        except Exception:
            continue
        if not (CARD_X_MIN <= x <= CARD_X_MAX):
            continue
        if not (CARD_Y_MIN <= y <= CARD_Y_MAX):
            continue
        # Keep the header and italic sub (they start at y ~ 2.10 / 2.60)
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            # Header & italic sub — keep
            if txt.startswith("2.  Endorsement strength") or \
               txt.startswith("Sponsors used a 0"):
                continue
            # Card background (the rounded white rect) — keep
            if txt == "":
                # Is this the card background? Its width is ~ 4.05
                try:
                    w = shape.width / 914400
                except Exception:
                    continue
                if w > 3.80:
                    continue
            # Old mini-bar fills (no text, narrower)
            if txt == "" and 3.80 > (shape.width / 914400) > 0.50:
                spTree.remove(shape._element)
                dropped += 1
                continue
            # Old labels
            if txt in SLIDE17_STRENGTH_DELETE_TEXTS:
                spTree.remove(shape._element)
                dropped += 1
                continue
        else:
            # Non-text shapes inside the card's stim area — drop
            try:
                w = shape.width / 914400
            except Exception:
                continue
            if w < 3.80:
                spTree.remove(shape._element)
                dropped += 1

    print(f"  slide 17: dropped {dropped} old strength stim element(s)")

    # Add new range badge inside the card, below the italic sub
    # Card at x=4.64, y=1.95, w=4.05, h=4.85
    # Italic sub is at y ~ 2.60-3.00. Badge starts y=3.35.
    badge_x = 4.64 + 0.20
    badge_y = 3.50
    badge_w = 4.05 - 0.40
    badge_h = 1.60
    draw_strength_range_badge(slide, badge_x, badge_y, badge_w, badge_h,
                              show_frame=False, header_size=11,
                              label_size=15, range_size=11)
    print("  slide 17: added combined strength-range badge")


# ---------------------------------------------------------------------------
# Phase 2: Slides 23, 24 — re-add badge top-left, shrink title to fit
# ---------------------------------------------------------------------------

def patch_rq3_slide_with_badge(slide, slide_label):
    """Shrink the descriptive title, add the range badge top-left."""
    # Find the title and shrink it. Current width is 8.60 (after v3).
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt.startswith("RQ3") and ("endorsements" in txt):
            # Shrink title: start at x=2.95, width=6.30
            shape.left = Inches(2.95)
            shape.width = Inches(6.30)
            shape.top = Inches(0.25)
            shape.height = Inches(1.40)
            break

    # Add badge at top-left
    draw_strength_range_badge(slide, 0.55, 0.25, 2.30, 1.50,
                              show_frame=True, header_size=9,
                              label_size=12, range_size=9)
    print(f"  {slide_label}: shrunk title + added range badge top-left")


# ---------------------------------------------------------------------------
# Phase 3: Slide 21 (RQ1) — rebuild as 2-panel composite
# ---------------------------------------------------------------------------

def patch_slide21_rq1_two_panel(slide):
    """Remove the current RQ1 figure; add two pictures side-by-side with
    OVERALL / BY GENDER captions so the audience sees the overall first,
    then the gender split.
    """
    # Drop the existing picture
    from pptx.shapes.picture import Picture
    pics = [s for s in slide.shapes if isinstance(s, Picture)]
    for p in pics:
        slide.shapes._spTree.remove(p._element)
    print(f"  slide 21: removed {len(pics)} existing picture(s)")

    # Layout: left panel small (overall) + right panel big (by gender)
    # Slide content area: y ~ 1.80 to 7.00, x ~ 0.10 to 13.23
    # Left (overall): x=0.30, y=2.00, w=4.20, h=4.80
    # Right (by gender): x=4.80, y=2.00, w=8.30, h=4.80
    # Captions above each panel at y=1.80 to 2.00

    # Left caption
    add_textbox(slide, 0.30, 1.75, 4.20, 0.28,
                "OVERALL",
                size=13, bold=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # Left picture (overall)
    if RQ1_OVERALL.exists():
        slide.shapes.add_picture(
            str(RQ1_OVERALL),
            left=Inches(0.50), top=Inches(2.10),
            width=Inches(3.80), height=Inches(4.90),
        )
        print("  slide 21: added OVERALL panel (rq1_overall.png)")

    # Right caption
    add_textbox(slide, 4.80, 1.75, 8.30, 0.28,
                "SAME DATA  ·  SPLIT BY SPONSOR GENDER",
                size=13, bold=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # Right picture (by gender - the transposed chart)
    if RQ1_SPLIT.exists():
        slide.shapes.add_picture(
            str(RQ1_SPLIT),
            left=Inches(4.70), top=Inches(2.10),
            width=Inches(8.55), height=Inches(4.90),
        )
        print("  slide 21: added BY GENDER panel (rq3_initial_trust.png)")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides")

    s_17 = slides[16]
    s_21 = slides[20]
    s_23 = slides[22]
    s_24 = slides[23]

    print("\n[Phase 1] Slide 17: replace strength card stim with range badge")
    patch_slide17_strength_card(s_17)

    print("\n[Phase 2] Slides 23 & 24: add range badge top-left")
    patch_rq3_slide_with_badge(s_23, "slide 23")
    patch_rq3_slide_with_badge(s_24, "slide 24")

    print("\n[Phase 3] Slide 21: RQ1 two-panel composite")
    patch_slide21_rq1_two_panel(s_21)

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
