"""
patch_2026_04_08_polish_v4c.py — add the instructions appendix.

Adds one divider slide + N picture slides at the very end of the deck,
each showing an instruction screen from the Stage 2 or Stage 3 Qualtrics
survey. The PNGs are pre-rendered by render_appendix_screens.py.

All new slides are appended to the end — the existing 28-slide main
deck is untouched. Deck grows 28 -> 28 + 1 + N.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
SCREENS = PROJECT / "pilots" / "output" / "appendix_screens"


# Ordered appendix list — must match render_appendix_screens.TARGETS
APPENDIX_SLIDES = [
    ("stage2_01_role.png",        "Stage 2  ·  Endorser role"),
    ("stage2_02_earnings.png",    "Stage 2  ·  Earnings ($1 flat + $3 bonus)"),
    ("stage2_03_task.png",        "Stage 2  ·  Task objective + slider"),
    ("stage3_01_purpose.png",     "Stage 3  ·  Study purpose"),
    ("stage3_02_setup.png",       "Stage 3  ·  What evaluators see"),
    ("stage3_03_motivation.png",  "Stage 3  ·  Endorser motivation"),
    ("stage3_04_two_decisions.png", "Stage 3  ·  Two-decision framing"),
    ("stage3_05_wager.png",       "Stage 3  ·  $0.50 bank wager mechanics"),
    ("stage3_cq1_bank.png",       "CQ1  ·  Bank amount"),
    ("stage3_cq2_wager.png",      "CQ2  ·  Wager calculation"),
    ("stage3_cq3_bonus.png",      "CQ3  ·  Bonus calculation"),
    ("stage3_cq4_task.png",       "CQ4  ·  Task question"),
    ("stage3_cq5_comparison.png", "CQ5  ·  Comparison question"),
]


NAVY       = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP  = RGBColor(0x00, 0x14, 0x3D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY   = RGBColor(0x5B, 0x65, 0x7A)
BLACK_TEXT = RGBColor(0x1F, 0x1F, 0x1F)


def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_blank_no_title_slide(prs):
    """Add a slide on the 'No Title' layout; drop non-slidenum placeholders."""
    layout = None
    for lyt in prs.slide_layouts:
        if lyt.name == "No Title":
            layout = lyt
            break
    if layout is None:
        layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(layout)
    spTree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        try:
            idx = ph.placeholder_format.idx
        except Exception:
            continue
        if idx == 11:
            continue
        spTree.remove(ph._element)
    return slide


def build_divider_slide(slide):
    """Big 'Appendix' title, small subtitle explaining what's in it."""
    # Title
    box = slide.shapes.add_textbox(in_(0.55), in_(2.60), in_(12.23), in_(1.20))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = "Appendix"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 64, bold=True, color=NAVY_DEEP)

    # Subtitle
    box2 = slide.shapes.add_textbox(in_(0.55), in_(3.90), in_(12.23), in_(0.80))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.text = "Instruction screens from the Stage 2 & Stage 3 Qualtrics surveys"
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    style_run(p2.runs[0], 22, italic=True, color=MID_GRAY)


def build_appendix_slide(slide, png_name, title):
    """Title bar + centered screenshot image."""
    # Title
    box = slide.shapes.add_textbox(in_(0.55), in_(0.30), in_(12.23), in_(0.60))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    style_run(p.runs[0], 24, bold=True, color=NAVY_DEEP)

    # Screenshot — scale to fit the content area (y=1.10 to y=7.00, x=0.55 to 12.78)
    png_path = SCREENS / png_name
    if not png_path.exists():
        print(f"  WARN missing {png_path}")
        return

    # We don't know the PNG aspect ratio without loading it; use PIL to check
    try:
        from PIL import Image as PILImage
        img = PILImage.open(png_path)
        img_w, img_h = img.size
    except Exception:
        img_w, img_h = 1200, 1600

    avail_w = 12.23
    avail_h = 5.90
    img_ratio = img_w / img_h
    avail_ratio = avail_w / avail_h
    if img_ratio > avail_ratio:
        # Width-constrained
        pic_w = avail_w
        pic_h = pic_w / img_ratio
    else:
        pic_h = avail_h
        pic_w = pic_h * img_ratio

    # Center horizontally; anchor near top of content area
    pic_x = 0.55 + (avail_w - pic_w) / 2
    pic_y = 1.05 + (avail_h - pic_h) / 2

    slide.shapes.add_picture(
        str(png_path),
        left=in_(pic_x), top=in_(pic_y),
        width=in_(pic_w), height=in_(pic_h),
    )


def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    n_before = len(prs.slides)
    print(f"  {n_before} slides before")

    # 1. Divider slide
    print("\n[1] Adding Appendix divider slide")
    divider = add_blank_no_title_slide(prs)
    build_divider_slide(divider)

    # 2. Appendix picture slides
    print(f"\n[2] Adding {len(APPENDIX_SLIDES)} instruction-screen appendix slides")
    for png, title in APPENDIX_SLIDES:
        slide = add_blank_no_title_slide(prs)
        build_appendix_slide(slide, png, title)
        print(f"  + {title}")

    n_after = len(prs.slides)
    print(f"\n  {n_before} -> {n_after} slides")

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
