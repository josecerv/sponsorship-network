"""
patch_2026_04_08_polish_v3.py — third polish pass.

Addresses Jose's 2026-04-08 afternoon review:
  1. Slides 21 & 22 (RQ3 strong/weak): drop the strength badge, widen
     the title, rewrite titles to be descriptive and parallel to 19/20.
  2. Slide 15 (Stage 3 randomization): expand the italic captions on
     each variant card to explain what Strong/Weak actually mean (and
     make gender + outcome captions parallel in info density).
  3. Slide 19 (RQ1): swap in the transposed rq3_initial_trust.png
     (strength on x-axis, gender as fill, overall-by-strength
     annotations). This breaks parallelism with RQ2/RQ3 on purpose —
     the RQ1 story is "no gender effect" and that reads best when the
     Male/Female bars sit side-by-side within each strength group.
  4. Timeline interweaving: rebuild slide 9 with Stage 1 highlighted,
     then insert two new navigation slides (Stage 2 highlighted before
     slide 11; Stage 3 highlighted before slide 13). Reorder via
     sldIdLst. Deck: 26 -> 28 slides.

Hooks: calls clear_non_placeholder_shapes on slide 9 and 21/22 badges.
After running, run fix_fld_guids.py + fix_orphaned_timing.py.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
RQ1_FIG = PROJECT / "pilots" / "output" / "talk_figures" / "rq3_initial_trust.png"


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------

NAVY         = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP    = RGBColor(0x00, 0x14, 0x3D)
NAVY_TINT    = RGBColor(0xE5, 0xEA, 0xF3)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT   = RGBColor(0x1F, 0x1F, 0x1F)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY     = RGBColor(0x5B, 0x65, 0x7A)
LIGHT_GRAY   = RGBColor(0xB0, 0xB8, 0xC6)
MUTED_GRAY   = RGBColor(0xA0, 0xA8, 0xB8)
TITLE_GRAY   = RGBColor(0x4B, 0x55, 0x63)
LIGHT_RULE   = RGBColor(0x9A, 0xA7, 0xBD)
CARD_OUTLINE = RGBColor(0xC7, 0xCF, 0xDD)
CARD_MUTED   = RGBColor(0xDF, 0xE4, 0xEC)

FONT = "Calibri"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size,
                bold=False, italic=False, color=BLACK_TEXT,
                align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color)
    return box


def add_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
             rounded=False, corner=0.5):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, in_(x), in_(y), in_(w), in_(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    if rounded:
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


def clear_non_placeholder_shapes(slide):
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        spTree.remove(shape._element)
        removed += 1
    return removed


def find_picture(slide):
    from pptx.shapes.picture import Picture
    pics = [s for s in slide.shapes if isinstance(s, Picture)]
    if not pics:
        return None
    pics.sort(key=lambda s: s.width * s.height, reverse=True)
    return pics[0]


def swap_picture(slide, new_png):
    pic = find_picture(slide)
    if pic is None:
        return False
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    slide.shapes._spTree.remove(pic._element)
    slide.shapes.add_picture(str(new_png),
                             left=left, top=top, width=width, height=height)
    return True


# ---------------------------------------------------------------------------
# Phase 1: Slides 21 & 22 — descriptive titles, drop the badge
# ---------------------------------------------------------------------------

RQ3_NEW_TITLES = {
    "RQ3 — Strong endorsers":
        "RQ3  —  Strong endorsements: how did one outcome shift wagers?",
    "RQ3 — Weak endorsers":
        "RQ3  —  Weak endorsements: how did one outcome shift wagers?",
}

# The strength badge is composed of five shapes anchored to the top-middle.
# We drop all of them (background card, header text, bar track, bar fill,
# Strong/Weak label, strong/unsure caption) to free up the title space.
BADGE_SHAPE_MARKERS = {
    # Text markers
    "ENDORSER CONFIDENCE",
    "Strong", "Weak",
    "very confident", "unsure",
    "hedged / unsure",  # legacy — safety net
}


def _shape_in_badge_region(shape):
    """Return True if the shape's top-left falls inside the former badge
    footprint (roughly x in [5.4, 8.0], y in [0.25, 1.80])."""
    try:
        x = shape.left / 914400
        y = shape.top / 914400
    except (AttributeError, TypeError):
        return False
    return (5.40 <= x <= 8.00) and (0.25 <= y <= 1.80)


def patch_rq3_slide(slide, slide_label):
    """Drop the strength badge, widen + retitle the RQ3 slide."""
    spTree = slide.shapes._spTree
    dropped = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        # Match by text OR by badge-region geometry
        matched = False
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt in BADGE_SHAPE_MARKERS:
                matched = True
        if not matched and _shape_in_badge_region(shape):
            # Skip the stats corner at x=9.38
            try:
                if shape.left / 914400 >= 9.00:
                    matched = False
                else:
                    # Badge background card + bar track + bar fill are non-text
                    if not shape.has_text_frame or shape.text_frame.text.strip() == "":
                        matched = True
            except Exception:
                pass
        if matched:
            spTree.remove(shape._element)
            dropped += 1

    # Retitle + widen
    title_done = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt in RQ3_NEW_TITLES:
            new = RQ3_NEW_TITLES[txt]
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = new
            shape.width = Inches(8.60)
            shape.height = Inches(1.40)
            title_done = True
            print(f"  {slide_label}: dropped {dropped} badge shape(s); title -> {new!r}")
            break
    if not title_done:
        print(f"  {slide_label}: WARN could not find RQ3 title")


# ---------------------------------------------------------------------------
# Phase 2: Slide 15 — expanded variant-card captions
# ---------------------------------------------------------------------------

# Find the existing italic caption text under each card header and replace
# with a fuller explanation. This relies on the v2 patcher already having
# laid out the cards — so these exact strings must still be on the slide.
SLIDE15_CAPTION_REWRITES = {
    "(name + colored avatar)":
        "The sponsor's first name + colored avatar signal their gender.",
    "(how confident the sponsor was)":
        "Sponsors used a 0–100 slider to choose a candidate — pushed near one end = Strong (very sure), left near the middle = Weak (barely committed).",
    "(was the endorsement right?)":
        "An endorsement was right if the endorsed candidate actually outscored the other on logical-reasoning.",
}


def patch_slide15_captions(slide):
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt in SLIDE15_CAPTION_REWRITES:
            new = SLIDE15_CAPTION_REWRITES[txt]
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = new
                # Shrink font on the long one so it still fits the card
                run = para.runs[0]
                run.font.size = Pt(10.5) if len(new) > 60 else Pt(11)
                run.font.italic = True
            # Taller frame to let it wrap
            shape.height = Inches(1.00)
            n += 1
    print(f"  slide 15: {n} caption(s) expanded")


# ---------------------------------------------------------------------------
# Phase 3: Slide 19 — swap RQ1 figure with the transposed version
# ---------------------------------------------------------------------------

def patch_slide19_rq1_figure(slide):
    if not RQ1_FIG.exists():
        print(f"  slide 19: WARN {RQ1_FIG} missing — skipping")
        return
    ok = swap_picture(slide, RQ1_FIG)
    if ok:
        print(f"  slide 19: swapped RQ1 figure (transposed: strength x-axis)")
    else:
        print(f"  slide 19: WARN no picture found")


# ---------------------------------------------------------------------------
# Phase 4: Timeline interweaving
# ---------------------------------------------------------------------------

# Shared stage-card copy (matches the original slide 9 content)
STAGE_CARDS = [
    dict(
        number="Stage 1",
        heading="Protégé performance",
        body=(
            "Real participants take three "
            "performance tasks (Barron et al., 2025). "
            "We use the posted scores to "
            "build candidate profiles."
        ),
    ),
    dict(
        number="Stage 2",
        heading="Sponsor endorsements",
        body=(
            "Real sponsors view candidate "
            "pairs and rate which one will "
            "perform better, and how confident they are."
        ),
    ),
    dict(
        number="Stage 3",
        heading="Audience evaluation",
        body=(
            "Evaluators bet on a sponsor's "
            "endorsement, see the outcome, "
            "then bet on a second endorsement "
            "from the same sponsor."
        ),
    ),
]


def _build_stage_card(slide, x, y, w, h, card, is_active):
    """Draw one stage card, highlighted iff is_active."""
    if is_active:
        bg_fill = NAVY_TINT
        border = NAVY
        border_w = 3.5
        num_color = NAVY_DEEP
        heading_color = NAVY_DEEP
        body_color = DARK_GRAY
    else:
        bg_fill = WHITE
        border = CARD_MUTED
        border_w = 1.25
        num_color = MUTED_GRAY
        heading_color = MID_GRAY
        body_color = MUTED_GRAY

    add_rect(slide, x, y, w, h,
             fill=bg_fill, outline=border, outline_width=border_w,
             rounded=True, corner=0.05)

    # Stage number (big)
    add_textbox(slide, x + 0.25, y + 0.30, w - 0.50, 0.80,
                card["number"],
                size=32, bold=True, color=num_color)

    # Heading
    add_textbox(slide, x + 0.25, y + 1.15, w - 0.50, 0.55,
                card["heading"],
                size=22, bold=True, color=heading_color)

    # Body
    body = slide.shapes.add_textbox(
        in_(x + 0.25), in_(y + 1.85), in_(w - 0.50), in_(h - 2.05)
    )
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = card["body"]
    for para in tf.paragraphs:
        para.line_spacing = 1.15
        for run in para.runs:
            style_run(run, 15, color=body_color)


def build_timeline_slide(slide, active_idx):
    """Clear non-placeholder shapes and rebuild the slide as a Stage
    navigation timeline with `active_idx` (0=Stage 1, 1=Stage 2, 2=Stage 3)
    highlighted."""
    clear_non_placeholder_shapes(slide)

    # Title
    add_textbox(slide, 0.55, 0.35, 12.23, 0.70,
                "Three-stage lab experiment",
                size=32, bold=True, color=NAVY_DEEP)

    # 3 stage cards
    card_y = 1.55
    card_w = 3.95
    card_h = 4.85
    gap = 0.32
    total_w = 3 * card_w + 2 * gap
    left = (13.33 - total_w) / 2
    xs = [left, left + card_w + gap, left + 2 * (card_w + gap)]

    for i, (x, card) in enumerate(zip(xs, STAGE_CARDS)):
        _build_stage_card(slide, x, card_y, card_w, card_h,
                          card, is_active=(i == active_idx))


def add_blank_slide_no_title(prs):
    """Add a new slide using the 'No Title' layout, keeping the slide-
    number placeholder and dropping any content/footer placeholders."""
    layout = None
    for lyt in prs.slide_layouts:
        if lyt.name == "No Title":
            layout = lyt
            break
    if layout is None:
        layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(layout)
    spTree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        try:
            idx = ph.placeholder_format.idx
        except Exception:
            continue
        if idx == 11:
            continue
        spTree.remove(ph._element)
    return slide


def move_slide_to(prs, slide, target_index):
    sldIdLst = prs.slides._sldIdLst
    sid = slide.slide_id
    match = None
    for sldId in sldIdLst:
        if int(sldId.get("id")) == sid:
            match = sldId
            break
    if match is None:
        raise ValueError(f"sldId not found for slide_id={sid}")
    sldIdLst.remove(match)
    children = list(sldIdLst)
    if target_index >= len(children):
        sldIdLst.append(match)
    else:
        sldIdLst.insert(target_index, match)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides before patch")

    s_9  = slides[8]    # Three-stage timeline
    s_15 = slides[14]   # Stage 3 randomization
    s_19 = slides[18]   # RQ1
    s_21 = slides[20]   # RQ3 strong
    s_22 = slides[21]   # RQ3 weak

    print("\n[Phase 1] Slides 21 & 22: descriptive titles, drop badge")
    patch_rq3_slide(s_21, "slide 21")
    patch_rq3_slide(s_22, "slide 22")

    print("\n[Phase 2] Slide 15: expanded variant captions")
    patch_slide15_captions(s_15)

    print("\n[Phase 3] Slide 19: swap RQ1 figure (transposed)")
    patch_slide19_rq1_figure(s_19)

    print("\n[Phase 4] Timeline interweaving")
    # 4a: rebuild slide 9 with Stage 1 highlighted
    print("  [4a] Rebuild slide 9 (Stage 1 highlighted)")
    build_timeline_slide(s_9, active_idx=0)

    # 4b: add Stage 2 highlight slide, move to position before current slide 11
    print("  [4b] Add Stage 2 highlight slide")
    s_stage2 = add_blank_slide_no_title(prs)
    build_timeline_slide(s_stage2, active_idx=1)

    # 4c: add Stage 3 highlight slide
    print("  [4c] Add Stage 3 highlight slide")
    s_stage3 = add_blank_slide_no_title(prs)
    build_timeline_slide(s_stage3, active_idx=2)

    # 4d: reorder via sldIdLst
    # Current sldIdLst order BEFORE move:
    #   0..25  (original 26 slides)
    #   26     = s_stage2 (appended)
    #   27     = s_stage3 (appended)
    # We want s_stage2 to sit at position 10 (= just before original slide 11),
    # and s_stage3 at position 12 (= just before original slide 13). After
    # inserting s_stage2 at 10, the original slide 11 shifts to position 11;
    # everything downstream shifts by one. So s_stage3 should go to position
    # 13 (not 12) to land before the original slide 13. Let me recount:
    #
    #   After inserting s_stage2 at index 10:
    #     0 .. 9   = original 0..9   (covers cover..slide 10)
    #     10       = s_stage2  (new)
    #     11 .. 25 = original 10..24 (original slide 11..25)
    #     26       = s_stage3  (appended)
    #
    #   Now we want s_stage3 right before the original slide 13, which is
    #   now at index 13 (was 12, shifted by +1). So move s_stage3 to 13.
    move_slide_to(prs, s_stage2, 10)
    move_slide_to(prs, s_stage3, 13)

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print(f"Done. {len(prs.slides)} slides after patch.")


if __name__ == "__main__":
    main()
