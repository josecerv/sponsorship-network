"""
patch_2026_04_08_polish_v3b.py — Stage 2 sample size / decision count.

Small follow-up to polish_v3:
  - Slide 12 (Stage 2 content, "Real sponsors made real endorsements"):
    bullet "Each sponsor sees one candidate pair from Stage 1" was
    inaccurate. Actual Stage 2 sponsors each saw 10 candidate pairs and
    made 10 endorsements. Rewrite the bullet. Source:
    pilots/old/output/stage2_main_data.csv has columns t1_* through
    t10_* per sponsor.
  - Slide 13 (Stage 2 incentive): add a small "n = 200 sponsors · 10
    endorsements each" annotation right below the sub-title so the
    audience knows the scale before they read the payout table. Keeps
    the existing layout.

No shape clearing. Pure surgical edits. fix_fld_guids.py not needed.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"

NAVY_DEEP = RGBColor(0x00, 0x14, 0x3D)
MID_GRAY  = RGBColor(0x5B, 0x65, 0x7A)
DARK_GRAY = RGBColor(0x1F, 0x29, 0x37)


def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=DARK_GRAY):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Slide 12 — fix the "one pair" bullet
# ---------------------------------------------------------------------------

def patch_slide12_bullets(slide):
    """Find the bullet-list textbox and replace the inaccurate 'Each
    sponsor sees one candidate pair from Stage 1' bullet with the
    accurate 10-pair version. Preserves all other bullets and formatting.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full_text = shape.text_frame.text
        if "Each sponsor sees one candidate pair" not in full_text:
            continue
        # Walk paragraphs, find the matching one, and edit its runs
        for para in shape.text_frame.paragraphs:
            p_text = "".join(r.text for r in para.runs)
            if "Each sponsor sees one candidate pair" in p_text:
                # Keep the bullet character ("•  ") in run[0] if present,
                # otherwise rewrite run[0] whole.
                runs = para.runs
                if not runs:
                    continue
                if runs[0].text.strip() == "•":
                    # Paragraph starts with "• " then content
                    pass
                # Concentrate the new text in the first run, blank the rest
                new_text = "Each sponsor sees 10 candidate pairs and makes 10 endorsements"
                if runs[0].text.startswith("•"):
                    # Preserve the bullet prefix
                    prefix_match = ""
                    for ch in runs[0].text:
                        if ch in "• ":
                            prefix_match += ch
                        else:
                            break
                    runs[0].text = prefix_match + new_text
                else:
                    runs[0].text = new_text
                for r in runs[1:]:
                    r.text = ""
                print("  slide 12: bullet updated (one pair -> 10 pairs)")
                return
    print("  slide 12: WARN no matching bullet found")


# ---------------------------------------------------------------------------
# Slide 13 — add sample size annotation
# ---------------------------------------------------------------------------

def patch_slide13_sample_size(slide):
    """Add a small 'n = 200 sponsors · 10 endorsements each' annotation
    between the sub-line and the outcome cards. Also update the bottom
    footer to reference the decision count for context."""
    # 1. Add sample size annotation at y=2.20 (between sub at 1.78 and
    #    cards at 2.30). Height 0.08 is tight — use a thin text box.
    ann = slide.shapes.add_textbox(in_(0.55), in_(2.20), in_(12.20), in_(0.28))
    tf = ann.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.text = "n = 200 sponsors  ·  10 endorsements each"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 12, italic=True, bold=True, color=NAVY_DEEP)
    print("  slide 13: added sample size annotation")

    # Need to nudge the existing outcome cards down slightly so they don't
    # collide with the new annotation. Cards are at y=2.30 — move to 2.50.
    # Actually, the cards are at y=2.30 (card_y from polish v2) and the
    # annotation ends at 2.48 — they collide. Shift cards + downstream
    # elements down by 0.22".
    # Identify cards by exact y=2.30 position.
    dy = Inches(0.22)
    # Everything at y in [2.30, 2.35) on original layout gets shifted
    # Cards: (_,2.30,_,0.92) — shift
    # We want to keep the slider + reveals + table + footer UNCHANGED
    # so only shift the cards. Cards are the 2 rectangle shapes at y=2.30.
    shifted = 0
    for shape in slide.shapes:
        try:
            y_in = shape.top / 914400
        except (AttributeError, TypeError):
            continue
        # Card rectangles are at y=2.30, h=0.92
        if 2.28 <= y_in <= 2.32 and shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt.startswith("If right") or txt.startswith("If wrong"):
                shape.top = shape.top + dy
                shifted += 1
    print(f"  slide 13: shifted {shifted} outcome card(s) down by 0.22 in")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides")

    s_12 = slides[11]  # Stage 2 content
    s_13 = slides[12]  # Stage 2 incentive

    print("\n[Phase 1] Slide 12 bullet fix")
    patch_slide12_bullets(s_12)

    print("\n[Phase 2] Slide 13 sample size annotation")
    patch_slide13_sample_size(s_13)

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
