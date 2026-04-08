"""
patch_2026_04_08_b.py — follow-up to patch_2026_04_08.py.

Two fixes:
  1. Insert a NEW slide right after hook 2a (Robert -> Nick) showing
     Susan Keller -> Nick Calder in the same centered single-letter layout.
     This is the pre-outcome gender flip — it directly teases RQ1
     ("are male and female sponsors trusted equally BEFORE any outcome?").
     The deck grows from 25 -> 26 slides.

  2. Fix missing slide-number placeholders on slides 3, 4, 5 (i.e. the
     three hook beats that sit between the cover and "What is sponsorship?").
     python-pptx's add_slide() does NOT copy the slide-number placeholder
     from the layout into the new slide — it leaves placeholder inheritance
     to the renderer. That's fine in theory, but in practice our deck's
     existing slides each carry their OWN cloned slide-number placeholder
     (because the original template was a custom Wharton deck), so the
     new slides end up with nothing and render without a number.

     Fix: clone the slide-number placeholder from slide 2 (which has a
     working one left over from the user-inserted template slide 10) onto
     every hook slide that's missing one. Uses the same deepcopy + id
     refresh pattern patch_deck.clone_slide_number_placeholder uses.

Run AFTER patch_2026_04_08.py has been applied. Back up the deck first.
"""

from copy import deepcopy
from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

# Reuse all the letter-card + hook helpers from the prior patcher so the
# visual language stays identical.
from patch_2026_04_08 import (
    build_letter_card,
    add_blank_hook_slide,
    move_slide_to,
)


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"


# ---------------------------------------------------------------------------
# New single-letter builder for Susan Keller -> Nick Calder
# ---------------------------------------------------------------------------

def build_hook_2a_susan(slide):
    """Mirror of build_hook_2a but with Susan Keller as the signer,
    plus a subtle pink highlight behind the signer name so the flip
    from Robert -> Susan is visible at a glance."""
    build_letter_card(
        slide,
        x_in=3.92, y_in=0.90, w_in=5.50, h_in=6.20,
        student="Nick Calder",
        signer="Susan Keller",
        body_font_size=14,
        signer_highlight=True,
    )


# ---------------------------------------------------------------------------
# Slide-number placeholder cloning (borrowed pattern from patch_deck.py)
# ---------------------------------------------------------------------------

def find_slide_number_placeholder(slide):
    """Return the idx=11 SLIDE_NUMBER placeholder element, or None."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == 11:
                return ph._element
        except (AttributeError, ValueError):
            continue
    return None


def clone_slide_number_placeholder(slide, template_element, next_id):
    """Deep-copy `template_element` and append it to `slide`'s spTree.

    - Rewrites the cNvPr id so it's unique within the destination slide.
    - Generates a fresh a16:creationId GUID so renderers that de-dupe
      by creationId don't choke.
    """
    new_el = deepcopy(template_element)

    cNvPr = new_el.find(".//" + qn("p:cNvPr"))
    if cNvPr is not None:
        cNvPr.set("id", str(next_id))

    a16_ns = "http://schemas.microsoft.com/office/drawing/2014/main"
    for ext in new_el.findall(".//{%s}creationId" % a16_ns):
        ext.set("id", "{%s}" % str(uuid4()).upper())

    slide.shapes._spTree.append(new_el)
    return new_el


def fix_missing_slide_numbers(prs):
    """Find a working template slide-number placeholder somewhere in the
    deck (slide 2 has one), then walk every slide and clone it onto any
    slide that doesn't already have an idx=11 placeholder. Skips slide 1
    (the cover is intentionally unnumbered)."""
    template_el = None
    for i, slide in enumerate(prs.slides):
        el = find_slide_number_placeholder(slide)
        if el is not None:
            template_el = el
            print(f"  template slide-number placeholder found on slide {i+1}")
            break

    if template_el is None:
        print("  WARN: no template slide-number placeholder found — skipping")
        return

    # Give cloned placeholders fresh IDs starting well above anything
    # python-pptx is likely to auto-allocate.
    next_id = 500
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            # Cover slide — intentionally unnumbered
            continue
        if find_slide_number_placeholder(slide) is None:
            clone_slide_number_placeholder(slide, template_el, next_id)
            next_id += 1
            print(f"  slide {i:2d}: cloned slide-number placeholder")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    print(f"  {len(prs.slides)} slides before patch")

    # --- Phase 1: add new Susan-Keller single-letter slide --------------------
    print("\n[Phase 1] Adding new hook slide: Susan Keller -> Nick Calder")
    s_new = add_blank_hook_slide(prs)
    build_hook_2a_susan(s_new)
    print("  built single centered letter card (Susan variant)")

    # Move the new slide to index 2 (0-based), i.e. slide 3 in 1-based display,
    # which sits right after cover+hook-2a and before the two-letter slides.
    move_slide_to(prs, s_new, 2)
    print("  moved new slide -> index 2 (display slide 3)")

    # --- Phase 2: clone slide-number placeholders onto every hook slide -----
    print("\n[Phase 2] Fixing missing slide-number placeholders")
    fix_missing_slide_numbers(prs)

    # --- Save -----------------------------------------------------------------
    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print(f"Done.  Deck now has {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
