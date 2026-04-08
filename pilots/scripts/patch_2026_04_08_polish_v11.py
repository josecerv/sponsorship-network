"""
patch_2026_04_08_polish_v11.py — slide 27 perception-vs-fact tweak.

Refines v10. Bullets 2 and 3 previously read as factual claims about
the actual endorsement distributions ("Women's endorsements cluster
high + narrow", "Men's endorsements spread wide"). Jose has separate
data showing the actual endorsement strengths don't differ by gender,
so the running story has to be explicit that this is a perception
about endorsements, not a fact about them.

v11 adds "are assumed" to the bullet 2 + 3 headers and reframes the
bullet 3 body so "carries information" is explicitly the evaluator's
read, not a property of men's calls.

Bullets 1, 4, the subtitle, the figure, and the caption are unchanged.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR = PROJECT / "pilots" / "output" / "talk_figures"
PRIOR_FIG = FIG_DIR / "prior_distributions.png"


NAVY       = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP  = RGBColor(0x00, 0x14, 0x3D)
DARK_GRAY  = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY   = RGBColor(0x5B, 0x65, 0x7A)
BLACK_TEXT = RGBColor(0x11, 0x18, 0x27)

FONT = "Arial"


SUBTITLE = "The prior leaks in once there's an outcome to read."

BULLETS = [
    (
        "Without a track record, people stay reserved.",
        "RQ1 is null (b = -1.06, p = .81). With nothing but a "
        "strength label, there's no basis to discriminate yet. "
        "Unlike most discrimination findings, the gap leaks in "
        "only at the outcome.",
    ),
    (
        "Women's endorsements are assumed high + narrow.",
        "Once the outcome arrives, \"she would have said it "
        "either way\" makes it barely informative, so trust "
        "barely moves (consistent with the muted effects in the "
        "stronger confidence cells).",
    ),
    (
        "Men's endorsements are assumed to spread wide.",
        "Evaluators read men's confident calls as carrying "
        "information that can be both praised or penalized, so "
        "trust moves accordingly.",
    ),
    (
        "Next direction: telegraph their standards.",
        "Signaling \"I only recommend the top 10%\" breaks the "
        "assumed clustering, so the outcome can be re-read as "
        "informative.",
    ),
]


# Layout — same as v7b/v8/v9/v10
LEFT_X     = 0.55
LEFT_Y     = 1.55
LEFT_W     = 7.20
BULLET_H   = 1.20
BULLET_GAP = 0.10
HEADER_PT  = 19
BODY_PT    = 16

FIG_X = 7.85
FIG_Y = 1.55
FIG_W = 5.40
FIG_H = 3.32  # 5.40 / (7.8/4.8 source aspect)


def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
                color=BLACK_TEXT, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color, font=font)
    return box


def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)


def main():
    prs = Presentation(str(DECK))
    slide = prs.slides[26]  # slide 27

    dropped = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        remove_shape(slide, shape)
        dropped += 1
    print(f"  slide 27: dropped {dropped} shape(s) to rebuild")

    # Title
    add_textbox(
        slide, 0.55, 0.30, 12.23, 0.70,
        "Running story",
        size=32, bold=True, color=NAVY,
    )

    # Subtitle
    add_textbox(
        slide, 0.55, 1.05, 12.23, 0.42,
        SUBTITLE,
        size=20, italic=True, color=MID_GRAY,
    )

    # 4 bullets, left column
    for i, (head, body) in enumerate(BULLETS):
        y = LEFT_Y + i * (BULLET_H + BULLET_GAP)

        mark = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            in_(LEFT_X), in_(y + 0.14),
            in_(0.16), in_(0.16),
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = NAVY
        mark.line.fill.background()
        mark.text_frame.text = ""

        box = slide.shapes.add_textbox(
            in_(LEFT_X + 0.28), in_(y),
            in_(LEFT_W - 0.28), in_(BULLET_H),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_top = Inches(0.0)

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = head
        style_run(r1, HEADER_PT, bold=True, color=NAVY_DEEP)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = body
        style_run(r2, BODY_PT, bold=False, color=DARK_GRAY)

    print(f"  slide 27: added {len(BULLETS)} bullet(s) at {HEADER_PT}/{BODY_PT}pt")

    if PRIOR_FIG.exists():
        slide.shapes.add_picture(
            str(PRIOR_FIG),
            left=in_(FIG_X), top=in_(FIG_Y),
            width=in_(FIG_W), height=in_(FIG_H),
        )
        add_textbox(
            slide, FIG_X, FIG_Y + FIG_H + 0.02, FIG_W, 0.30,
            "Hypothesized prior distributions",
            size=12, italic=True, color=MID_GRAY,
            align=PP_ALIGN.CENTER,
        )
        print("  slide 27: prior_distributions.png + caption restored")

    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
