"""
patch_2026_04_08_polish_v5.py — UChicago deck polish pass 5.

Six-slide polish pass per Jose's 2026-04-08 evening feedback:

 1. Slide 2 — Goals for this talk. Polish Jose's manual text; add the
    Wharton logo at top-right so slide 2 shares visual DNA with the
    cover slide.

 2. Slides 5 & 6 — YEAR 1 → YEAR 2 hook beats. Drop the clunky
    "ONE YEAR LATER" caption + arrow + SUCCEEDED pill. Replace with a
    case-file-style treatment:
      a. On the Year 1 (Nick) letter: rotated "GOOD STUDENT ✓" stamp
         overlay and gray-out the letter body/signature text.
      b. On the Year 2 (Jose) letter: large circled "?" overlay marking
         the open question.

 3. Slide 17 — Stage 3 setup. Enlarge body fonts + widen spacing so the
    slide fills more of the content area.

 4. Slide 22 — RQ1 result. Replace the two fuzzy side-by-side pictures
    (rq1_overall.png + rq3_initial_trust.png) with the single
    rq1_combined.png 2-panel composite that includes reference bands
    showing what the sponsor said (strong 77-88, weak 12-23).

 5. Slides 24 & 25 — RQ3 strong/weak. Left-align the title to match
    the rest of the deck, replace the top-left dual-band confidence
    badge with a compact single-band badge (strong only on s24, weak
    only on s25) sitting between the title and the stats corner, and
    shrink the stats corner itself.

Surgical edits only. No clear_non_placeholder_shapes calls, so
fix_fld_guids.py / fix_orphaned_timing.py are NOT required after this
run.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Paths & config
# ---------------------------------------------------------------------------

PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR = PROJECT / "pilots" / "output" / "talk_figures"
RQ1_COMBINED = FIG_DIR / "rq1_combined.png"
WHARTON_LOGO = PROJECT / "pilots" / "output" / "wharton_logo.png"

STRONG_LO, STRONG_HI = 77, 88
WEAK_LO, WEAK_HI = 12, 23


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------

NAVY         = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP    = RGBColor(0x00, 0x14, 0x3D)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT   = RGBColor(0x11, 0x18, 0x27)
DARK_GRAY    = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY     = RGBColor(0x5B, 0x65, 0x7A)
BODY_GRAY    = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY   = RGBColor(0x9C, 0xA3, 0xAF)
VERY_LIGHT   = RGBColor(0xC7, 0xCF, 0xDD)
CARD_OUTLINE = RGBColor(0xC7, 0xCF, 0xDD)
BAR_TRACK    = RGBColor(0xE5, 0xEA, 0xF3)
WRONG_RED    = RGBColor(0xB1, 0x21, 0x21)
STAMP_RED    = RGBColor(0xBA, 0x1F, 0x1F)
AMBER_LINE   = RGBColor(0xB4, 0x53, 0x09)
AMBER_SOFT   = RGBColor(0xFE, 0xF3, 0xC7)

FONT = "Arial"


# ---------------------------------------------------------------------------
# Small helpers
# ---------------------------------------------------------------------------

def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
                color=BLACK_TEXT, align=PP_ALIGN.LEFT, anchor=None, font=FONT):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color, font=font)
    return box


def add_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
             rounded=False, corner=0.5):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, in_(x), in_(y), in_(w), in_(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    if rounded:
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)


# ---------------------------------------------------------------------------
# Phase 1 — Slide 2 (Goals) polish + Wharton logo
# ---------------------------------------------------------------------------

GOALS_TITLE = "Goals for this talk"

GOALS_LEAD = "Two things we'd love feedback on today:"

GOALS_BULLETS = [
    (
        "Experimental design",
        "anything to flag before we scale up and preregister?",
    ),
    (
        "Interpretation of the Stage 3 pilot",
        "we have a running theory, but alternatives are very welcome.",
    ),
]


def patch_slide2_goals(slide):
    """Drop the existing content placeholder, add a polished title/body,
    then place the Wharton wordmark at top-right to echo slide 1."""
    # Drop every non-slidenum placeholder so we start from a clean slide.
    dropped = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            try:
                idx = shape.placeholder_format.idx
            except Exception:
                idx = None
            # Keep the slide number placeholder (idx 11 in this deck).
            if idx == 11:
                continue
            remove_shape(slide, shape)
            dropped += 1
    print(f"  slide 2: dropped {dropped} placeholder(s)")

    # Title at left
    add_textbox(
        slide, 0.55, 0.35, 8.80, 0.90,
        GOALS_TITLE, size=32, bold=True, color=NAVY,
    )

    # Wharton logo at top-right (mirrors slide 1, which has it at top-left)
    if WHARTON_LOGO.exists():
        logo_w = 3.10
        logo_h = 0.88
        logo_x = 13.33 - logo_w - 0.30  # 9.93
        logo_y = 0.45
        slide.shapes.add_picture(
            str(WHARTON_LOGO),
            left=in_(logo_x), top=in_(logo_y),
            width=in_(logo_w), height=in_(logo_h),
        )
        print("  slide 2: Wharton logo placed top-right")
    else:
        print(f"  slide 2: WARNING no logo at {WHARTON_LOGO}")

    # Thin navy accent rule under the title to echo the cover slide
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, in_(0.55), in_(1.50), in_(1.60), in_(0.06)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = NAVY
    rule.line.fill.background()

    # Lead-in line
    add_textbox(
        slide, 0.55, 1.85, 12.23, 0.60,
        GOALS_LEAD, size=22, italic=True, color=MID_GRAY,
    )

    # Two bullet cards — rounded-rect accent + bold header + body
    card_y = 2.80
    card_h = 1.65
    card_gap = 0.35
    for i, (head, body) in enumerate(GOALS_BULLETS):
        y = card_y + i * (card_h + card_gap)
        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            in_(0.55), in_(y + 0.15), in_(0.18), in_(card_h - 0.30),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        try:
            bar.adjustments[0] = 0.5
        except Exception:
            pass
        bar.text_frame.text = ""

        # Header + body in a single textbox so PowerPoint keeps them grouped
        card = slide.shapes.add_textbox(
            in_(0.95), in_(y + 0.05), in_(12.00), in_(card_h),
        )
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_top = Inches(0.05)

        # Header paragraph
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = head
        style_run(r1, 24, bold=True, color=NAVY_DEEP)

        # Body paragraph
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = body
        style_run(r2, 21, bold=False, color=DARK_GRAY)

    print("  slide 2: polished goals body in place")


# ---------------------------------------------------------------------------
# Phase 2 — Slides 5 & 6 stamp/question-mark overlay
# ---------------------------------------------------------------------------

# Nick's letter card (left half) ~ x=0.55, w=4.70, y=0.55, h=5.55
# Jose's letter card (right half) ~ x=8.08, w=4.70, y=0.55, h=5.55
LEFT_CARD_X_MIN,  LEFT_CARD_X_MAX  = 0.30, 5.50
RIGHT_CARD_X_MIN, RIGHT_CARD_X_MAX = 8.00, 13.00


def _shape_xy(shape):
    try:
        return shape.left / 914400, shape.top / 914400
    except Exception:
        return None, None


def _shape_wh(shape):
    try:
        return shape.width / 914400, shape.height / 914400
    except Exception:
        return None, None


def _gray_text_shape(shape, gray=LIGHT_GRAY):
    """Recolor every run in a text-bearing shape to the passed gray."""
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = gray


def patch_hook_year_overlay(slide, label):
    """Remove ONE YEAR LATER + arrow + SUCCEEDED pill, stamp the Year 1
    letter as resolved, and mark the Year 2 letter with a question
    circle."""
    # ---- 1. Drop the old interstitial elements ----
    to_drop = []
    for shape in list(slide.shapes):
        if not shape.has_text_frame and shape.shape_type != 1:
            continue
        x, y = _shape_xy(shape)
        if x is None:
            continue
        # Interstitial zone: between the two cards, y < 4.2
        if not (5.2 <= x <= 8.00):
            continue
        if y > 4.20:
            continue
        txt = shape.text_frame.text.strip() if shape.has_text_frame else ""
        if txt in {"ONE YEAR LATER", "SUCCEEDED"}:
            to_drop.append(shape)
            continue
        # Non-text interstitial shape: the arrow
        if txt == "" and shape.shape_type == 1:
            w, h = _shape_wh(shape)
            if 1.50 <= w <= 2.50 and 0.5 <= h <= 1.10:
                to_drop.append(shape)
                continue

    for shape in to_drop:
        remove_shape(slide, shape)
    print(f"  {label}: dropped {len(to_drop)} interstitial shape(s)")

    # ---- 2. Gray out the Year 1 (Nick) letter body + signature ----
    grayed = 0
    for shape in list(slide.shapes):
        x, y = _shape_xy(shape)
        if x is None:
            continue
        if not (LEFT_CARD_X_MIN <= x <= LEFT_CARD_X_MAX):
            continue
        if not (0.55 <= y <= 6.20):
            continue
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        # Skip the YEAR 1 badge (it floats above the card at y=0.10)
        if txt == "YEAR 1":
            continue
        # Skip empty / background-only shapes
        if not txt:
            continue
        _gray_text_shape(shape, gray=LIGHT_GRAY)
        grayed += 1
    print(f"  {label}: grayed {grayed} Year 1 text element(s)")

    # ---- 3. Stamp the Year 1 letter ----
    stamp_w = 3.40
    stamp_h = 0.95
    stamp_cx = 2.90
    stamp_cy = 3.20
    stamp_x = stamp_cx - stamp_w / 2
    stamp_y = stamp_cy - stamp_h / 2

    stamp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        in_(stamp_x), in_(stamp_y), in_(stamp_w), in_(stamp_h),
    )
    stamp.fill.background()
    stamp.line.color.rgb = STAMP_RED
    stamp.line.width = Pt(4.0)
    stamp.rotation = -11.0

    tf = stamp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = "GOOD STUDENT  \u2713"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.runs[0], 28, bold=True, color=STAMP_RED, font="Arial")

    # Sub-label under the stamp ("ONE YEAR LATER")
    sub = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        in_(stamp_x), in_(stamp_y + stamp_h + 0.02),
        in_(stamp_w), in_(0.30),
    )
    sub.fill.background()
    sub.line.fill.background()
    sub.rotation = -11.0
    tfs = sub.text_frame
    tfs.margin_left = Inches(0.04)
    tfs.margin_right = Inches(0.04)
    tfs.margin_top = Inches(0.02)
    tfs.margin_bottom = Inches(0.02)
    tfs.word_wrap = False
    tfs.text = "one year later"
    ps = tfs.paragraphs[0]
    ps.alignment = PP_ALIGN.CENTER
    style_run(ps.runs[0], 13, bold=True, italic=True, color=STAMP_RED, font="Arial")

    print(f"  {label}: stamped Year 1 letter")

    # ---- 4. Question-mark overlay on Year 2 (Jose) letter ----
    # Oval: centered on Jose's card mid-body
    q_w = 2.20
    q_h = 2.20
    q_cx = 10.43
    q_cy = 3.00
    q_x = q_cx - q_w / 2
    q_y = q_cy - q_h / 2

    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        in_(q_x), in_(q_y), in_(q_w), in_(q_h),
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = AMBER_SOFT
    oval.line.color.rgb = AMBER_LINE
    oval.line.width = Pt(3.5)
    tfq = oval.text_frame
    tfq.margin_left = Inches(0.02)
    tfq.margin_right = Inches(0.02)
    tfq.margin_top = Inches(0.02)
    tfq.margin_bottom = Inches(0.02)
    tfq.vertical_anchor = MSO_ANCHOR.MIDDLE
    tfq.word_wrap = False
    tfq.text = "?"
    pq = tfq.paragraphs[0]
    pq.alignment = PP_ALIGN.CENTER
    style_run(pq.runs[0], 96, bold=True, color=AMBER_LINE, font="Arial")

    # Caption under the "?": "year 2 — new case"
    qcap = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        in_(q_cx - 1.40), in_(q_y + q_h + 0.05),
        in_(2.80), in_(0.32),
    )
    qcap.fill.background()
    qcap.line.fill.background()
    tfqc = qcap.text_frame
    tfqc.margin_left = Inches(0.02)
    tfqc.margin_right = Inches(0.02)
    tfqc.margin_top = Inches(0.02)
    tfqc.margin_bottom = Inches(0.02)
    tfqc.word_wrap = False
    tfqc.text = "how do you read this one?"
    pqc = tfqc.paragraphs[0]
    pqc.alignment = PP_ALIGN.CENTER
    style_run(pqc.runs[0], 13, bold=True, italic=True, color=AMBER_LINE, font="Arial")

    print(f"  {label}: ? overlay placed on Year 2 letter")


# ---------------------------------------------------------------------------
# Phase 3 — Slide 17 (Stage 3 setup) enlargement
# ---------------------------------------------------------------------------

# New y / size plan per shape signature (matched by text prefix or by
# (x, y) approximate position).

SLIDE17_TEXT_UPDATES = [
    # (match-prefix, new_x, new_y, new_w, new_h, new_size, bold_override)
    ("n = 403 Prolific",                 0.55, 1.15, 12.20, 0.45, 18, None),
    ("Each evaluator is randomly",       0.55, 1.68, 12.20, 0.50, 20, None),
    ("the sponsor's name",               0.95, 2.28, 11.80, 0.50, 20, None),
    ("that sponsor's endorsement",       0.95, 2.80, 11.80, 0.50, 20, None),
    ("how confident the sponsor",        0.95, 3.32, 11.80, 0.50, 20, None),
    ("a slider to wager part",           0.95, 3.84, 11.80, 0.50, 20, None),
    ("Each evaluator sees TWO",          0.55, 4.48, 12.20, 0.55, 20, None),
    ("trust change",                     0.55, 6.18, 12.20, 0.70, 28, True),
]


SLIDE17_FLOW_BOX_LABELS = {"First wager", "outcome revealed", "Second wager"}


def _resize_font_runs(shape, new_size, force_bold=None):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip() == "":
                continue
            r.font.size = Pt(new_size)
            if force_bold is not None:
                r.font.bold = force_bold


def _matches_prefix(text, prefix):
    return text.lstrip().startswith(prefix)


def patch_slide17_enlarge(slide):
    """Bump body/bullet/flow-box/formula fonts and reflow y positions
    so the slide fills more of the content area."""
    moved = 0

    # Update textual body + bullets + formula
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if not txt:
            continue
        for prefix, nx, ny, nw, nh, sz, bold_override in SLIDE17_TEXT_UPDATES:
            if _matches_prefix(txt, prefix):
                shape.left = in_(nx)
                shape.top = in_(ny)
                shape.width = in_(nw)
                shape.height = in_(nh)
                _resize_font_runs(shape, sz, force_bold=bold_override)
                moved += 1
                break

    # Flow boxes — relocate + enlarge
    # New flow-box row at y=5.05, wider boxes (2.50") at x=2.57, 5.42, 8.27
    flow_positions = {
        "First wager":      (2.57, 5.05, 2.50, 0.85),
        "outcome revealed": (5.42, 5.05, 2.50, 0.85),
        "Second wager":     (8.27, 5.05, 2.50, 0.85),
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt in flow_positions:
            nx, ny, nw, nh = flow_positions[txt]
            shape.left = in_(nx)
            shape.top = in_(ny)
            shape.width = in_(nw)
            shape.height = in_(nh)
            _resize_font_runs(shape, 20, force_bold=True)
            moved += 1

    # Flow arrows between boxes — reposition at the row midline
    # Original: y=4.62 w=0.55 h=0.20. Move to new row center.
    arrow_positions = [
        (4.96, 5.43, 0.55, 0.20),  # between box 1 and 2
        (7.81, 5.43, 0.55, 0.20),  # between box 2 and 3
    ]
    arrow_idx = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            continue
        if shape.shape_type != 1:  # auto_shape
            continue
        try:
            w = shape.width / 914400
            h = shape.height / 914400
        except Exception:
            continue
        # Match the ~0.55 x 0.20 arrows only
        if not (0.45 <= w <= 0.65 and 0.15 <= h <= 0.25):
            continue
        if arrow_idx >= len(arrow_positions):
            break
        nx, ny, nw, nh = arrow_positions[arrow_idx]
        shape.left = in_(nx)
        shape.top = in_(ny)
        shape.width = in_(nw)
        shape.height = in_(nh)
        arrow_idx += 1
        moved += 1

    print(f"  slide 17: relocated/enlarged {moved} shape(s)")


# ---------------------------------------------------------------------------
# Phase 4 — Slide 22 (RQ1) combined figure
# ---------------------------------------------------------------------------

def patch_slide22_rq1_combined(slide):
    """Drop the two side-by-side pictures + OVERALL/SPLIT captions and
    add the single rq1_combined.png composite."""
    from pptx.shapes.picture import Picture

    # 1. Remove captions (OVERALL and SAME DATA ... captions) and both pics
    to_drop = []
    for shape in list(slide.shapes):
        if isinstance(shape, Picture):
            to_drop.append(shape)
            continue
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == "OVERALL" or txt.startswith("SAME DATA"):
                to_drop.append(shape)
    for shape in to_drop:
        remove_shape(slide, shape)
    print(f"  slide 22: dropped {len(to_drop)} old figure element(s)")

    # 2. Add the combined picture. rq1_combined.png is 15x6 at 260 dpi ->
    #    aspect ratio 2.5. Slide content area x=0.10..13.23 (13.13 wide)
    #    matches the RQ2/RQ3 positioning.
    if not RQ1_COMBINED.exists():
        print(f"  slide 22: WARNING no figure at {RQ1_COMBINED}")
        return

    slide.shapes.add_picture(
        str(RQ1_COMBINED),
        left=in_(0.10), top=in_(1.80),
        width=in_(13.13), height=in_(5.25),
    )
    print("  slide 22: rq1_combined.png placed")


# ---------------------------------------------------------------------------
# Phase 5 — Slides 24 & 25: single-band badge + left title + compact stats
# ---------------------------------------------------------------------------

# Badge element match: existing badge lives inside [0.30 .. 3.00] x [0.20 .. 1.80]
# on both slides. We delete all shapes in that rect and rebuild as a compact
# single-band badge to the RIGHT of the title.

BADGE_DELETE_RECT = (0.30, 0.20, 3.00, 1.80)  # x_min, y_min, x_max, y_max


def _shape_in_rect(shape, x_min, y_min, x_max, y_max):
    x, y = _shape_xy(shape)
    if x is None:
        return False
    return x_min <= x <= x_max and y_min <= y <= y_max


def draw_single_band_badge(slide, x, y, w, h, band="strong"):
    """Compact endorser-confidence badge showing only ONE highlighted
    band (strong OR weak). Header + scale bar track + one filled band
    + endpoint labels + a single "Strong 77-88" / "Weak 12-23" label."""
    # Frame
    add_rect(slide, x, y, w, h,
             fill=WHITE, outline=CARD_OUTLINE, outline_width=1.0,
             rounded=True, corner=0.10)

    margin = 0.14
    inner_x = x + margin
    inner_w = w - 2 * margin

    # Header
    header_y = y + 0.06
    header_h = 0.22
    add_textbox(slide, inner_x, header_y, inner_w, header_h,
                "ENDORSER CONFIDENCE",
                size=10, bold=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    # Scale bar track
    bar_y = header_y + header_h + 0.10
    bar_h = 0.22
    track = add_rect(slide, inner_x, bar_y, inner_w, bar_h,
                     fill=BAR_TRACK, rounded=True, corner=0.5)
    track.line.color.rgb = VERY_LIGHT
    track.line.width = Pt(0.75)

    # Only one highlighted band
    if band == "strong":
        lo, hi = STRONG_LO, STRONG_HI
        fill_color = NAVY
        label_color = NAVY_DEEP
        label = "Strong"
    else:
        lo, hi = WEAK_LO, WEAK_HI
        fill_color = WRONG_RED
        label_color = WRONG_RED
        label = "Weak"

    band_x = inner_x + inner_w * (lo / 100.0)
    band_w = inner_w * ((hi - lo) / 100.0)
    add_rect(slide, band_x, bar_y, band_w, bar_h,
             fill=fill_color, rounded=True, corner=0.5)

    # Endpoint scale labels "0" and "100"
    endpoint_y = bar_y + bar_h + 0.02
    add_textbox(slide, inner_x - 0.10, endpoint_y, 0.35, 0.22,
                "0", size=9, color=MID_GRAY, align=PP_ALIGN.LEFT)
    add_textbox(slide, inner_x + inner_w - 0.25, endpoint_y, 0.35, 0.22,
                "100", size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    # Single band label: "Strong" / "Weak" + range — stacked and centered
    # under the colored band.
    band_center = inner_x + inner_w * ((lo + hi) / 2 / 100.0)
    label_w = 1.30
    label_y = endpoint_y + 0.22
    add_textbox(slide, band_center - label_w / 2, label_y, label_w, 0.32,
                label, size=14, bold=True, color=label_color,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, band_center - label_w / 2, label_y + 0.30,
                label_w, 0.26,
                f"{lo}\u2013{hi}",
                size=11, italic=True, color=MID_GRAY,
                align=PP_ALIGN.CENTER)


def patch_rq3_slide(slide, band, label):
    """Shared patcher for slides 24 (strong) and 25 (weak): drop the
    old dual-band badge, left-align the title, shrink the stats corner,
    and add a new single-band badge in the gap."""
    # 1. Drop old badge shapes
    dropped = 0
    for shape in list(slide.shapes):
        if not _shape_in_rect(shape, *BADGE_DELETE_RECT):
            continue
        # Don't drop the stats corner (it's at x=9.38, outside the rect)
        if shape.is_placeholder:
            continue
        # Skip the slidenum placeholder (y~7.09, already excluded by rect)
        remove_shape(slide, shape)
        dropped += 1
    print(f"  {label}: dropped {dropped} old badge shape(s)")

    # 2. Left-align title + shrink it + shrink font a hair so it fits
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt.startswith("RQ3") and "endorsements" in txt:
            shape.left = in_(0.55)
            shape.top = in_(0.30)
            shape.width = in_(6.70)
            shape.height = in_(1.40)
            # Shrink font to 22pt so the two-line title fits w=6.70
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and r.font.size.pt > 22:
                        r.font.size = Pt(22)
            break

    # 3. Shrink the stats corner slightly
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if "\u0394 Male" in txt or "Δ Male" in txt:
            shape.left = in_(10.08)
            shape.top = in_(0.30)
            shape.width = in_(2.90)
            shape.height = in_(1.40)
            # Tighten font sizes — deltas 15/17 -> 13/15, b/p 13/14 -> 12/13
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is None:
                        continue
                    cur = r.font.size.pt
                    if cur >= 17:
                        r.font.size = Pt(15)
                    elif cur >= 15:
                        r.font.size = Pt(13)
                    elif cur >= 14:
                        r.font.size = Pt(13)
                    elif cur >= 13:
                        r.font.size = Pt(12)
            break

    # 4. New compact single-band badge in the gap between title and stats
    # Title ends at x=7.25, stats starts at x=10.08 -> badge 7.40 .. 9.95
    draw_single_band_badge(
        slide, x=7.40, y=0.32, w=2.55, h=1.36, band=band,
    )
    print(f"  {label}: new single-band ({band}) badge placed")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides in deck")

    # Slide 2 — Goals
    print("\n[1] Slide 2 — Goals polish + Wharton logo")
    patch_slide2_goals(slides[1])

    # Slides 5 & 6 — year-hook overlay
    print("\n[2] Slides 5 & 6 — year-hook stamp + ? overlay")
    patch_hook_year_overlay(slides[4], "slide 5")
    patch_hook_year_overlay(slides[5], "slide 6")

    # Slide 17 — enlarge Stage 3 setup
    print("\n[3] Slide 17 — enlarge Stage 3 setup")
    patch_slide17_enlarge(slides[16])

    # Slide 22 — RQ1 combined figure
    print("\n[4] Slide 22 — combined RQ1 figure")
    patch_slide22_rq1_combined(slides[21])

    # Slides 24 & 25 — RQ3 badge + title + stats
    print("\n[5] Slides 24 & 25 — RQ3 single-band badge + left title")
    patch_rq3_slide(slides[23], band="strong", label="slide 24")
    patch_rq3_slide(slides[24], band="weak",   label="slide 25")

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
