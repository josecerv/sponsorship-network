"""
fix_orphaned_timing.py — remove orphaned <p:timing> blocks that reference
shape IDs no longer present on the slide.

Background
==========
patch_2026_04_07_content.py and patch_2026_04_08.py both rebuild certain
slides in place by calling `clear_non_placeholder_shapes()` and then
adding new shapes. That helper removes <p:sp> elements but never touches
<p:timing>, which lives at the slide root. The timing block contains
<p:bldP spid="X"/> and <p:spTgt spid="X"/> elements that reference the
OLD shape IDs by number — so after a rebuild those spids point at
nothing, and modern PowerPoint flags the slide as corrupt and offers to
repair it.

Fix
===
For every slide, collect the set of existing cNvPr ids, then walk the
timing block and check every bldP/spTgt spid. If any spid isn't in the
set, the entire <p:timing> block is stale — drop it wholesale. We don't
use any intra-slide animations on the rebuilt hook or Stage-3-randomization
slides, so deleting the timing block is safe (advance is click-to-next-slide).

Idempotent: re-running is a no-op because the next run won't find any
timing blocks with orphans.
"""

from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"


def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    print(f"  {len(prs.slides)} slides")

    fixed = 0
    for i, slide in enumerate(prs.slides, start=1):
        sld_el = slide._element  # the <p:sld> element
        timing = sld_el.find(qn("p:timing"))
        if timing is None:
            continue

        # Collect all real cNvPr ids on the slide
        real_ids = set()
        for cnvpr in sld_el.iter(qn("p:cNvPr")):
            real_ids.add(cnvpr.get("id"))

        # Walk timing targets
        orphans = set()
        for bldp in timing.iter(qn("p:bldP")):
            spid = bldp.get("spid")
            if spid and spid not in real_ids:
                orphans.add(spid)
        for spt in timing.iter(qn("p:spTgt")):
            spid = spt.get("spid")
            if spid and spid not in real_ids:
                orphans.add(spid)

        if orphans:
            print(
                f"  slide {i:2d}: {len(orphans)} orphaned spid(s) "
                f"{sorted(orphans)[:8]}{'...' if len(orphans) > 8 else ''}; "
                f"dropping timing block"
            )
            sld_el.remove(timing)
            fixed += 1

    print(f"\nRemoved {fixed} stale timing block(s)")

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
