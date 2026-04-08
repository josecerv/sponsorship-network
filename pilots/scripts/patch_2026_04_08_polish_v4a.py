"""
patch_2026_04_08_polish_v4a.py — fix Stage 2 incentive to match actual survey.

Background: I had the Stage 2 scheme written as just "$3 × evaluator's wager".
The actual Stage 2 survey (SV_54rmw8wULxvqS46) instructions said:
  - $1 flat fee for completion (guaranteed)
  - PLUS a bonus of up to $3 based on ONE randomly selected endorsement
  - Bonus conditions:
      (a) endorsed candidate must have actually scored higher
      (b) slider can't be at 50 (must have actually favored a candidate)
      (c) bonus amount = evaluator's stake % × $3

This patch surgically rewrites the lead, sub, and "If wrong" card text on
slide 13 to reflect the real instructions. The slider + reveals + table are
unchanged. The redundant bottom footer ("Paid on one randomly selected
endorsement.") is removed since the sub now carries that info.

Surgical edits only. No shape clearing. fix_fld_guids/orphaned_timing not
needed.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"

NAVY_DEEP = RGBColor(0x00, 0x14, 0x3D)
DARK_GRAY = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY  = RGBColor(0x5B, 0x65, 0x7A)
WRONG_RED = RGBColor(0xB1, 0x21, 0x21)


def patch_slide13_incentive(slide):
    changes = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()

        # 1. Replace the lead line
        if txt == "Each endorsement could earn the sponsor up to $3.00.":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = "Each sponsor earns "
            r1.font.name = "Calibri"; r1.font.size = Pt(26); r1.font.color.rgb = DARK_GRAY
            r2 = p.add_run()
            r2.text = "$1 flat"
            r2.font.name = "Calibri"; r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = NAVY_DEEP
            r3 = p.add_run()
            r3.text = " plus up to "
            r3.font.name = "Calibri"; r3.font.size = Pt(26); r3.font.color.rgb = DARK_GRAY
            r4 = p.add_run()
            r4.text = "$3 bonus"
            r4.font.name = "Calibri"; r4.font.size = Pt(26); r4.font.bold = True; r4.font.color.rgb = NAVY_DEEP
            r5 = p.add_run()
            r5.text = " per session."
            r5.font.name = "Calibri"; r5.font.size = Pt(26); r5.font.color.rgb = DARK_GRAY
            changes += 1
            print("  slide 13: lead rewritten")
            continue

        # 2. Replace the sub-line
        if "The amount depends on how confidently a downstream evaluator" in txt:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = (
                "The bonus is based on ONE randomly selected endorsement — paid only "
                "if the endorsed candidate actually scored higher, scaling with how "
                "much a downstream evaluator wagered on it."
            )
            run.font.name = "Calibri"
            run.font.size = Pt(15)
            run.font.italic = True
            run.font.color.rgb = MID_GRAY
            changes += 1
            print("  slide 13: sub rewritten")
            continue

        # 3. Update "If wrong → nothing" -> "If wrong → no bonus" for clarity
        if txt == "If wrong →  nothing":
            tf = shape.text_frame
            # Preserve styling — rewrite runs
            p = tf.paragraphs[0]
            # Find the "nothing" run and update
            for run in p.runs:
                if "nothing" in run.text:
                    run.text = run.text.replace("nothing", "no bonus")
                    changes += 1
                    print("  slide 13: 'If wrong' card updated (nothing -> no bonus)")
                    break
            continue

        # 4. Remove redundant footer (sub already says "one randomly selected")
        if txt == "Paid on one randomly selected endorsement.":
            slide.shapes._spTree.remove(shape._element)
            changes += 1
            print("  slide 13: redundant footer removed")
            continue

    return changes


def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    s_13 = prs.slides[12]
    print("\n[Phase 1] Slide 13: fix Stage 2 incentive to match actual instructions")
    n = patch_slide13_incentive(s_13)
    print(f"  {n} edits total")
    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
