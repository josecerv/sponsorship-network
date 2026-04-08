"""
patch_2026_04_08_polish_v6.py — polish follow-ups on v5.

Preserves Jose's manual edits (two-line RQ titles, grouped badges on
slides 24/25, shrunk stats corners). Surgical edits only:

 1. Slides 5 & 6 — fix the GOOD STUDENT stamp:
    - widen the stamp rect from 3.40"->4.00" and deepen from 0.95"->1.30"
      so the rotated 28pt text no longer clips the G and the check mark
    - change stamp color from red (BA1F1F) to green (15803D) because
      red reads as "bad" and we're marking a SUCCEED event
    - recenter the "one year later" sub-caption under the new stamp
      and flip it to green
    - make the "?" overlay on the Year 2 letter recede into the
      background: remove fill, switch to a very thin light-gray outline,
      shrink "?" text from 96pt->72pt and lighten it, gray the
      "how do you read this one?" caption

 2. Slide 22 — swap in the new rq1_combined.png (no sponsor-said bands,
    CIs on OVERALL panel, tighter y-range, no clipping). Only touch the
    picture — leave Jose's two-line title and stats corner alone.

 3. Slide 27 — rewrite the Running story slide:
    - keep the title "Running story"
    - drop the current body textbox
    - build a 2-column layout: left = theory bullets, right =
      prior_distributions.png (men wide/critical, women narrow/forgiving)
    - bottom strip: a "Next directions" pill with the telegraph-
      standards framing
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR = PROJECT / "pilots" / "output" / "talk_figures"
RQ1_COMBINED = FIG_DIR / "rq1_combined.png"
PRIOR_FIG = FIG_DIR / "prior_distributions.png"


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------

NAVY          = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP     = RGBColor(0x00, 0x14, 0x3D)
RED_LINE      = RGBColor(0x99, 0x00, 0x00)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TEXT    = RGBColor(0x11, 0x18, 0x27)
DARK_GRAY     = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY      = RGBColor(0x5B, 0x65, 0x7A)
BODY_GRAY     = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
VERY_LIGHT    = RGBColor(0xE5, 0xE7, 0xEB)

STAMP_GREEN   = RGBColor(0x15, 0x80, 0x3D)  # forest green — "good" verdict
OLD_STAMP_RED = RGBColor(0xBA, 0x1F, 0x1F)

Q_LIGHT_BORDER = RGBColor(0xD1, 0xD5, 0xDB)  # very light gray outline
Q_LIGHT_TEXT   = RGBColor(0xB0, 0xB7, 0xC3)  # light slate gray for ?
Q_CAPTION_GRAY = RGBColor(0x6B, 0x72, 0x80)

PRIOR_TITLE   = RGBColor(0x00, 0x14, 0x3D)
PILL_FILL     = RGBColor(0xFE, 0xF3, 0xC7)
PILL_OUTLINE  = RGBColor(0xB4, 0x53, 0x09)
PILL_TEXT     = RGBColor(0x7C, 0x2D, 0x12)

FONT = "Arial"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
                color=BLACK_TEXT, align=PP_ALIGN.LEFT, anchor=None, font=FONT):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color, font=font)
    return box


def add_rounded_rect(slide, x, y, w, h, fill=None, outline=None, outline_width=None,
                     corner=0.3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, in_(x), in_(y), in_(w), in_(h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        if outline_width is not None:
            shape.line.width = Pt(outline_width)
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    shape.text_frame.text = ""
    return shape


def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)


def _set_all_run_colors(shape, color):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = color


# ---------------------------------------------------------------------------
# Phase 1 — slides 5 & 6 stamp fix
# ---------------------------------------------------------------------------

# New stamp geometry — centered at (2.90, 3.20) on Nick's letter body,
# widened from 3.40 -> 4.00 and deepened from 0.95 -> 1.30 so a rotated
# 28pt "GOOD STUDENT ✓" no longer clips.
STAMP_CX   = 2.90
STAMP_CY   = 3.20
STAMP_W    = 4.00
STAMP_H    = 1.30
STAMP_X    = STAMP_CX - STAMP_W / 2
STAMP_Y    = STAMP_CY - STAMP_H / 2
STAMP_ROT  = 349.0  # -11 degrees


def patch_hook_stamp(slide, label):
    """Green stamp + fitted rect + green sub-caption + faded '?' overlay."""
    stamp_shape  = None
    sub_shape    = None
    q_oval       = None
    q_caption    = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if "GOOD STUDENT" in txt:
            stamp_shape = shape
        elif txt == "one year later":
            sub_shape = shape
        elif txt == "?":
            q_oval = shape
        elif txt == "how do you read this one?":
            q_caption = shape

    # --- Stamp ---
    if stamp_shape is not None:
        stamp_shape.left   = in_(STAMP_X)
        stamp_shape.top    = in_(STAMP_Y)
        stamp_shape.width  = in_(STAMP_W)
        stamp_shape.height = in_(STAMP_H)
        stamp_shape.rotation = STAMP_ROT
        stamp_shape.line.color.rgb = STAMP_GREEN
        stamp_shape.line.width = Pt(4.0)
        _set_all_run_colors(stamp_shape, STAMP_GREEN)
        print(f"  {label}: GOOD STUDENT stamp widened + recolored green")
    else:
        print(f"  {label}: WARNING no GOOD STUDENT shape found")

    # --- Sub-caption "one year later" ---
    if sub_shape is not None:
        sub_y = STAMP_Y + STAMP_H + 0.03
        sub_shape.left   = in_(STAMP_X)
        sub_shape.top    = in_(sub_y)
        sub_shape.width  = in_(STAMP_W)
        sub_shape.height = in_(0.32)
        sub_shape.rotation = STAMP_ROT
        _set_all_run_colors(sub_shape, STAMP_GREEN)
        print(f"  {label}: 'one year later' sub recentered + green")
    else:
        print(f"  {label}: WARNING no sub-caption found")

    # --- Question oval — fade into background ---
    if q_oval is not None:
        q_oval.fill.background()
        q_oval.line.color.rgb = Q_LIGHT_BORDER
        q_oval.line.width = Pt(1.25)
        # Shrink the ? text from 96pt -> 72pt and lighten
        for p in q_oval.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() == "?":
                    r.font.size = Pt(72)
                    r.font.color.rgb = Q_LIGHT_TEXT
                    r.font.bold = True
        print(f"  {label}: ? oval faded into background")
    else:
        print(f"  {label}: WARNING no ? oval found")

    # --- Question caption gray ---
    if q_caption is not None:
        _set_all_run_colors(q_caption, Q_CAPTION_GRAY)
        print(f"  {label}: ? caption grayed")


# ---------------------------------------------------------------------------
# Phase 2 — slide 22 figure swap (preserve title + stats)
# ---------------------------------------------------------------------------

def patch_slide22_figure(slide):
    from pptx.shapes.picture import Picture

    # Drop the existing picture only — title + stats corner stay
    pics = [s for s in slide.shapes if isinstance(s, Picture)]
    for p in pics:
        remove_shape(slide, p)
    print(f"  slide 22: removed {len(pics)} old picture(s)")

    if not RQ1_COMBINED.exists():
        print(f"  slide 22: WARNING no figure at {RQ1_COMBINED}")
        return

    slide.shapes.add_picture(
        str(RQ1_COMBINED),
        left=in_(0.10), top=in_(1.80),
        width=in_(13.13), height=in_(5.25),
    )
    print("  slide 22: new rq1_combined.png placed")


# ---------------------------------------------------------------------------
# Phase 3 — slide 27 Running story rewrite
# ---------------------------------------------------------------------------

RUNNING_TITLE = "Running story"

RUNNING_SUB = "Different priors \u2014 not a gender penalty"

# 4 theory bullets — header + body line
THEORY_BULLETS = [
    (
        "People hold different priors for men vs women sponsors.",
        "Men are assumed critical; women are assumed nicer on average.",
    ),
    (
        "Women's endorsements cluster high + narrow.",
        "A generous prior is hard to read signal from \u2014 "
        "evaluators can't tell a real rave from a polite rave.",
    ),
    (
        "Men's endorsements spread wide.",
        "When a critical sponsor praises, it carries information.",
    ),
    (
        "Result: women's endorsements are discounted at both ends.",
        "\"If she says 10, she doesn't know what she's talking about. "
        "If she says 1, she doesn't know either.\"",
    ),
]

NEXT_DIR_HEAD = "NEXT DIRECTIONS"
NEXT_DIR_BODY = (
    "Can women close the gap by telegraphing their standards "
    "\u2014 not by waiting on seniority?"
)


def patch_slide27_running_story(slide):
    """Drop the old body textbox. Keep the title. Build a new 2-column
    layout: theory bullets on the left, prior-distribution figure on
    the right, with a 'Next directions' pill anchored to the bottom."""
    # Keep only the title + slide number placeholder; drop the rest
    to_drop = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue  # keep slidenum placeholder
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt.startswith("Running story") and "Let" not in txt:
                # That's the title — keep it, but move it to its standard
                # template position just in case
                shape.left = in_(0.55)
                shape.top = in_(0.30)
                shape.width = in_(12.23)
                shape.height = in_(0.70)
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(32)
                        r.font.bold = True
                        r.font.name = FONT
                        r.font.color.rgb = NAVY
                continue
        to_drop.append(shape)

    for shape in to_drop:
        remove_shape(slide, shape)
    print(f"  slide 27: dropped {len(to_drop)} old body shape(s)")

    # Sub-line under the title
    add_textbox(
        slide, 0.55, 1.05, 12.23, 0.45,
        RUNNING_SUB, size=20, italic=True, color=MID_GRAY,
    )

    # ---- Left column: theory bullets ----
    left_x = 0.55
    left_y = 1.70
    left_w = 6.10
    bullet_gap = 0.15
    bullet_h = 1.08

    for i, (head, body) in enumerate(THEORY_BULLETS):
        y = left_y + i * (bullet_h + bullet_gap)

        # Navy accent square
        mark = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            in_(left_x), in_(y + 0.09),
            in_(0.14), in_(0.14),
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = NAVY
        mark.line.fill.background()
        mark.text_frame.text = ""

        box = slide.shapes.add_textbox(
            in_(left_x + 0.26), in_(y),
            in_(left_w - 0.26), in_(bullet_h),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_top = Inches(0.0)

        # Header
        p1 = tf.paragraphs[0]
        p1.space_before = Pt(0)
        r1 = p1.add_run()
        r1.text = head
        style_run(r1, 16, bold=True, color=NAVY_DEEP)

        # Body
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = body
        style_run(r2, 14, bold=False, color=DARK_GRAY)

    # ---- Right column: prior distributions figure ----
    if PRIOR_FIG.exists():
        fig_x = 6.85
        fig_y = 1.65
        fig_w = 6.30
        fig_h = 3.88
        slide.shapes.add_picture(
            str(PRIOR_FIG),
            left=in_(fig_x), top=in_(fig_y),
            width=in_(fig_w), height=in_(fig_h),
        )
        # Small italic caption under the figure
        add_textbox(
            slide, fig_x, fig_y + fig_h + 0.02, fig_w, 0.30,
            "Hypothesized prior distributions",
            size=12, italic=True, color=MID_GRAY,
            align=PP_ALIGN.CENTER,
        )
        print("  slide 27: prior_distributions.png placed in right column")
    else:
        print(f"  slide 27: WARNING no prior figure at {PRIOR_FIG}")

    # ---- Next directions pill, bottom strip ----
    pill_x = 0.55
    pill_y = 6.25
    pill_w = 12.23
    pill_h = 0.78

    pill = add_rounded_rect(
        slide, pill_x, pill_y, pill_w, pill_h,
        fill=PILL_FILL, outline=PILL_OUTLINE, outline_width=1.5,
        corner=0.5,
    )

    # Inline text inside the pill: bold label + body
    box = slide.shapes.add_textbox(
        in_(pill_x + 0.22), in_(pill_y + 0.10),
        in_(pill_w - 0.44), in_(pill_h - 0.20),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = NEXT_DIR_HEAD + "  "
    style_run(r1, 16, bold=True, color=PILL_TEXT)

    r2 = p.add_run()
    r2.text = NEXT_DIR_BODY
    style_run(r2, 16, bold=False, italic=True, color=DARK_GRAY)
    print("  slide 27: next directions pill added")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides in deck")

    print("\n[1] Slides 5 & 6 stamp/? fix")
    patch_hook_stamp(slides[4], "slide 5")
    patch_hook_stamp(slides[5], "slide 6")

    print("\n[2] Slide 22 figure swap")
    patch_slide22_figure(slides[21])

    print("\n[3] Slide 27 Running story rewrite")
    patch_slide27_running_story(slides[26])

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
