"""
patch_2026_04_08_polish_v5b.py — follow-up fix for v5.

The v5 patch missed the four bullet text boxes on slide 17 because
their text frame leads with "•  " (a separate run), so a prefix match
against "the sponsor's name" failed. This patch uses a contains match
and bumps the bullets' y-positions + font size (18 -> 20pt) to match
the rest of the enlarged slide.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"


BULLET_UPDATES = [
    ("the sponsor's name",        0.95, 2.28, 11.80, 0.52),
    ("that sponsor's endorsement", 0.95, 2.82, 11.80, 0.52),
    ("how confident the sponsor",  0.95, 3.36, 11.80, 0.52),
    ("a slider to wager part",     0.95, 3.90, 11.80, 0.52),
]


def main():
    prs = Presentation(str(DECK))
    slide = prs.slides[16]  # slide 17

    hits = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        for needle, nx, ny, nw, nh in BULLET_UPDATES:
            if needle in txt:
                shape.left = Inches(nx)
                shape.top = Inches(ny)
                shape.width = Inches(nw)
                shape.height = Inches(nh)
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text.strip():
                            r.font.size = Pt(20)
                hits += 1
                break

    print(f"  slide 17: reflowed {hits} bullet(s) to 20pt")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
