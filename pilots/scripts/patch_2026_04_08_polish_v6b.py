"""
patch_2026_04_08_polish_v6b.py — slide 27 layout fix.

The v6 layout overshot the slide: 4 bullets at h=1.08 + gaps=0.15 from
y=1.70 ran to y=6.47, colliding with the NEXT DIRECTIONS pill at y=6.25.
This pass rebuilds slide 27 in place with tighter vertical rhythm so
everything fits above the pill, and the pill sits cleanly in the
remaining bottom strip.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR = PROJECT / "pilots" / "output" / "talk_figures"
PRIOR_FIG = FIG_DIR / "prior_distributions.png"


NAVY          = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP     = RGBColor(0x00, 0x14, 0x3D)
DARK_GRAY     = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY      = RGBColor(0x5B, 0x65, 0x7A)
LIGHT_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
PILL_FILL     = RGBColor(0xFE, 0xF3, 0xC7)
PILL_OUTLINE  = RGBColor(0xB4, 0x53, 0x09)
PILL_TEXT     = RGBColor(0x7C, 0x2D, 0x12)

FONT = "Arial"


THEORY_BULLETS = [
    (
        "People hold different priors for men vs women sponsors.",
        "Men are assumed critical; women are assumed nicer on average.",
    ),
    (
        "Women's endorsements cluster high + narrow.",
        "A generous prior is hard to read signal from \u2014 "
        "evaluators can't tell a real rave from a polite rave.",
    ),
    (
        "Men's endorsements spread wide.",
        "When a critical sponsor praises, it carries information.",
    ),
    (
        "Women's endorsements get discounted at both ends.",
        "\"If she says 10 she doesn't know; if she says 1, she doesn't either.\"",
    ),
]

NEXT_DIR_HEAD = "NEXT DIRECTIONS"
NEXT_DIR_BODY = (
    "Can women close the gap by telegraphing their standards \u2014 "
    "not by waiting on seniority?"
)


def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=DARK_GRAY, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
                color=DARK_GRAY, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color, font=font)
    return box


def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)


def main():
    prs = Presentation(str(DECK))
    slide = prs.slides[26]  # slide 27

    # Drop everything except the slidenum placeholder — we rebuild from
    # scratch (title included) so the positions are consistent.
    dropped = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        remove_shape(slide, shape)
        dropped += 1
    print(f"  slide 27: dropped {dropped} shape(s) to rebuild")

    # --- Title ---
    add_textbox(
        slide, 0.55, 0.30, 12.23, 0.70,
        "Running story",
        size=32, bold=True, color=NAVY,
    )

    # --- Subtitle ---
    add_textbox(
        slide, 0.55, 1.05, 12.23, 0.42,
        "Different priors \u2014 not a gender penalty",
        size=20, italic=True, color=MID_GRAY,
    )

    # --- Left column: theory bullets (tighter rhythm) ---
    left_x = 0.55
    left_y = 1.55
    left_w = 6.10
    bullet_h = 0.96
    bullet_gap = 0.12

    for i, (head, body) in enumerate(THEORY_BULLETS):
        y = left_y + i * (bullet_h + bullet_gap)

        # Navy accent square
        mark = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            in_(left_x), in_(y + 0.08),
            in_(0.14), in_(0.14),
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = NAVY
        mark.line.fill.background()
        mark.text_frame.text = ""

        box = slide.shapes.add_textbox(
            in_(left_x + 0.26), in_(y),
            in_(left_w - 0.26), in_(bullet_h),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_top = Inches(0.0)

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = head
        style_run(r1, 15, bold=True, color=NAVY_DEEP)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = body
        style_run(r2, 13, bold=False, color=DARK_GRAY)

    # --- Right column: prior distributions figure ---
    if PRIOR_FIG.exists():
        fig_x = 6.85
        fig_y = 1.55
        fig_w = 6.30
        fig_h = 3.88
        slide.shapes.add_picture(
            str(PRIOR_FIG),
            left=in_(fig_x), top=in_(fig_y),
            width=in_(fig_w), height=in_(fig_h),
        )
        add_textbox(
            slide, fig_x, fig_y + fig_h + 0.02, fig_w, 0.30,
            "Hypothesized prior distributions",
            size=12, italic=True, color=MID_GRAY,
            align=PP_ALIGN.CENTER,
        )
        print("  slide 27: prior figure + caption added")

    # --- Next directions pill (bottom strip) ---
    pill_x = 0.55
    pill_y = 6.05
    pill_w = 12.23
    pill_h = 0.85

    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        in_(pill_x), in_(pill_y), in_(pill_w), in_(pill_h),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = PILL_FILL
    pill.line.color.rgb = PILL_OUTLINE
    pill.line.width = Pt(1.5)
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    pill.text_frame.text = ""

    # Pill text
    box = slide.shapes.add_textbox(
        in_(pill_x + 0.28), in_(pill_y + 0.12),
        in_(pill_w - 0.56), in_(pill_h - 0.24),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = NEXT_DIR_HEAD + "   "
    style_run(r1, 16, bold=True, color=PILL_TEXT)
    r2 = p.add_run()
    r2.text = NEXT_DIR_BODY
    style_run(r2, 16, bold=False, italic=True, color=DARK_GRAY)

    print("  slide 27: next directions pill added")

    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
