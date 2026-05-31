"""
patch_deck.py — surgical in-place edits to docs/UChicago-0410.pptx.

Why this script exists: the user has been editing the deck directly in
PowerPoint and we MUST preserve those edits. We cannot re-run
build_uchicago_deck.py because that rebuilds from scratch and clobbers
their changes. This script opens the existing file and only touches the
specific shapes we need to change.

Tasks:
  1. Strip every custom 'X / 20' slide-number textbox the original
     builder added.
  2. On every slide except slide 1, clone the template's native slide-
     number placeholder (from slide 10, where the user inserted a fresh
     template slide that already has it) so the auto-numbered field
     appears on every slide and updates with future inserts.
  3. Re-render the result figures (RQ1, RQ2 strong, RQ2 weak, RQ3) and
     swap the embedded image in place on slides 14-17. Position is
     preserved; nothing else on the slide is touched.
  4. On slides 15 and 16, center the 'endorser confidence' badge
     horizontally and shrink the title width so it doesn't overlap.
  5. On slide 18, replace the synthesis chart with the wider re-render
     and give it more room on the slide.
"""

import re
from copy import deepcopy
from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.oxml.ns import qn

PROJECT      = Path(__file__).resolve().parent.parent.parent
DECK         = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR      = PROJECT / "pilots" / "output" / "talk_figures"

RQ1_PNG      = FIG_DIR / "rq1_gender_x_outcome.png"
RQ2_STRONG   = FIG_DIR / "rq2_strong_only.png"
RQ2_WEAK     = FIG_DIR / "rq2_weak_only.png"
RQ3_PNG      = FIG_DIR / "rq3_initial_trust.png"
SYNTH_PNG    = FIG_DIR / "synthesis_outcome_sensitivity.png"

# Slide numbers in the CURRENT (post-edit) deck.
# HISTORY: these were 14/15/16/17/18 when the deck had 20 slides. On
# 2026-04-07 Jose inserted 2 extra incentive slides (positions 9 & 11)
# plus an appendix slide at the end, so the result/RQ slides shifted by
# +2. Always re-verify via inspect_deck.py before running this script;
# Jose is actively doing manual PowerPoint edits and the deck layout
# may drift again. See project_uchicago_talk.md for the current map.
SLIDE_RQ1            = 16
SLIDE_RQ2_STRONG     = 17
SLIDE_RQ2_WEAK       = 18
SLIDE_RQ3            = 19
SLIDE_SYNTHESIS      = 20

SLIDE_NUM_PATTERN = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def find_slide_number_placeholder(slide):
    """Return the SLIDE_NUMBER (idx=11) placeholder element if present."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == 11:
                return ph._element
        except (AttributeError, ValueError):
            continue
    return None

def strip_custom_slide_number_textboxes(slide):
    """Remove every shape whose text content matches 'N / N'."""
    spTree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if SLIDE_NUM_PATTERN.match(txt):
                spTree.remove(shape._element)
                removed += 1
    return removed

def clone_slide_number_placeholder(slide, template_element, next_id):
    """Deep-copy `template_element` and append it to `slide`'s spTree.

    Updates the cNvPr id so it's unique within the destination slide.
    Also generates a fresh creationId GUID. The placeholder type is
    sldNum so the field auto-evaluates to the slide's actual position.
    """
    new_el = deepcopy(template_element)

    # Update the cNvPr id (must be unique within slide)
    cNvPr = new_el.find(".//" + qn("p:cNvPr"))
    if cNvPr is not None:
        cNvPr.set("id", str(next_id))

    # Update creationId GUID (some renderers complain about duplicates)
    a16_ns = "http://schemas.microsoft.com/office/drawing/2014/main"
    for ext in new_el.findall(".//{%s}creationId" % a16_ns):
        ext.set("id", "{%s}" % str(uuid4()).upper())

    slide.shapes._spTree.append(new_el)
    return new_el

def replace_picture(slide, old_shape, new_png_path,
                    left=None, top=None, width=None, height=None):
    """Remove `old_shape` and add a new picture at the same (or specified)
    geometry. Returns the new picture shape."""
    # Capture geometry from the old shape unless overridden
    if left   is None: left   = old_shape.left
    if top    is None: top    = old_shape.top
    if width  is None: width  = old_shape.width
    if height is None: height = old_shape.height
    spTree = slide.shapes._spTree
    spTree.remove(old_shape._element)
    return slide.shapes.add_picture(str(new_png_path), left, top,
                                     width=width, height=height)

def find_largest_picture(slide):
    """Return the picture shape with the largest area on the slide."""
    best = None
    best_area = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            area = (shape.width or 0) * (shape.height or 0)
            if area > best_area:
                best = shape
                best_area = area
    return best

def shift_shape(shape, dx_emu):
    """Move a shape horizontally by dx_emu (positive = right, negative = left)."""
    shape.left = shape.left + dx_emu


# ---------------------------------------------------------------------------
# Phase 1: Slide-number conversion
# ---------------------------------------------------------------------------

def convert_slide_numbers(prs):
    """Replace every custom 'X / 20' textbox with a cloned native
    slide-number placeholder."""
    # Find a template slide-number placeholder to clone (slide 10 has one
    # because the user inserted a fresh template slide there)
    template_el = None
    for slide in prs.slides:
        el = find_slide_number_placeholder(slide)
        if el is not None:
            template_el = el
            print(f"  template slide-number placeholder found")
            break

    if template_el is None:
        print("  WARN: no slide-number placeholder found anywhere; skipping")
        return

    next_id = 200
    for i, slide in enumerate(prs.slides, 1):
        removed = strip_custom_slide_number_textboxes(slide)
        if removed:
            print(f"  slide {i:2d}: removed {removed} custom slide-number textbox(es)")

        if i == 1:
            # Cover slide stays unnumbered
            continue

        if find_slide_number_placeholder(slide) is None:
            clone_slide_number_placeholder(slide, template_el, next_id)
            next_id += 1
            print(f"  slide {i:2d}: cloned native slide-number placeholder")


# ---------------------------------------------------------------------------
# Phase 2: Swap result figures (slides 14-17)
# ---------------------------------------------------------------------------

def swap_result_figures(prs):
    swaps = [
        (SLIDE_RQ1,        RQ1_PNG,    "RQ1"),
        (SLIDE_RQ2_STRONG, RQ2_STRONG, "RQ2 strong"),
        (SLIDE_RQ2_WEAK,   RQ2_WEAK,   "RQ2 weak"),
        (SLIDE_RQ3,        RQ3_PNG,    "RQ3"),
    ]
    for slide_idx, png_path, label in swaps:
        slide = prs.slides[slide_idx - 1]
        pic = find_largest_picture(slide)
        if pic is None:
            print(f"  slide {slide_idx} ({label}): NO picture found, skipping")
            continue
        # Use the slide-fill geometry the user wants regardless of the old
        # picture's exact bounds (so all four are visually consistent)
        left   = Inches(0.10)
        top    = Inches(1.80)
        width  = Inches(13.13)
        height = Inches(5.25)
        replace_picture(slide, pic, png_path,
                        left=left, top=top, width=width, height=height)
        print(f"  slide {slide_idx} ({label}): swapped figure")


# ---------------------------------------------------------------------------
# Phase 3: Center strength badge on slides 15-16
# ---------------------------------------------------------------------------

def center_strength_badges(prs):
    """Move badge components left to horizontally center them on the slide.

    The badge is composed of:
      - background rounded-rect at x=7.90, w=2.40
      - header textbox 'ENDORSER CONFIDENCE' at x=8.08
      - bar track + bar fill at x=8.08
      - label textbox ('Strong' / 'Weak') at x=8.08
      - caption textbox ('very confident' / 'hedged / unsure') at x=8.08

    The stats corner ALSO sits in the top row at x=9.38, w=3.50, so a naive
    'shift everything in [7.85, 10.35]' filter would catch it. We instead
    identify badge components by *narrow* x range [7.85, 8.15] AND/OR by
    text content ('ENDORSER CONFIDENCE' / 'Strong' / 'Weak' / 'very
    confident' / 'hedged / unsure').
    """
    BADGE_LEFT_MIN = Inches(7.85)
    BADGE_LEFT_MAX = Inches(8.20)
    BADGE_TEXT_MARKERS = {
        "ENDORSER CONFIDENCE",
        "Strong", "Weak",
        "very confident",
        "hedged / unsure",
    }
    DX = Inches(5.465) - Inches(7.90)  # negative

    for slide_idx in (SLIDE_RQ2_STRONG, SLIDE_RQ2_WEAK):
        slide = prs.slides[slide_idx - 1]

        badge_shapes = []
        for shape in slide.shapes:
            try:
                left = shape.left
            except (AttributeError, TypeError):
                continue
            # Skip the figure (way to the left)
            if left < Inches(5.0):
                continue
            # Stats corner sits at x=9.38 — keep it in place
            if left >= Inches(9.0):
                continue
            # Either narrow-x range OR text-content match qualifies
            if BADGE_LEFT_MIN <= left <= BADGE_LEFT_MAX:
                badge_shapes.append(shape)
                continue
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt in BADGE_TEXT_MARKERS:
                    badge_shapes.append(shape)

        for shape in badge_shapes:
            shift_shape(shape, DX)

        # Shrink the title to make room (title is anchored at x=0.55, w=7.20)
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt.startswith("RQ2"):
                    shape.width = Inches(4.80)

        print(f"  slide {slide_idx}: shifted {len(badge_shapes)} badge shapes "
              f"left by {round(DX/914400, 2)}\"")


# ---------------------------------------------------------------------------
# Phase 4: Enlarge synthesis chart on slide 18
# ---------------------------------------------------------------------------

def enlarge_synthesis_chart(prs):
    slide = prs.slides[SLIDE_SYNTHESIS - 1]
    pic = find_largest_picture(slide)
    if pic is None:
        print(f"  slide {SLIDE_SYNTHESIS}: no picture found, skipping")
        return

    # New layout: bullets on left (compressed), big chart on right
    new_left   = Inches(6.20)
    new_top    = Inches(1.45)
    new_width  = Inches(7.00)
    new_height = Inches(5.50)
    replace_picture(slide, pic, SYNTH_PNG,
                    left=new_left, top=new_top,
                    width=new_width, height=new_height)
    print(f"  slide {SLIDE_SYNTHESIS}: replaced + enlarged synthesis chart")

    # Shrink the bullets text box (look for the one with 'endorsements' content)
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if "endorsements" in txt.lower() and shape.shape_type != 13:
                if shape.width > Inches(6.0):
                    shape.width = Inches(5.55)
                    print(f"  slide {SLIDE_SYNTHESIS}: narrowed bullets to make room")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("=" * 72)
    print(f"  Patching {DECK}")
    print("=" * 72)

    prs = Presentation(str(DECK))
    print(f"  Loaded ({len(prs.slides)} slides)")

    print("\n[1/4] Converting slide numbers to template-native placeholders...")
    convert_slide_numbers(prs)

    print("\n[2/4] Swapping result figures (slides 14-17)...")
    swap_result_figures(prs)

    print("\n[3/4] Centering strength badges (slides 15-16)...")
    center_strength_badges(prs)

    print("\n[4/4] Enlarging synthesis chart (slide 18)...")
    enlarge_synthesis_chart(prs)

    prs.save(str(DECK))
    print(f"\n  Saved {DECK}")
    print("=" * 72)


if __name__ == "__main__":
    main()
