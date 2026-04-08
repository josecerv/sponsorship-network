"""
fix_fld_guids.py — fix the duplicate <a:fld id="..."> GUIDs that were
causing PowerPoint to prompt to repair the deck on open.

Background
==========
Per ECMA-376 part 1 section 20.1.2.2.9, the `a:fld id` attribute is a
unique identifier for a field within the document. When `patch_deck.py`
cloned the template-slide slide-number placeholder onto every slide to
get auto-updating slide numbers, its `clone_slide_number_placeholder`
helper only freshened the cNvPr id and the a16:creationId — it never
touched the inner `<a:fld id="{...}">` GUID. The result: 25 slides all
sharing one field GUID, which modern PowerPoint flags as corruption and
offers to repair.

This script walks every slide, locates every <a:fld> descendant, and
assigns each one a fresh UUID4 GUID. It's idempotent-ish — re-running
just re-randomizes the GUIDs, which is harmless.

Uses python-pptx + lxml. Back up the deck first.
"""

from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.oxml.ns import qn

PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
A_FLD = "{%s}fld" % A_NS


def fresh_guid():
    return "{%s}" % str(uuid4()).upper()


def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    print(f"  {len(prs.slides)} slides")

    total_fields = 0
    for i, slide in enumerate(prs.slides, start=1):
        # Walk the entire spTree looking for any <a:fld> element
        sp_tree = slide.shapes._spTree
        flds = sp_tree.findall(".//" + A_FLD)
        for fld in flds:
            new = fresh_guid()
            old = fld.get("id", "")
            fld.set("id", new)
            total_fields += 1
        if flds:
            print(f"  slide {i:2d}: updated {len(flds)} field GUID(s)")

    print(f"\nTotal field GUIDs freshened: {total_fields}")

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
