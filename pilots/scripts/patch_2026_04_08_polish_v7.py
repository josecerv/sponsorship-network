"""
patch_2026_04_08_polish_v7.py — em dash cleanup + slide 27 fold-in.

Two related goals from Jose's evening pass:

 1. Remove EM dashes (\u2014) wherever they appear in the deck. Jose
    doesn't want them anywhere; he thought slide 27 was the only one
    but a scan turned up four others (slides 7, 14, 18, 20). All five
    get an in-place text replacement to comma or period as fits the
    sentence. (En dashes \u2013, used for ranges like 0\u2013100, are
    preserved.)

 2. Slide 27 — fold the NEXT DIRECTIONS pill into the bullet list as a
    fourth bullet, bump the bullet font sizes (15/13 -> 20/16), and
    keep Jose's manual edits to bullets 2 and 3.

Surgical edits only. The figure (prior_distributions.png) and Jose's
edited bullet 2/3 text are preserved.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PROJECT = Path(__file__).resolve().parent.parent.parent
DECK = PROJECT / "docs" / "UChicago-0410.pptx"
FIG_DIR = PROJECT / "pilots" / "output" / "talk_figures"
PRIOR_FIG = FIG_DIR / "prior_distributions.png"


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------

NAVY       = RGBColor(0x01, 0x1F, 0x5B)
NAVY_DEEP  = RGBColor(0x00, 0x14, 0x3D)
DARK_GRAY  = RGBColor(0x1F, 0x29, 0x37)
MID_GRAY   = RGBColor(0x5B, 0x65, 0x7A)
BLACK_TEXT = RGBColor(0x11, 0x18, 0x27)

FONT = "Arial"

EM_DASH = "\u2014"


# ---------------------------------------------------------------------------
# Em dash replacements (whole-string match -> replacement)
# ---------------------------------------------------------------------------

# (slide_index, find substring containing em dash, replacement substring)
EM_DASH_FIXES = [
    # Slide 7 — "What is sponsorship?"
    (
        6,
        " \u2014 with the intention of",
        ", with the intention of",
    ),
    # Slide 14 — Stage 2 incentive sub
    (
        13,
        " \u2014 paid only if",
        ", paid only if",
    ),
    # Slide 18 — Stage 3 randomization (strength card sub)
    (
        17,
        " \u2014 pushed near one end",
        ". Pushed near one end",
    ),
    # Slide 20 — Step 2 outcome reveal
    (
        19,
        " \u2014 the evaluator's bank goes up",
        ". The evaluator's bank goes up",
    ),
]


def fix_em_dashes(slides):
    """Walk only the targeted slides; rewrite any matching run text."""
    fixed = 0
    for slide_idx, needle, replacement in EM_DASH_FIXES:
        slide = slides[slide_idx]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if needle in run.text:
                        run.text = run.text.replace(needle, replacement)
                        fixed += 1
                        print(f"  slide {slide_idx+1}: replaced em dash in run")
                        break
    print(f"  total em-dash fixes: {fixed}")


# ---------------------------------------------------------------------------
# Slide 27 — fold pill into bullet list, bump font sizes
# ---------------------------------------------------------------------------

# Layout constants — must keep the figure at the same position so the
# right-column visual stays exactly where Jose likes it.
LEFT_X     = 0.55
LEFT_Y     = 1.55
LEFT_W     = 6.10
BULLET_H   = 1.20
BULLET_GAP = 0.10

HEADER_PT  = 20
BODY_PT    = 16

FIG_X = 6.85
FIG_Y = 1.55
FIG_W = 6.30
FIG_H = 3.88


def in_(v): return Inches(v)


def style_run(run, size, bold=False, italic=False, color=BLACK_TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
                color=BLACK_TEXT, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.runs[0], size, bold=bold, italic=italic, color=color, font=font)
    return box


def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)


def _read_paragraphs(shape):
    """Return list of paragraph strings from a textbox shape."""
    out = []
    for p in shape.text_frame.paragraphs:
        out.append("".join(r.text for r in p.runs))
    return out


def patch_slide27(slide):
    # ---- 1. Read existing bullet content so we preserve Jose's edits ----
    # Existing bullets are at shape indices 3 (b1), 5 (b2), 7 (b3) — text
    # boxes after the navy accent square. We find them by walking the
    # text-bearing shapes in y order and pulling the ones in the left
    # column (x ~ 0.55..6.65).
    bullet_texts = []
    for shape in sorted(slide.shapes, key=lambda s: s.top if s.top else 0):
        if not shape.has_text_frame:
            continue
        try:
            x = shape.left / 914400
            y = shape.top / 914400
        except Exception:
            continue
        # Skip title (y < 1.5), figure caption (x > 6.5), pill text
        if y < 1.50:
            continue
        if x > 6.50:
            continue
        # Skip the pill text — it lives at y > 6.0
        if y > 5.90:
            continue
        paras = _read_paragraphs(shape)
        if len(paras) >= 2 and paras[0].strip() and paras[1].strip():
            bullet_texts.append((paras[0], paras[1]))

    print(f"  slide 27: read {len(bullet_texts)} existing bullet(s)")

    # Fall back to defaults if anything is missing (3 expected)
    DEFAULT_BULLETS = [
        ("People hold different priors for men vs women sponsors.",
         "Men are assumed critical; women are assumed nicer on average."),
        ("Women's endorsements cluster high + narrow.",
         "A generous prior is hard to read signal from and evaluators "
         "discount their endorsements."),
        ("Men's endorsements spread wide.",
         "When men endorse with confidence it carries information which "
         "can be both praised or penalized."),
    ]
    while len(bullet_texts) < 3:
        bullet_texts.append(DEFAULT_BULLETS[len(bullet_texts)])

    # The new fourth bullet replaces the NEXT DIRECTIONS pill.
    # NO em dash anywhere.
    bullet_texts.append((
        "Next direction: can women telegraph their standards?",
        "Signaling \"I only recommend the top 10%\" could close the gap "
        "without years of seniority.",
    ))

    # Ensure no em dashes leak in from Jose's edits — strip them defensively.
    bullet_texts = [
        (h.replace(EM_DASH, ", "), b.replace(EM_DASH, ", "))
        for (h, b) in bullet_texts
    ]

    # ---- 2. Drop everything except title + slidenum + figure + caption ----
    keep_text_prefixes = ("Running story", "Hypothesized prior")
    dropped = 0
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        # Keep the figure (picture)
        if shape.shape_type == 13:  # PICTURE
            continue
        # Keep title + caption text
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if any(txt.startswith(p) for p in keep_text_prefixes):
                continue
        remove_shape(slide, shape)
        dropped += 1
    print(f"  slide 27: dropped {dropped} old body shape(s)")

    # ---- 3. Re-add subtitle (it was a textbox we just dropped) ----
    add_textbox(
        slide, 0.55, 1.05, 12.23, 0.42,
        "Different priors, not a gender penalty",
        size=20, italic=True, color=MID_GRAY,
    )

    # ---- 4. Re-add 4 bullets with bigger fonts ----
    for i, (head, body) in enumerate(bullet_texts):
        y = LEFT_Y + i * (BULLET_H + BULLET_GAP)

        # Navy accent square (slightly larger to match the bumped fonts)
        mark = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            in_(LEFT_X), in_(y + 0.16),
            in_(0.18), in_(0.18),
        )
        mark.fill.solid()
        mark.fill.fore_color.rgb = NAVY
        mark.line.fill.background()
        mark.text_frame.text = ""

        box = slide.shapes.add_textbox(
            in_(LEFT_X + 0.30), in_(y),
            in_(LEFT_W - 0.30), in_(BULLET_H),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_top = Inches(0.0)

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = head
        style_run(r1, HEADER_PT, bold=True, color=NAVY_DEEP)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = body
        style_run(r2, BODY_PT, bold=False, color=DARK_GRAY)

    print(f"  slide 27: added {len(bullet_texts)} bullet(s) at {HEADER_PT}/{BODY_PT}pt")


def main():
    print(f"Reading {DECK}")
    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    print(f"  {len(slides)} slides in deck")

    print("\n[1] Em dash cleanup across slides 7, 14, 18, 20")
    fix_em_dashes(slides)

    print("\n[2] Slide 27 fold pill into bullet list + bigger fonts")
    patch_slide27(slides[26])

    print(f"\nSaving {DECK}")
    prs.save(str(DECK))
    print("Done.")


if __name__ == "__main__":
    main()
